"""
WebShare Pro - Media Routes
미디어 스트리밍, 썸네일, 갤러리, 플레이리스트, 문서 미리보기
"""

import os
import re
import io
import mimetypes
from collections import OrderedDict
from markupsafe import escape
from flask import Blueprint, jsonify, request, send_file, abort, session

from config import (
    conf, cache_lock,
    IMAGE_EXTENSIONS, VIDEO_EXTENSIONS, AUDIO_EXTENSIONS
)
from utils.api_errors import api_exception
from utils.log_manager import logger
from utils.file_utils import validate_path, get_real_ip, get_file_type
from utils.request_policy import ensure_mutation_allowed, ensure_path_access, parse_json_body
from security.auth import login_required
from features.audit_log import log_audit
from utils.helpers import add_recent_file, atomic_write_bytes, build_recent_owner_key

media_bp = Blueprint('media', __name__)

# 썸네일 캐시
THUMBNAIL_CACHE = OrderedDict()
MAX_THUMBNAIL_CACHE = 200
MAX_TEXT_EDIT_SIZE = 10 * 1024 * 1024


def _recent_owner_key() -> str:
    return build_recent_owner_key(
        session_id=session.get('session_id', '') or '',
        role=session.get('role', 'guest') or 'guest',
        ip=get_real_ip(),
    )


# ==========================================
# 미디어 스트리밍
# ==========================================

@media_bp.route('/stream/<path:filepath>')
@login_required()
def stream_media(filepath):
    """HTTP Range 요청을 지원하는 미디어 스트리밍"""
    from flask import current_app
    from utils.helpers import check_download_limit, track_download

    ok, message, status_code = ensure_path_access(filepath, 'read')
    if not ok:
        return jsonify({'error': message}), status_code
    
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return abort(404)
    
    file_size = os.path.getsize(full_path)
    mime_type, _ = mimetypes.guess_type(full_path)
    if not mime_type:
        mime_type = 'application/octet-stream'

    client_ip = get_real_ip()
    range_header = request.headers.get('Range')
    
    if range_header:
        # Range 요청 파싱
        byte_start = 0
        byte_end = file_size - 1
        
        match = re.match(r'bytes=(\d+)-(\d*)', range_header)
        if match:
            byte_start = int(match.group(1))
            if match.group(2):
                byte_end = int(match.group(2))
        
        byte_end = min(byte_end, file_size - 1)
        content_length = byte_end - byte_start + 1
        if content_length <= 0:
            return abort(416)
        allowed, limit_msg = check_download_limit(client_ip, False, projected_bytes=content_length)
        if not allowed:
            return jsonify({'error': limit_msg}), 429
        track_download(client_ip, content_length, False)
        if byte_start == 0:
            add_recent_file(
                filepath,
                os.path.basename(full_path),
                get_file_type(os.path.splitext(full_path)[1]),
                owner_key=_recent_owner_key(),
            )
        
        def generate():
            with open(full_path, 'rb') as f:
                f.seek(byte_start)
                remaining = content_length
                chunk_size = 1024 * 1024  # 1MB chunks
                while remaining > 0:
                    read_size = min(chunk_size, remaining)
                    data = f.read(read_size)
                    if not data:
                        break
                    remaining -= len(data)
                    yield data
        
        response = current_app.response_class(
            generate(),
            status=206,
            mimetype=mime_type,
            direct_passthrough=True
        )
        response.headers['Content-Range'] = f'bytes {byte_start}-{byte_end}/{file_size}'
        response.headers['Content-Length'] = content_length
        response.headers['Accept-Ranges'] = 'bytes'
        return response
    else:
        allowed, limit_msg = check_download_limit(client_ip, False, projected_bytes=file_size)
        if not allowed:
            return jsonify({'error': limit_msg}), 429
        track_download(client_ip, file_size, False)
        add_recent_file(
            filepath,
            os.path.basename(full_path),
            get_file_type(os.path.splitext(full_path)[1]),
            owner_key=_recent_owner_key(),
        )
        # 전체 파일 스트리밍
        def generate_full():
            with open(full_path, 'rb') as f:
                while True:
                    data = f.read(1024 * 1024)
                    if not data:
                        break
                    yield data
        
        response = current_app.response_class(
            generate_full(),
            mimetype=mime_type,
            direct_passthrough=True
        )
        response.headers['Content-Length'] = file_size
        response.headers['Accept-Ranges'] = 'bytes'
        return response


# ==========================================
# 이미지 썸네일
# ==========================================

@media_bp.route('/thumbnail/<path:filepath>')
@login_required()
def get_thumbnail(filepath):
    """이미지 썸네일 생성 (LRU 캐시)"""
    try:
        from PIL import Image
    except ImportError:
        return abort(500)

    ok, message, status_code = ensure_path_access(filepath, 'read')
    if not ok:
        return jsonify({'error': message}), status_code
    
    is_valid, full_path, _ = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.exists(full_path):
        return abort(404)
    
    cache_key = f"{filepath}_{os.path.getmtime(full_path)}"
    
    # 캐시 확인
    with cache_lock:
        if cache_key in THUMBNAIL_CACHE:
            return send_file(io.BytesIO(THUMBNAIL_CACHE[cache_key]), mimetype='image/jpeg')
    
    try:
        img = Image.open(full_path)
        img.thumbnail((150, 150), Image.Resampling.LANCZOS)
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        buffer = io.BytesIO()
        img.save(buffer, format='JPEG', quality=70)
        buffer.seek(0)
        
        # 캐시 저장
        with cache_lock:
            if len(THUMBNAIL_CACHE) >= MAX_THUMBNAIL_CACHE:
                THUMBNAIL_CACHE.popitem(last=False)
            THUMBNAIL_CACHE[cache_key] = buffer.getvalue()
        
        buffer.seek(0)
        return send_file(buffer, mimetype='image/jpeg')
    except Exception as e:
        logger.add(f"썸네일 생성 실패: {e}", "ERROR")
        return abort(500)


# ==========================================
# 동영상 썸네일
# ==========================================

@media_bp.route('/video_thumbnail/<path:filepath>')
@login_required()
def video_thumbnail(filepath):
    """동영상 썸네일 반환"""
    from features.metadata import generate_video_thumbnail

    ok, message, status_code = ensure_path_access(filepath, 'read')
    if not ok:
        return jsonify({'error': message}), status_code
    
    is_valid, full_path, _ = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return abort(404)
    
    thumb_path = generate_video_thumbnail(full_path)
    if thumb_path and os.path.exists(thumb_path):
        return send_file(thumb_path, mimetype='image/jpeg')
    
    return abort(404)


# ==========================================
# 오디오 플레이리스트
# ==========================================

@media_bp.route('/playlist/<path:folder_path>')
@login_required()
def get_playlist(folder_path):
    """폴더 내 오디오 파일 플레이리스트"""
    ok, message, status_code = ensure_path_access(folder_path, 'read')
    if not ok:
        return jsonify({'error': message}), status_code

    is_valid, full_path, error = validate_path(conf.get('folder'), folder_path)
    if not is_valid or not os.path.isdir(full_path):
        return jsonify({'error': '폴더를 찾을 수 없습니다.'}), 404
    
    audio_extensions = {'.mp3', '.wav', '.ogg', '.m4a', '.flac', '.aac', '.wma'}
    tracks = []
    
    for name in sorted(os.listdir(full_path)):
        ext = os.path.splitext(name)[1].lower()
        if ext in audio_extensions:
            rel_path = os.path.join(folder_path, name).replace('\\', '/')
            ok, _, _ = ensure_path_access(rel_path, 'read')
            if not ok:
                continue
            tracks.append({
                'name': name,
                'path': rel_path,
                'stream_url': f'/stream/{rel_path}'
            })
    
    return jsonify({'folder': folder_path, 'tracks': tracks, 'count': len(tracks)})


# ==========================================
# 이미지 갤러리
# ==========================================

@media_bp.route('/gallery/<path:folder_path>')
@login_required()
def get_gallery(folder_path):
    """폴더 내 이미지 갤러리"""
    ok, message, status_code = ensure_path_access(folder_path, 'read')
    if not ok:
        return jsonify({'error': message}), status_code

    is_valid, full_path, error = validate_path(conf.get('folder'), folder_path)
    if not is_valid or not os.path.isdir(full_path):
        return jsonify({'error': '폴더를 찾을 수 없습니다.'}), 404
    
    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg'}
    images = []
    
    for name in sorted(os.listdir(full_path)):
        ext = os.path.splitext(name)[1].lower()
        if ext in image_extensions:
            rel_path = os.path.join(folder_path, name).replace('\\', '/')
            ok, _, _ = ensure_path_access(rel_path, 'read')
            if not ok:
                continue
            images.append({
                'name': name,
                'path': rel_path,
                'url': f'/download/{rel_path}',
                'thumbnail': f'/thumbnail/{rel_path}'
            })
    
    return jsonify({'folder': folder_path, 'images': images, 'count': len(images)})


# ==========================================
# 문서 미리보기
# ==========================================

@media_bp.route('/preview/<path:filepath>')
@login_required()
def document_preview(filepath):
    """문서 미리보기 (Word, Excel, PowerPoint, CSV, JSON)"""
    import json as json_module

    ok, message, status_code = ensure_path_access(filepath, 'read')
    if not ok:
        return jsonify({'error': message}), status_code
    
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return jsonify({'error': '파일을 찾을 수 없습니다.'}), 404
    
    ext = os.path.splitext(full_path)[1].lower()
    content = ""
    preview_type = "text"
    
    try:
        # Word (.docx)
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(full_path)
                paragraphs = []
                for para in doc.paragraphs[:100]:
                    if para.text.strip():
                        paragraphs.append(f"<p>{escape(para.text)}</p>")
                content = "\n".join(paragraphs) if paragraphs else "<p>문서가 비어있습니다.</p>"
                preview_type = "html"
            except ImportError:
                content = "python-docx 라이브러리가 필요합니다. pip install python-docx"
        
        # Excel (.xlsx)
        elif ext in ['.xlsx', '.xls']:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(full_path, read_only=True, data_only=True)
                sheet = wb.active
                rows = []
                if sheet is not None:
                    for i, row in enumerate(sheet.iter_rows(max_row=50, values_only=True)):
                        if i >= 50:
                            break
                        cells = "".join([f"<td>{escape(str(cell)) if cell is not None else ''}</td>" for cell in row[:20]])
                        rows.append(f"<tr>{cells}</tr>")
                content = f"<table border='1' style='border-collapse:collapse; width:100%;'>{''.join(rows)}</table>"
                preview_type = "html"
                wb.close()
            except ImportError:
                content = "openpyxl 라이브러리가 필요합니다. pip install openpyxl"
        
        # PowerPoint (.pptx)
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(full_path)
                slides_content = []
                for i, slide in enumerate(list(prs.slides)[:20]):
                    slide_text = []
                    for shape in slide.shapes:
                        shape_text = getattr(shape, "text", "")
                        if isinstance(shape_text, str) and shape_text.strip():
                            slide_text.append(shape_text)
                    if slide_text:
                        escaped_text = '<br>'.join([str(escape(t)) for t in slide_text])
                        slides_content.append(f"<div style='border:1px solid #ccc; padding:15px; margin:10px 0; border-radius:8px;'><strong>슬라이드 {i+1}</strong><br>{escaped_text}</div>")
                content = "".join(slides_content) if slides_content else "<p>프레젠테이션이 비어있습니다.</p>"
                preview_type = "html"
            except ImportError:
                content = "python-pptx 라이브러리가 필요합니다. pip install python-pptx"
        
        # CSV
        elif ext == '.csv':
            import csv
            with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = []
                for i, row in enumerate(reader):
                    if i >= 100: break
                    cells = "".join([f"<td>{escape(cell)}</td>" for cell in row[:20]])
                    rows.append(f"<tr>{cells}</tr>")
                content = f"<table border='1' style='border-collapse:collapse; width:100%;'>{''.join(rows)}</table>"
                preview_type = "html"
        
        # JSON
        elif ext == '.json':
            with open(full_path, 'r', encoding='utf-8') as f:
                data = json_module.load(f)
                safe_json = escape(json_module.dumps(data, ensure_ascii=False, indent=2)[:10000])
                content = f"<pre>{safe_json}</pre>"
                preview_type = "html"
        
        else:
            content = "지원하지 않는 파일 형식입니다."
        
        return jsonify({
            'success': True,
            'content': content,
            'type': preview_type,
            'filename': os.path.basename(full_path),
            'safe_html': preview_type == 'html',
        })
        
    except Exception as exc:
        return api_exception('문서 미리보기 오류', exc, extra={'success': False})


# ==========================================
# 텍스트 파일 읽기/쓰기
# ==========================================

@media_bp.route('/get_content/<path:path>')
@login_required()
def get_content(path):
    """텍스트 파일 내용 읽기"""
    ok, message, status_code = ensure_path_access(path, 'read')
    if not ok:
        return jsonify({'error': message}), status_code

    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid:
        return jsonify({'error': error}), 403
    if not os.path.isfile(full_path):
        return jsonify({'error': '파일을 찾을 수 없습니다.'}), 404

    file_size = os.path.getsize(full_path)
    if file_size > MAX_TEXT_EDIT_SIZE:
        return jsonify({
            'error': f'파일 크기가 너무 큽니다. 최대 {MAX_TEXT_EDIT_SIZE // (1024 * 1024)}MB까지 편집할 수 있습니다.',
            'max_bytes': MAX_TEXT_EDIT_SIZE,
            'file_size': file_size,
        }), 413
    
    try:
        with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        add_recent_file(
            path,
            os.path.basename(full_path),
            get_file_type(os.path.splitext(full_path)[1]),
            owner_key=_recent_owner_key(),
        )
        return jsonify({'content': content})
    except Exception as exc:
        return api_exception('파일 읽기 오류', exc)


@media_bp.route('/save_content/<path:path>', methods=['POST'])
@login_required()
def save_content(path):
    """텍스트 파일 저장"""
    from utils.helpers import create_file_version

    role = session.get('role', 'guest')
    allowed, message, status_code = ensure_mutation_allowed(role)
    if not allowed:
        return jsonify({'success': False, 'error': message}), status_code

    ok, message, status_code = ensure_path_access(path, 'write', role=role)
    if not ok:
        return jsonify({'success': False, 'error': message}), status_code
    
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid:
        return jsonify({'success': False, 'error': error}), 403
    
    try:
        data = parse_json_body(request)
        content = data.get('content', '')
        if not isinstance(content, str):
            content = str(content)
        content_bytes = len(content.encode('utf-8'))
        if content_bytes > MAX_TEXT_EDIT_SIZE:
            return jsonify({
                'success': False,
                'error': f'저장할 내용이 너무 큽니다. 최대 {MAX_TEXT_EDIT_SIZE // (1024 * 1024)}MB까지 저장할 수 있습니다.',
                'max_bytes': MAX_TEXT_EDIT_SIZE,
                'content_size': content_bytes,
            }), 413

        # 수정 전 버전 백업
        create_file_version(full_path)

        atomic_write_bytes(full_path, content.encode('utf-8'))
        logger.add(f"파일수정: {path}")
        
        # 감사 로그 기록
        log_audit(
            user=session.get('role', 'unknown'),
            action='file_edit',
            target=path,
            details=f"크기: {content_bytes} 바이트",
            ip=get_real_ip()
        )
        
        return jsonify({'success': True})
    except Exception as exc:
        return api_exception('파일 저장 오류', exc, extra={'success': False})


# ==========================================
# HLS 트랜스코딩 스트리밍 (v7.2.3)
# ==========================================

@media_bp.route('/stream/hls/<path:filepath>/index.m3u8')
@login_required()
def stream_hls_playlist(filepath):
    """HLS 플레이리스트 반환 (트랜스코딩 시작)"""
    from features.transcoder import get_transcoder
    from utils.helpers import check_download_limit, track_download

    ok, message, status_code = ensure_path_access(filepath, 'read')
    if not ok:
        return jsonify({'error': message}), status_code
    
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.exists(full_path):
        return abort(404)
        
    try:
        client_ip = get_real_ip()
        transcoder = get_transcoder(full_path)
        
        # 파일이 생성될 때까지 잠시 대기
        for _ in range(20):
            if os.path.exists(transcoder.playlist_path):
                playlist_size = os.path.getsize(transcoder.playlist_path)
                allowed, limit_msg = check_download_limit(client_ip, False, projected_bytes=playlist_size)
                if not allowed:
                    return jsonify({'error': limit_msg}), 429
                track_download(client_ip, playlist_size, False)
                add_recent_file(
                    filepath,
                    os.path.basename(full_path),
                    get_file_type(os.path.splitext(full_path)[1]),
                    owner_key=_recent_owner_key(),
                )
                return send_file(transcoder.playlist_path, mimetype='application/vnd.apple.mpegurl')
            import time
            time.sleep(0.5)
            
        return abort(503, description="Transcoding timeout")
    except Exception as exc:
        logger.add(f"트랜스코딩 오류: {exc}", "ERROR")
        return abort(500)

@media_bp.route('/stream/hls/<path:filepath>/<segment>')
@login_required()
def stream_hls_segment(filepath, segment):
    """HLS 세그먼트 반환"""
    from features.transcoder import get_transcoder
    from utils.helpers import check_download_limit, track_download
    
    # 세그먼트 파일명 검증
    if not re.match(r'segment_\d+\.ts', segment):
        return abort(404)

    ok, message, status_code = ensure_path_access(filepath, 'read')
    if not ok:
        return jsonify({'error': message}), status_code
        
    is_valid, full_path, _ = validate_path(conf.get('folder'), filepath)
    if not is_valid: 
        return abort(404)
        
    try:
        # 세션 찾기 (이미 생성되어 있어야 함)
        transcoder = get_transcoder(full_path)
        seg_path = os.path.join(transcoder.output_dir, segment)
        
        if os.path.exists(seg_path):
            file_size = os.path.getsize(seg_path)
            
            # 대역폭 제한 확인
            client_ip = get_real_ip()
            allowed, limit_msg = check_download_limit(client_ip, False, projected_bytes=file_size)
            if not allowed:
                 # HLS는 429를 받으면 재생이 멈출 수 있음.
                 # 하지만 정책상 차단해야 함.
                 return jsonify({'error': limit_msg}), 429
            
            # 통계 및 기록
            track_download(client_ip, file_size, False)
            
            return send_file(seg_path, mimetype='video/MP2T')
            
        return abort(404)
    except Exception:
        return abort(404)

