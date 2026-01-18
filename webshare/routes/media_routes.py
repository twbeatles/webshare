"""
WebShare Pro - Media Routes
미디어 스트리밍, 썸네일, 갤러리, 플레이리스트, 문서 미리보기
"""

import os
import re
import io
import mimetypes
from collections import OrderedDict
from markupsafe import escape
from flask import Blueprint, jsonify, request, send_file, abort, session

from ..config import (
    conf, cache_lock,
    IMAGE_EXTENSIONS, VIDEO_EXTENSIONS, AUDIO_EXTENSIONS
)
from ..utils.log_manager import logger
from ..utils.file_utils import validate_path, get_real_ip
from ..security.auth import login_required
from ..features.audit_log import log_audit

media_bp = Blueprint('media', __name__)

# 썸네일 캐시
THUMBNAIL_CACHE = OrderedDict()
MAX_THUMBNAIL_CACHE = 200


# ==========================================
# 미디어 스트리밍
# ==========================================

@media_bp.route('/stream/<path:filepath>')
@login_required()
def stream_media(filepath):
    """HTTP Range 요청을 지원하는 미디어 스트리밍"""
    from flask import current_app
    
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.exists(full_path):
        return abort(404)
    
    file_size = os.path.getsize(full_path)
    mime_type, _ = mimetypes.guess_type(full_path)
    if not mime_type:
        mime_type = 'application/octet-stream'
    
    range_header = request.headers.get('Range')
    
    if range_header:
        # Range 요청 파싱
        byte_start = 0
        byte_end = file_size - 1
        
        match = re.match(r'bytes=(\d+)-(\d*)', range_header)
        if match:
            byte_start = int(match.group(1))
            if match.group(2):
                byte_end = int(match.group(2))
        
        byte_end = min(byte_end, file_size - 1)
        content_length = byte_end - byte_start + 1
        
        def generate():
            with open(full_path, 'rb') as f:
                f.seek(byte_start)
                remaining = content_length
                chunk_size = 1024 * 1024  # 1MB chunks
                while remaining > 0:
                    read_size = min(chunk_size, remaining)
                    data = f.read(read_size)
                    if not data:
                        break
                    remaining -= len(data)
                    yield data
        
        response = current_app.response_class(
            generate(),
            status=206,
            mimetype=mime_type,
            direct_passthrough=True
        )
        response.headers['Content-Range'] = f'bytes {byte_start}-{byte_end}/{file_size}'
        response.headers['Content-Length'] = content_length
        response.headers['Accept-Ranges'] = 'bytes'
        return response
    else:
        # 전체 파일 스트리밍
        def generate_full():
            with open(full_path, 'rb') as f:
                while True:
                    data = f.read(1024 * 1024)
                    if not data:
                        break
                    yield data
        
        response = current_app.response_class(
            generate_full(),
            mimetype=mime_type,
            direct_passthrough=True
        )
        response.headers['Content-Length'] = file_size
        response.headers['Accept-Ranges'] = 'bytes'
        return response


# ==========================================
# 이미지 썸네일
# ==========================================

@media_bp.route('/thumbnail/<path:filepath>')
@login_required()
def get_thumbnail(filepath):
    """이미지 썸네일 생성 (LRU 캐시)"""
    try:
        from PIL import Image
    except ImportError:
        return abort(500)
    
    is_valid, full_path, _ = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.exists(full_path):
        return abort(404)
    
    cache_key = f"{filepath}_{os.path.getmtime(full_path)}"
    
    # 캐시 확인
    with cache_lock:
        if cache_key in THUMBNAIL_CACHE:
            return send_file(io.BytesIO(THUMBNAIL_CACHE[cache_key]), mimetype='image/jpeg')
    
    try:
        img = Image.open(full_path)
        img.thumbnail((150, 150), Image.Resampling.LANCZOS)
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        buffer = io.BytesIO()
        img.save(buffer, format='JPEG', quality=70)
        buffer.seek(0)
        
        # 캐시 저장
        with cache_lock:
            if len(THUMBNAIL_CACHE) >= MAX_THUMBNAIL_CACHE:
                THUMBNAIL_CACHE.popitem(last=False)
            THUMBNAIL_CACHE[cache_key] = buffer.getvalue()
        
        buffer.seek(0)
        return send_file(buffer, mimetype='image/jpeg')
    except Exception as e:
        logger.add(f"썸네일 생성 실패: {e}", "ERROR")
        return abort(500)


# ==========================================
# 동영상 썸네일
# ==========================================

@media_bp.route('/video_thumbnail/<path:filepath>')
@login_required()
def video_thumbnail(filepath):
    """동영상 썸네일 반환"""
    from ..features.metadata import generate_video_thumbnail
    
    is_valid, full_path, _ = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return abort(404)
    
    thumb_path = generate_video_thumbnail(full_path)
    if thumb_path and os.path.exists(thumb_path):
        return send_file(thumb_path, mimetype='image/jpeg')
    
    return abort(404)


# ==========================================
# 오디오 플레이리스트
# ==========================================

@media_bp.route('/playlist/<path:folder_path>')
@login_required()
def get_playlist(folder_path):
    """폴더 내 오디오 파일 플레이리스트"""
    is_valid, full_path, error = validate_path(conf.get('folder'), folder_path)
    if not is_valid or not os.path.isdir(full_path):
        return jsonify({'error': '폴더를 찾을 수 없습니다.'}), 404
    
    audio_extensions = {'.mp3', '.wav', '.ogg', '.m4a', '.flac', '.aac', '.wma'}
    tracks = []
    
    for name in sorted(os.listdir(full_path)):
        ext = os.path.splitext(name)[1].lower()
        if ext in audio_extensions:
            rel_path = os.path.join(folder_path, name).replace('\\', '/')
            tracks.append({
                'name': name,
                'path': rel_path,
                'stream_url': f'/stream/{rel_path}'
            })
    
    return jsonify({'folder': folder_path, 'tracks': tracks, 'count': len(tracks)})


# ==========================================
# 이미지 갤러리
# ==========================================

@media_bp.route('/gallery/<path:folder_path>')
@login_required()
def get_gallery(folder_path):
    """폴더 내 이미지 갤러리"""
    is_valid, full_path, error = validate_path(conf.get('folder'), folder_path)
    if not is_valid or not os.path.isdir(full_path):
        return jsonify({'error': '폴더를 찾을 수 없습니다.'}), 404
    
    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg'}
    images = []
    
    for name in sorted(os.listdir(full_path)):
        ext = os.path.splitext(name)[1].lower()
        if ext in image_extensions:
            rel_path = os.path.join(folder_path, name).replace('\\', '/')
            images.append({
                'name': name,
                'path': rel_path,
                'url': f'/download/{rel_path}',
                'thumbnail': f'/thumbnail/{rel_path}'
            })
    
    return jsonify({'folder': folder_path, 'images': images, 'count': len(images)})


# ==========================================
# 문서 미리보기
# ==========================================

@media_bp.route('/preview/<path:filepath>')
@login_required()
def document_preview(filepath):
    """문서 미리보기 (Word, Excel, PowerPoint, CSV, JSON)"""
    import json as json_module
    
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return jsonify({'error': '파일을 찾을 수 없습니다.'}), 404
    
    ext = os.path.splitext(full_path)[1].lower()
    content = ""
    preview_type = "text"
    
    try:
        # Word (.docx)
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(full_path)
                paragraphs = []
                for para in doc.paragraphs[:100]:
                    if para.text.strip():
                        paragraphs.append(f"<p>{escape(para.text)}</p>")
                content = "\n".join(paragraphs) if paragraphs else "<p>문서가 비어있습니다.</p>"
                preview_type = "html"
            except ImportError:
                content = "python-docx 라이브러리가 필요합니다. pip install python-docx"
        
        # Excel (.xlsx)
        elif ext in ['.xlsx', '.xls']:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(full_path, read_only=True, data_only=True)
                sheet = wb.active
                rows = []
                for i, row in enumerate(sheet.iter_rows(max_row=50, values_only=True)):
                    if i >= 50: break
                    cells = "".join([f"<td>{escape(str(cell)) if cell is not None else ''}</td>" for cell in row[:20]])
                    rows.append(f"<tr>{cells}</tr>")
                content = f"<table border='1' style='border-collapse:collapse; width:100%;'>{''.join(rows)}</table>"
                preview_type = "html"
                wb.close()
            except ImportError:
                content = "openpyxl 라이브러리가 필요합니다. pip install openpyxl"
        
        # PowerPoint (.pptx)
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(full_path)
                slides_content = []
                for i, slide in enumerate(prs.slides[:20]):
                    if i >= 20: break
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        escaped_text = '<br>'.join([str(escape(t)) for t in slide_text])
                        slides_content.append(f"<div style='border:1px solid #ccc; padding:15px; margin:10px 0; border-radius:8px;'><strong>슬라이드 {i+1}</strong><br>{escaped_text}</div>")
                content = "".join(slides_content) if slides_content else "<p>프레젠테이션이 비어있습니다.</p>"
                preview_type = "html"
            except ImportError:
                content = "python-pptx 라이브러리가 필요합니다. pip install python-pptx"
        
        # CSV
        elif ext == '.csv':
            import csv
            with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = []
                for i, row in enumerate(reader):
                    if i >= 100: break
                    cells = "".join([f"<td>{escape(cell)}</td>" for cell in row[:20]])
                    rows.append(f"<tr>{cells}</tr>")
                content = f"<table border='1' style='border-collapse:collapse; width:100%;'>{''.join(rows)}</table>"
                preview_type = "html"
        
        # JSON
        elif ext == '.json':
            with open(full_path, 'r', encoding='utf-8') as f:
                data = json_module.load(f)
                content = f"<pre>{json_module.dumps(data, ensure_ascii=False, indent=2)[:10000]}</pre>"
                preview_type = "html"
        
        else:
            content = "지원하지 않는 파일 형식입니다."
        
        return jsonify({
            'success': True,
            'content': content,
            'type': preview_type,
            'filename': os.path.basename(full_path)
        })
        
    except Exception as e:
        logger.add(f"문서 미리보기 오류: {e}", "ERROR")
        return jsonify({'success': False, 'error': str(e)})


# ==========================================
# 텍스트 파일 읽기/쓰기
# ==========================================

@media_bp.route('/get_content/<path:path>')
@login_required()
def get_content(path):
    """텍스트 파일 내용 읽기"""
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid:
        return jsonify({'error': error}), 403
    
    try:
        with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
            return jsonify({'content': f.read()})
    except Exception as e:
        logger.add(f"파일 읽기 오류: {e}", "ERROR")
        return jsonify({'error': str(e)})


@media_bp.route('/save_content/<path:path>', methods=['POST'])
@login_required('admin')
def save_content(path):
    """텍스트 파일 저장"""
    from ..utils.helpers import create_file_version
    
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid:
        return jsonify({'success': False, 'error': error}), 403
    
    try:
        # 수정 전 버전 백업
        create_file_version(full_path)
        
        content = request.get_json().get('content', '')
        with open(full_path, 'w', encoding='utf-8') as f:
            f.write(content)
        logger.add(f"파일수정: {path}")
        
        # 감사 로그 기록
        log_audit(
            user=session.get('role', 'unknown'),
            action='file_edit',
            target=path,
            details=f"크기: {len(content)} 바이트",
            ip=get_real_ip()
        )
        
        return jsonify({'success': True})
    except Exception as e:
        logger.add(f"파일 저장 오류: {e}", "ERROR")
        return jsonify({'success': False, 'error': str(e)})
