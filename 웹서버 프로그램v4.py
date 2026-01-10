import os
import sys
import socket
import threading
import webbrowser
import mimetypes
import json
import shutil
import zipfile
import io
import time
import logging
import queue
import re
import hashlib
import secrets
from datetime import datetime, timedelta
from functools import wraps

# Windows DPI Awareness 설정 (PyQt6 import 전에 설정해야 함)
# import sys, import os 중복 제거됨

# Qt DPI 경고 메시지 억제
os.environ['QT_LOGGING_RULES'] = 'qt.qpa.window=false'

if sys.platform == 'win32':
    try:
        import ctypes
        ctypes.windll.shcore.SetProcessDpiAwareness(2)  # PROCESS_PER_MONITOR_DPI_AWARE_V2
    except (AttributeError, OSError, Exception):
        try:
            ctypes.windll.user32.SetProcessDPIAware()
        except (AttributeError, OSError, Exception):
            pass

# GUI Imports - PyQt6
try:
    from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, 
                                  QPushButton, QLabel, QLineEdit, QComboBox, QCheckBox, QTabWidget,
                                  QTextEdit, QFileDialog, QMessageBox, QGroupBox, QFrame, QSizePolicy,
                                  QSpacerItem, QDialog, QDialogButtonBox, QScrollArea, QSystemTrayIcon, QMenu)
    from PyQt6.QtCore import Qt, QTimer, pyqtSignal, QThread, QSize
    from PyQt6.QtGui import QFont, QIcon, QPalette, QColor, QPixmap, QImage, QAction
    PYQT6_AVAILABLE = True
except ImportError:
    PYQT6_AVAILABLE = False
    # Fallback to Tkinter if PyQt6 not installed
    import tkinter as tk
    from tkinter import ttk, messagebox, scrolledtext, filedialog

from PIL import Image  # Requires: pip install pillow
if not PYQT6_AVAILABLE:
    from PIL import ImageTk  # Tkinter GUI에서만 필요

# Server Imports
from flask import Flask, request, send_from_directory, render_template_string, redirect, url_for, session, abort, send_file, jsonify, g
from werkzeug.serving import make_server
from werkzeug.security import generate_password_hash, check_password_hash
# secure_filename은 한글을 지원하지 않으므로 직접 구현한 함수를 사용합니다.

# ==========================================
# 1. 설정 및 상수 (Constants)
# ==========================================
APP_TITLE = "WebShare Pro v6.0"
CONFIG_FILE = "webshare_config.json"
DEFAULT_PORT = 5000
TEXT_EXTENSIONS = {
    '.txt', '.py', '.html', '.css', '.js', '.json', '.md', '.log', '.xml', '.ini', '.conf', 
    '.sh', '.bat', '.c', '.cpp', '.h', '.java', '.sql', '.yaml', '.yml',
    # v6.0: 추가 확장자
    '.rs', '.go', '.kt', '.swift', '.ts', '.tsx', '.jsx', '.rb', '.php', '.pl',
    '.dockerfile', '.env', '.toml', '.csv', '.scss', '.less', '.sass', '.vue', '.svelte',
    '.r', '.m', '.mm', '.scala', '.clj', '.ex', '.exs', '.hs', '.lua', '.ps1', '.psm1'
}
MAX_LOG_LINES = 1000
SESSION_TIMEOUT_MINUTES = 30  # 세션 만료 시간 (분)
VERSION_FOLDER_NAME = ".webshare_versions"  # 파일 버전 저장 폴더
MAX_FILE_VERSIONS = 5  # 최대 버전 수

# ==========================================
# 스레드 동기화 락 (Thread Locks)
# ==========================================
import threading
from collections import OrderedDict  # v7.1: LRU 캐시용
_stats_lock = threading.Lock()
_share_links_lock = threading.Lock()
_access_log_lock = threading.Lock()
_login_attempts_lock = threading.Lock()
_metadata_lock = threading.Lock()  # 태그, 즐겨찾기, 메모용
_cache_lock = threading.Lock()  # 썸네일, 다운로드 추적용
_session_lock = threading.Lock()  # 활성 세션 추적용
# v7.1: 추가 락
_download_tracker_lock = threading.Lock()
_download_tracker_lock = threading.Lock()
_recent_files_lock = threading.Lock()
_upload_session_lock = threading.Lock()  # v7.1: 업로드 세션 동기화

# 서버 통계 전역 변수
SERVER_START_TIME = datetime.now()
STATS = {
    'requests': 0,
    'bytes_sent': 0,
    'bytes_received': 0,
    'errors': 0,
    'active_connections': 0  # 현재 접속자 수
}

# 공유 링크 저장소 (메모리 저장, 서버 재시작 시 초기화)
# 형식: {token: {'path': 경로, 'expires': 만료시간, 'created_by': 생성자}}
SHARE_LINKS = {}

# 북마크 저장소
BOOKMARKS = []

# 접속 기록 (최대 100개)
ACCESS_LOG = []
MAX_ACCESS_LOG = 100

# 썸네일 캐시 (메모리, 최대 200개)
THUMBNAIL_CACHE = OrderedDict()
MAX_THUMBNAIL_CACHE = 200

# 휴지통 폴더명
TRASH_FOLDER_NAME = ".webshare_trash"

# v5.1: 활성 세션 추적 {session_id: {'ip': IP, 'login_time': datetime, 'role': 역할, 'last_active': datetime}}
ACTIVE_SESSIONS = {}

# v5.1: 최근 파일 목록 (전역, 최대 20개)
RECENT_FILES = []
MAX_RECENT_FILES = 20

# v5.1: 다운로드 추적 {ip: {'count': 횟수, 'bytes': 용량, 'date': 날짜}}
DOWNLOAD_TRACKER = {}

# ==========================================
# v7.0 신규 전역 변수
# ==========================================

# v7.0: 로그인 시도 추적 {ip: {'attempts': 실패횟수, 'blocked_until': 차단해제시간}}
LOGIN_ATTEMPTS = {}
MAX_LOGIN_ATTEMPTS = 5  # 최대 로그인 실패 횟수
LOGIN_BLOCK_MINUTES = 15  # IP 차단 시간 (분)

# v7.0: 파일 태그 저장소 {경로: [{'tag': 태그명, 'color': 색상}]}
FILE_TAGS = {}

# v7.0: 즐겨찾기 폴더 저장소
FAVORITE_FOLDERS = []

# v7.0: 파일 메모 저장소 {경로: {'memo': 내용, 'updated': 시간}}
FILE_MEMOS = {}

# v7.0: 휴지통 자동 비우기 설정
TRASH_AUTO_DELETE_DAYS = 7  # 기본 7일

# v7.0: 동영상 썸네일 캐시 폴더명
VIDEO_THUMB_FOLDER = ".webshare_thumbs"


# v5.1: 다국어 지원
I18N = {
    'ko': {
        'login': '접속하기',
        'logout': '로그아웃',
        'upload': '업로드',
        'download': '다운로드',
        'delete': '삭제',
        'rename': '이름 변경',
        'new_folder': '새 폴더',
        'search': '파일 검색...',
        'empty_folder': '폴더가 비어있습니다',
        'drag_hint': '파일을 드래그하거나 업로드 버튼을 클릭하세요',
        'recent_files': '최근 파일',
        'no_recent': '최근 파일이 없습니다',
        'settings': '설정',
        'server_status': '서버 상태',
        'active_users': '접속자',
        'disk_warning': '디스크 용량 경고!',
        'download_limit': '다운로드 한도 초과',
        'ip_blocked': 'IP가 허용 목록에 없습니다',
        'admin': '관리자',
        'guest': '게스트',
        'save': '저장',
        'cancel': '취소',
        'close': '닫기',
    },
    'en': {
        'login': 'Login',
        'logout': 'Logout',
        'upload': 'Upload',
        'download': 'Download',
        'delete': 'Delete',
        'rename': 'Rename',
        'new_folder': 'New Folder',
        'search': 'Search files...',
        'empty_folder': 'Folder is empty',
        'drag_hint': 'Drag files here or click upload',
        'recent_files': 'Recent Files',
        'no_recent': 'No recent files',
        'settings': 'Settings',
        'server_status': 'Server Status',
        'active_users': 'Active Users',
        'disk_warning': 'Low disk space warning!',
        'download_limit': 'Download limit exceeded',
        'ip_blocked': 'Your IP is not allowed',
        'admin': 'Admin',
        'guest': 'Guest',
        'save': 'Save',
        'cancel': 'Cancel',
        'close': 'Close',
    }
}

# ==========================================
# 2. 유틸리티 함수 (Utility Functions)
# ==========================================

def safe_filename(filename):
    """
    Werkzeug의 secure_filename은 한글을 모두 삭제하므로,
    한글을 지원하는 안전한 파일명 변환 함수를 구현합니다.
    """
    # 1. 경로 구분자 제거 (보안)
    filename = filename.replace('/', '').replace('\\', '')
    
    # 2. 상위 디렉토리 탐색(..) 방지
    filename = re.sub(r'\.\.+', '.', filename)
    
    # 3. 윈도우/리눅스 예약 문자 제거 또는 치환
    # < > : " / \ | ? *
    filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
    
    # 4. 공백 및 제어 문자 처리
    filename = filename.strip()
    
    # 5. 빈 파일명 방지
    if not filename:
        filename = "unnamed_file"
        
    return filename

def validate_path(base_dir: str, path: str) -> tuple:
    """
    경로 탐색 공격을 방지하기 위한 경로 검증 함수.
    
    Args:
        base_dir: 기본 허용 디렉토리
        path: 검증할 상대 경로
        
    Returns:
        tuple: (is_valid: bool, full_path: str, error_msg: str)
    """
    try:
        full_path = os.path.normpath(os.path.join(base_dir, path))
        base_dir_normalized = os.path.normpath(os.path.abspath(base_dir))
        
        # 경로가 기본 디렉토리 내에 있는지 확인
        if not os.path.abspath(full_path).startswith(base_dir_normalized):
            return (False, None, "잘못된 경로입니다.")
        
        return (True, full_path, None)
    except (OSError, ValueError) as e:
        return (False, None, f"경로 검증 오류: {str(e)}")

class LogManager:
    def __init__(self):
        self.queue = queue.Queue()

    def add(self, msg, level="INFO"):
        timestamp = datetime.now().strftime("%H:%M:%S")
        formatted_msg = f"[{timestamp}] [{level}] {msg}"
        self.queue.put(formatted_msg)
        print(formatted_msg) 

logger = LogManager()

# ==========================================
# 2.1 보안 헬퍼 함수 (v4 추가)
# ==========================================
# hashlib은 상단에서 import됨

def hash_password(password: str) -> str:
    """
    비밀번호 해싱 (Werkzeug의 pbkdf2:sha256 사용)
    v7.1: 단순 SHA256에서 솔트가 포함된 안전한 해시로 변경
    """
    return generate_password_hash(password)

def verify_password(stored_password: str, provided_password: str) -> bool:
    """비밀번호 검증 (구버전 호환성 포함)"""
    try:
        # 1. v7.1 신규 해시 (pbkdf2/scrypt, 형식: method$salt$hash)
        if '$' in stored_password:
            return check_password_hash(stored_password, provided_password)
            
        # 2. v4~v7.0 구버전 해시 (SHA256, 64자)
        if len(stored_password) == 64:
            # 구버전 해시 검증
            legacy_hash = hashlib.sha256(provided_password.encode('utf-8')).hexdigest()
            if legacy_hash == stored_password:
                # [TODO] 여기서 자동 마이그레이션을 할 수도 있음 (설정 파일 쓰기 권한 필요)
                return True
            return False
            
        # 3. 평문 비밀번호 (v3 이하 또는 설정파일 직접 수정 시)
        return stored_password == provided_password
    except Exception:
        return False

def generate_csrf_token():
    """CSRF 토큰 생성"""
    if '_csrf_token' not in session:
        session['_csrf_token'] = secrets.token_hex(16)
    return session['_csrf_token']

def log_access(ip: str, action: str, details: str = ""):
    """접속 기록 저장 (스레드 안전)"""
    global ACCESS_LOG
    entry = {
        'time': datetime.now().isoformat(),
        'ip': ip,
        'action': action,
        'details': details
    }
    with _access_log_lock:
        ACCESS_LOG.insert(0, entry)
        # 최대 개수 제한
        if len(ACCESS_LOG) > MAX_ACCESS_LOG:
            ACCESS_LOG = ACCESS_LOG[:MAX_ACCESS_LOG]

def create_file_version(file_path: str):
    """파일 수정 전 버전 자동 백업"""
    if not os.path.exists(file_path):
        return
    
    base_dir = conf.get('folder')
    version_dir = os.path.join(base_dir, VERSION_FOLDER_NAME)
    os.makedirs(version_dir, exist_ok=True)
    
    filename = os.path.basename(file_path)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    version_name = f"{timestamp}_{filename}"
    version_path = os.path.join(version_dir, version_name)
    
    try:
        shutil.copy2(file_path, version_path)
        
        # 최대 버전 수 유지 (오래된 것 삭제)
        versions = sorted([
            f for f in os.listdir(version_dir) 
            if f.endswith(f'_{filename}')
        ], reverse=True)
        
        for old_version in versions[MAX_FILE_VERSIONS:]:
            os.remove(os.path.join(version_dir, old_version))
            
    except (OSError, IOError, shutil.Error) as e:
        logger.add(f"버전 생성 실패: {e}", "ERROR")

def extract_original_name_from_trash(trash_name: str) -> str:
    """
    휴지통 파일명에서 원본 파일명 추출.
    형식: YYYYMMDD_HHMMSS_원본파일명
    정규식을 사용하여 타임스탬프 부분만 정확히 제거.
    """
    # 타임스탬프 패턴: YYYYMMDD_HHMMSS_
    pattern = r'^\d{8}_\d{6}_'
    match = re.match(pattern, trash_name)
    if match:
        return trash_name[match.end():]
    return trash_name

# ==========================================
# 2.2 v5.1 신규 헬퍼 함수
# ==========================================

def check_ip_whitelist(ip: str) -> bool:
    """v5.1: IP 화이트리스트 확인"""
    whitelist = conf.get('ip_whitelist') or []
    if not whitelist:  # 빈 배열이면 모두 허용
        return True
    return ip in whitelist or ip == '127.0.0.1'

def check_download_limit(ip: str) -> tuple:
    """v5.1: 다운로드 제한 확인. (허용여부, 메시지)"""
    global DOWNLOAD_TRACKER
    today = datetime.now().strftime('%Y-%m-%d')
    
    with _download_tracker_lock:
        # 기존 트래커 확인
        if ip not in DOWNLOAD_TRACKER or DOWNLOAD_TRACKER[ip].get('date') != today:
            DOWNLOAD_TRACKER[ip] = {'count': 0, 'bytes': 0, 'date': today}
        
        tracker = DOWNLOAD_TRACKER[ip]
        limit_count = conf.get('daily_download_limit') or 0
        limit_mb = conf.get('daily_bandwidth_limit_mb') or 0
        
        if limit_count > 0 and tracker['count'] >= limit_count:
            return (False, f"일일 다운로드 횟수 초과 ({limit_count}회)")
        
        if limit_mb > 0 and tracker['bytes'] >= limit_mb * 1024 * 1024:
            return (False, f"일일 대역폭 초과 ({limit_mb}MB)")
    
    return (True, "")

def track_download(ip: str, file_size: int):
    """v5.1: 다운로드 기록 추가"""
    global DOWNLOAD_TRACKER
    today = datetime.now().strftime('%Y-%m-%d')
    
    with _download_tracker_lock:
        if ip not in DOWNLOAD_TRACKER or DOWNLOAD_TRACKER[ip].get('date') != today:
            DOWNLOAD_TRACKER[ip] = {'count': 0, 'bytes': 0, 'date': today}
        
        DOWNLOAD_TRACKER[ip]['count'] += 1
        DOWNLOAD_TRACKER[ip]['bytes'] += file_size

def add_recent_file(path: str, name: str, file_type: str = 'file'):
    """v5.1: 최근 파일 목록에 추가"""
    global RECENT_FILES
    entry = {
        'path': path,
        'name': name,
        'type': file_type,
        'accessed': datetime.now().isoformat()
    }
    
    with _recent_files_lock:
        # 중복 제거
        RECENT_FILES = [f for f in RECENT_FILES if f['path'] != path]
        RECENT_FILES.insert(0, entry)
        # 최대 개수 유지
        if len(RECENT_FILES) > MAX_RECENT_FILES:
            RECENT_FILES = RECENT_FILES[:MAX_RECENT_FILES]

def get_text(key: str, lang: str = None) -> str:
    """v5.1: 다국어 텍스트 반환"""
    if lang is None:
        lang = conf.get('language') or 'ko'
    return I18N.get(lang, I18N['ko']).get(key, key)

def fmt_bytes(b: int) -> str:
    """바이트를 읽기 좋은 형식으로 변환 (전역 유틸리티)"""
    if b < 1024: return f"{b} B"
    elif b < 1024*1024: return f"{b/1024:.1f} KB"
    elif b < 1024*1024*1024: return f"{b/(1024*1024):.1f} MB"
    return f"{b/(1024*1024*1024):.1f} GB"

def cleanup_expired_sessions():
    """v7.0: 만료된 세션 정리 (스레드 안전)"""
    global ACTIVE_SESSIONS
    now = datetime.now()
    timeout_minutes = conf.get('session_timeout') or SESSION_TIMEOUT_MINUTES
    expired = []
    
    with _session_lock:
        for sid, info in list(ACTIVE_SESSIONS.items()):
            last_active = info.get('last_active')
            if last_active:
                age_minutes = (now - last_active).total_seconds() / 60
                if age_minutes > timeout_minutes:
                    expired.append(sid)
        
        for sid in expired:
            del ACTIVE_SESSIONS[sid]
    
    if expired:
        logger.add(f"만료 세션 정리: {len(expired)}개")
    return len(expired)

def cleanup_expired_share_links():
    """v7.0: 만료된 공유 링크 정리"""
    global SHARE_LINKS
    now = datetime.now()
    expired = []
    
    with _share_links_lock:
        for token, info in list(SHARE_LINKS.items()):
            expires = info.get('expires')
            if expires and now > expires:
                expired.append(token)
        
        for token in expired:
            del SHARE_LINKS[token]
    
    if expired:
        logger.add(f"만료 공유 링크 정리: {len(expired)}개")
    return len(expired)

def get_folder_size(folder_path: str) -> int:
    """v5.1: 폴더 크기 계산 (바이트)"""
    total_size = 0
    try:
        for dirpath, dirnames, filenames in os.walk(folder_path):
            # 숨김 폴더 제외
            dirnames[:] = [d for d in dirnames if not d.startswith('.')]
            for filename in filenames:
                filepath = os.path.join(dirpath, filename)
                if os.path.isfile(filepath):
                    total_size += os.path.getsize(filepath)
    except (OSError, IOError):
        pass
    return total_size

# ==========================================
# v7.0 신규 헬퍼 함수
# ==========================================

def check_ip_blocked(ip: str) -> tuple:
    """v7.0: IP 차단 상태 확인. (차단여부, 남은시간(분))"""
    if ip not in LOGIN_ATTEMPTS:
        return (False, 0)
    
    info = LOGIN_ATTEMPTS[ip]
    blocked_until = info.get('blocked_until')
    
    if blocked_until and datetime.now() < blocked_until:
        remaining = (blocked_until - datetime.now()).total_seconds() / 60
        return (True, round(remaining, 1))
    
    return (False, 0)

def record_login_attempt(ip: str, success: bool):
    """v7.0: 로그인 시도 기록 (스레드 안전)"""
    global LOGIN_ATTEMPTS
    
    with _login_attempts_lock:
        if success:
            # 성공 시 기록 초기화
            if ip in LOGIN_ATTEMPTS:
                del LOGIN_ATTEMPTS[ip]
            return
        
        # 실패 시 카운트 증가
        if ip not in LOGIN_ATTEMPTS:
            LOGIN_ATTEMPTS[ip] = {'attempts': 0, 'blocked_until': None}
        
        LOGIN_ATTEMPTS[ip]['attempts'] += 1
        
        # 최대 시도 초과 시 차단
        if LOGIN_ATTEMPTS[ip]['attempts'] >= MAX_LOGIN_ATTEMPTS:
            LOGIN_ATTEMPTS[ip]['blocked_until'] = datetime.now() + timedelta(minutes=LOGIN_BLOCK_MINUTES)
            logger.add(f"IP 차단됨: {ip} ({LOGIN_BLOCK_MINUTES}분)", "WARN")

def unblock_ip(ip: str) -> bool:
    """v7.0: IP 차단 해제 (스레드 안전)"""
    global LOGIN_ATTEMPTS
    with _login_attempts_lock:
        if ip in LOGIN_ATTEMPTS:
            del LOGIN_ATTEMPTS[ip]
            logger.add(f"IP 차단 해제: {ip}")
            return True
    return False

def get_blocked_ips() -> list:
    """v7.0: 현재 차단된 IP 목록"""
    now = datetime.now()
    blocked = []
    for ip, info in LOGIN_ATTEMPTS.items():
        blocked_until = info.get('blocked_until')
        if blocked_until and now < blocked_until:
            blocked.append({
                'ip': ip,
                'attempts': info.get('attempts', 0),
                'blocked_until': blocked_until.isoformat(),
                'remaining_minutes': round((blocked_until - now).total_seconds() / 60, 1)
            })
    return blocked

def encrypt_file_aes(file_path: str, password: str) -> tuple:
    """v7.0: AES 파일 암호화 (랜덤 salt 사용). (성공여부, 메시지)"""
    try:
        from cryptography.fernet import Fernet
        from cryptography.hazmat.primitives import hashes
        from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
        import base64
        
        # 랜덤 salt 생성 (16 바이트)
        salt = os.urandom(16)
        kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=100000)
        key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
        fernet = Fernet(key)
        
        # 파일 읽기
        with open(file_path, 'rb') as f:
            data = f.read()
        
        # 암호화
        encrypted = fernet.encrypt(data)
        
        # 암호화된 파일 저장 (salt + 암호화 데이터)
        # 형식: [16바이트 salt][암호화된 데이터]
        enc_path = file_path + '.enc'
        with open(enc_path, 'wb') as f:
            f.write(salt)  # salt 먼저 저장
            f.write(encrypted)
        
        # 원본 삭제
        os.remove(file_path)
        logger.add(f"파일 암호화됨: {os.path.basename(file_path)}")
        return (True, enc_path)
        
    except ImportError:
        return (False, "cryptography 라이브러리가 설치되지 않았습니다. pip install cryptography")
    except Exception as e:
        logger.add(f"암호화 오류: {e}", "ERROR")
        return (False, str(e))

def decrypt_file_aes(file_path: str, password: str) -> tuple:
    """v7.0: AES 파일 복호화 (파일에서 salt 읽기). (성공여부, 메시지)"""
    try:
        from cryptography.fernet import Fernet, InvalidToken
        from cryptography.hazmat.primitives import hashes
        from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
        import base64
        
        if not file_path.endswith('.enc'):
            return (False, "암호화된 파일이 아닙니다 (.enc 확장자 필요)")
        
        # 파일 읽기 (salt + 암호화 데이터)
        with open(file_path, 'rb') as f:
            salt = f.read(16)  # 처음 16바이트는 salt
            encrypted = f.read()
        
        if len(salt) < 16:
            return (False, "잘못된 암호화 파일 형식입니다")
        
        # salt에서 키 생성
        kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=salt, iterations=100000)
        key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
        fernet = Fernet(key)
        
        # 복호화
        try:
            decrypted = fernet.decrypt(encrypted)
        except InvalidToken:
            return (False, "비밀번호가 올바르지 않습니다")
        
        # 복호화된 파일 저장 (.enc 제거)
        dec_path = file_path[:-4]
        with open(dec_path, 'wb') as f:
            f.write(decrypted)
        
        # 암호화 파일 삭제
        os.remove(file_path)
        logger.add(f"파일 복호화됨: {os.path.basename(dec_path)}")
        return (True, dec_path)
        
    except ImportError:
        return (False, "cryptography 라이브러리가 설치되지 않았습니다")
    except Exception as e:
        logger.add(f"복호화 오류: {e}", "ERROR")
        return (False, str(e))

def auto_cleanup_trash():
    """v7.0: 휴지통 자동 비우기 (오래된 파일 삭제)"""
    base_dir = conf.get('folder')
    trash_dir = os.path.join(base_dir, TRASH_FOLDER_NAME)
    
    if not os.path.exists(trash_dir):
        return 0
    
    deleted_count = 0
    now = datetime.now()
    max_age_days = conf.get('trash_auto_delete_days') or TRASH_AUTO_DELETE_DAYS
    
    try:
        for item in os.listdir(trash_dir):
            item_path = os.path.join(trash_dir, item)
            # 타임스탬프에서 삭제 시간 추출 (형식: YYYYMMDD_HHMMSS_파일명)
            try:
                timestamp_str = item[:15]  # YYYYMMDD_HHMMSS
                deleted_time = datetime.strptime(timestamp_str, '%Y%m%d_%H%M%S')
                age_days = (now - deleted_time).days
                
                if age_days >= max_age_days:
                    if os.path.isdir(item_path):
                        shutil.rmtree(item_path)
                    else:
                        os.remove(item_path)
                    deleted_count += 1
                    logger.add(f"휴지통 자동 삭제: {item} ({age_days}일 경과)")
            except (ValueError, OSError):
                continue
    except Exception as e:
        logger.add(f"휴지통 자동 비우기 오류: {e}", "ERROR")
    
    return deleted_count

def generate_video_thumbnail(video_path: str) -> str:
    """v7.0: ffmpeg로 동영상 썸네일 생성. 썸네일 경로 반환 (없으면 빈 문자열)"""
    base_dir = conf.get('folder')
    thumb_dir = os.path.join(base_dir, VIDEO_THUMB_FOLDER)
    os.makedirs(thumb_dir, exist_ok=True)
    
    # 썸네일 파일명 생성 (영상 경로의 해시)
    video_hash = hashlib.md5(video_path.encode()).hexdigest()[:12]
    thumb_path = os.path.join(thumb_dir, f"{video_hash}.jpg")
    
    # 이미 존재하면 반환
    if os.path.exists(thumb_path):
        return thumb_path
    
    try:
        import subprocess
        # ffmpeg로 1초 지점에서 프레임 추출
        cmd = [
            'ffmpeg', '-y', '-i', video_path,
            '-ss', '00:00:01', '-vframes', '1',
            '-vf', 'scale=150:-1',
            '-q:v', '5', thumb_path
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=10)
        
        if result.returncode == 0 and os.path.exists(thumb_path):
            return thumb_path
    except FileNotFoundError:
        # ffmpeg가 설치되지 않음
        pass
    except Exception as e:
        logger.add(f"영상 썸네일 생성 실패: {e}", "ERROR")
    
    return ""

def save_metadata():
    """v7.0: 메타데이터(태그, 즐겨찾기, 메모) 파일로 저장 (스레드 안전)"""
    base_dir = conf.get('folder')
    meta_path = os.path.join(base_dir, '.webshare_meta.json')
    
    with _metadata_lock:
        data = {
            'tags': FILE_TAGS,
            'favorites': FAVORITE_FOLDERS,
            'memos': FILE_MEMOS,
            'bookmarks': BOOKMARKS,
            'updated': datetime.now().isoformat()
        }
        
        try:
            with open(meta_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.add(f"메타데이터 저장 실패: {e}", "ERROR")

def load_metadata():
    """v7.0: 메타데이터(태그, 즐겨찾기, 메모) 파일에서 로드 (스레드 안전)"""
    global FILE_TAGS, FAVORITE_FOLDERS, FILE_MEMOS, BOOKMARKS
    
    base_dir = conf.get('folder')
    meta_path = os.path.join(base_dir, '.webshare_meta.json')
    
    if not os.path.exists(meta_path):
        return
    
    with _metadata_lock:
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            FILE_TAGS = data.get('tags', {})
            FAVORITE_FOLDERS = data.get('favorites', [])
            FILE_MEMOS = data.get('memos', {})
            BOOKMARKS = data.get('bookmarks', [])
            logger.add("메타데이터 로드 완료")
        except Exception as e:
            logger.add(f"메타데이터 로드 실패: {e}", "ERROR")



class ConfigManager:
    def __init__(self):
        self.config = {
            'folder': os.path.abspath(os.path.join(os.getcwd(), 'shared_files')),
            'port': DEFAULT_PORT,
            'admin_pw': "1234",
            'guest_pw': "0000",
            'allow_guest_upload': False,
            'display_host': '0.0.0.0',
            'use_https': False,
            # v4 설정
            'session_timeout': SESSION_TIMEOUT_MINUTES,
            'enable_notifications': True,
            'enable_versioning': True,
            'minimize_to_tray': True,
            # v5.1 신규 설정
            'language': 'ko',
            'ip_whitelist': [],
            'daily_download_limit': 0,
            'daily_bandwidth_limit_mb': 0,
            'disk_warning_threshold': 90,
            # v5.1 추가 옵션
            'close_to_tray': True,              # 닫기 버튼 시 트레이로 최소화
            'autostart': False,                  # 윈도우 시작 시 자동 실행
        }
        self.load()

    def load(self):
        if not os.path.exists(self.config['folder']):
            try: 
                os.makedirs(self.config['folder'])
            except Exception as e:
                print(f"폴더 생성 실패: {e}")
        if os.path.exists(CONFIG_FILE):
            try:
                with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                    self.config.update(json.load(f))
            except (json.JSONDecodeError, IOError, KeyError) as e:
                logger.add(f"설정 로드 실패: {e}", "ERROR")

    def save(self):
        try:
            with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
                json.dump(self.config, f, indent=4)
        except (IOError, TypeError) as e:
            logger.add(f"설정 저장 실패: {e}", "ERROR")
            
    def get(self, key, default=None): return self.config.get(key, default)
    def set(self, key, value): self.config[key] = value

conf = ConfigManager()

# ==========================================
# 3. HTML 템플릿 (변경 없음)
# ==========================================

# v7.0: 공유 링크 비밀번호 입력 폼
SHARE_PASSWORD_TEMPLATE = """
<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>비밀번호 필요 - WebShare Pro</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <style>
        * { box-sizing: border-box; }
        body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); min-height: 100vh; display: flex; justify-content: center; align-items: center; margin: 0; }
        .card { background: white; padding: 40px; border-radius: 20px; box-shadow: 0 25px 50px rgba(0,0,0,0.2); text-align: center; max-width: 400px; width: 90%; }
        h2 { color: #1e293b; margin-bottom: 10px; }
        p { color: #64748b; margin-bottom: 25px; }
        input { width: 100%; padding: 15px; border: 2px solid #e2e8f0; border-radius: 12px; font-size: 1rem; margin-bottom: 15px; transition: border-color 0.3s; }
        input:focus { outline: none; border-color: #6366f1; }
        button { width: 100%; padding: 15px; background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%); color: white; border: none; border-radius: 12px; font-size: 1rem; font-weight: 600; cursor: pointer; transition: transform 0.2s, box-shadow 0.2s; }
        button:hover { transform: translateY(-2px); box-shadow: 0 10px 30px rgba(99,102,241,0.4); }
        .error { color: #ef4444; font-size: 0.9rem; margin-bottom: 15px; }
        .icon { font-size: 4rem; color: #6366f1; margin-bottom: 20px; }
    </style>
</head>
<body>
    <div class="card">
        <div class="icon"><i class="fa-solid fa-lock"></i></div>
        <h2>비밀번호 필요</h2>
        <p>이 파일에 접근하려면 비밀번호가 필요합니다.</p>
        {% if error %}<div class="error">{{ error }}</div>{% endif %}
        <form method="post">
            <input type="hidden" name="csrf_token" value="{{ csrf_token() }}">
            <input type="password" name="password" placeholder="비밀번호를 입력하세요" required autofocus>
            <button type="submit"><i class="fa-solid fa-unlock"></i> 확인</button>
        </form>
    </div>
</body>
</html>
"""

# v7.0: 공유 링크 만료/제한 초과 메시지
SHARE_EXPIRED_TEMPLATE = """
<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>접근 불가 - WebShare Pro</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <style>
        * { box-sizing: border-box; }
        body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; background: linear-gradient(135deg, #ef4444 0%, #dc2626 100%); min-height: 100vh; display: flex; justify-content: center; align-items: center; margin: 0; }
        .card { background: white; padding: 40px; border-radius: 20px; box-shadow: 0 25px 50px rgba(0,0,0,0.2); text-align: center; max-width: 400px; width: 90%; }
        h2 { color: #1e293b; margin-bottom: 10px; }
        p { color: #64748b; }
        .icon { font-size: 4rem; color: #ef4444; margin-bottom: 20px; }
    </style>
</head>
<body>
    <div class="card">
        <div class="icon"><i class="fa-solid fa-circle-xmark"></i></div>
        <h2>접근 불가</h2>
        <p>{{ message }}</p>
    </div>
</body>
</html>
"""

HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="ko" data-theme="light">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">
    <meta name="description" content="WebShare Pro - 파일 공유 및 관리 시스템">
    <title>WebShare Pro</title>
    <meta name="csrf-token" content="{{ csrf_token() }}">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    
    <script src="https://cdn.jsdelivr.net/npm/marked/marked.min.js"></script>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.7.0/styles/github-dark.min.css">
    <script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.7.0/highlight.min.js"></script>

    <style>
        :root {
            --primary: #6366f1; --primary-dark: #4f46e5; --primary-light: #a5b4fc;
            --bg: #f8fafc; --card: #ffffff; --text: #1e293b; 
            --text-secondary: #64748b; --border: #e2e8f0; --danger: #ef4444; --danger-light: #fecaca;
            --folder: #f59e0b; --hover: #f1f5f9;
            --success: #10b981; --success-dark: #059669; --success-light: #d1fae5;
            --warning: #f59e0b; --info: #3b82f6;
            --focus-ring: #818cf8;
            --gradient: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            --gradient-primary: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%);
            --gradient-success: linear-gradient(135deg, #10b981 0%, #059669 100%);
            --gradient-danger: linear-gradient(135deg, #ef4444 0%, #dc2626 100%);
            --gradient-subtle: linear-gradient(135deg, rgba(99, 102, 241, 0.1) 0%, rgba(139, 92, 246, 0.1) 100%);
            --input-bg: #ffffff;
            --shadow-sm: 0 1px 2px rgba(0,0,0,0.04);
            --shadow-md: 0 4px 12px rgba(0,0,0,0.08);
            --shadow-lg: 0 12px 40px rgba(0,0,0,0.12);
            --shadow-xl: 0 25px 50px rgba(0,0,0,0.15);
            --glow-primary: 0 0 30px rgba(99, 102, 241, 0.35);
            --glow-success: 0 0 30px rgba(16, 185, 129, 0.35);
            --glow-danger: 0 0 30px rgba(239, 68, 68, 0.35);
            --transition-fast: 0.15s cubic-bezier(0.4, 0, 0.2, 1);
            --transition-normal: 0.25s cubic-bezier(0.4, 0, 0.2, 1);
            --transition-slow: 0.4s cubic-bezier(0.4, 0, 0.2, 1);
            --glass-bg: rgba(255, 255, 255, 0.85);
            --glass-border: rgba(255, 255, 255, 0.2);
        }
        [data-theme="dark"] {
            --primary: #818cf8; --primary-dark: #6366f1; --primary-light: #c7d2fe;
            --bg: #0f172a; --card: #1e293b; --text: #f1f5f9;
            --text-secondary: #94a3b8; --border: #334155; --folder: #fbbf24; --hover: #334155;
            --danger-light: #7f1d1d; --success-light: #064e3b;
            --gradient: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            --gradient-primary: linear-gradient(135deg, #818cf8 0%, #a78bfa 100%);
            --gradient-subtle: linear-gradient(135deg, rgba(129, 140, 248, 0.15) 0%, rgba(167, 139, 250, 0.15) 100%);
            --input-bg: #1e293b;
            --shadow-sm: 0 1px 2px rgba(0,0,0,0.25);
            --shadow-md: 0 4px 12px rgba(0,0,0,0.35);
            --shadow-lg: 0 12px 40px rgba(0,0,0,0.45);
            --shadow-xl: 0 25px 50px rgba(0,0,0,0.5);
            --glow-primary: 0 0 35px rgba(129, 140, 248, 0.45);
            --glass-bg: rgba(30, 41, 59, 0.9);
            --glass-border: rgba(255, 255, 255, 0.08);
        }
        
        * { box-sizing: border-box; }
        
        body { 
            font-family: 'Pretendard', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; 
            background: var(--bg); 
            color: var(--text); 
            margin: 0; 
            transition: background var(--transition-slow), color var(--transition-normal); 
            padding-bottom: 80px; 
            -webkit-tap-highlight-color: transparent;
            line-height: 1.6;
            font-size: 15px;
        }
        
        /* Enhanced scrollbar for modern look */
        ::-webkit-scrollbar { width: 10px; height: 10px; }
        ::-webkit-scrollbar-track { background: transparent; }
        ::-webkit-scrollbar-thumb { 
            background: var(--border); 
            border-radius: 10px; 
            border: 2px solid var(--bg);
        }
        ::-webkit-scrollbar-thumb:hover { background: var(--text-secondary); }
        
        *:focus-visible { outline: 2px solid var(--focus-ring); outline-offset: 2px; border-radius: 4px; }

        .container { max-width: 1100px; margin: 0 auto; padding: 24px; }
        
        header { 
            display: flex; 
            justify-content: space-between; 
            align-items: center; 
            margin-bottom: 24px; 
            padding-bottom: 16px;
            border-bottom: 1px solid var(--border);
        }
        
        .card { 
            background: var(--glass-bg); 
            backdrop-filter: blur(20px);
            -webkit-backdrop-filter: blur(20px);
            border-radius: 20px; 
            box-shadow: var(--shadow-md); 
            border: 1px solid var(--glass-border); 
            overflow: hidden;
            transition: all var(--transition-normal);
        }
        .card:hover {
            box-shadow: var(--shadow-lg);
        }
        
        .toolbar { display: flex; gap: 12px; margin-bottom: 20px; flex-wrap: wrap; align-items: center; }
        .search-box { flex: 1; position: relative; min-width: 200px; }
        .search-box input { 
            width: 100%; 
            padding: 12px 12px 12px 44px; 
            border-radius: 14px; 
            border: 2px solid var(--border); 
            background: var(--input-bg); 
            color: var(--text); 
            box-sizing: border-box; 
            height: 48px;
            font-size: 0.95rem;
            transition: all var(--transition-fast);
        }
        .search-box input:focus { 
            border-color: var(--primary); 
            box-shadow: 0 0 0 4px rgba(99, 102, 241, 0.12); 
            background: var(--card);
        }
        .search-box input::placeholder { color: var(--text-secondary); opacity: 0.7; }
        .search-box i { 
            position: absolute; 
            left: 16px; 
            top: 50%; 
            transform: translateY(-50%); 
            color: var(--text-secondary); 
            font-size: 1rem;
            transition: color var(--transition-fast);
        }
        .search-box:focus-within i { color: var(--primary); }
        
        .sort-select { 
            padding: 0 14px; 
            height: 44px; 
            border-radius: 12px; 
            border: 1px solid var(--border); 
            background: var(--card); 
            color: var(--text); 
            cursor: pointer;
            font-size: 0.9rem;
        }

        .btn { 
            background: var(--gradient-primary); 
            color: white; 
            border: none; 
            padding: 10px 20px; 
            border-radius: 12px; 
            cursor: pointer; 
            font-weight: 600; 
            text-decoration: none; 
            display: inline-flex; 
            align-items: center; 
            gap: 8px; 
            transition: all var(--transition-normal); 
            font-size: 0.9rem; 
            height: 44px; 
            box-sizing: border-box;
            position: relative;
            overflow: hidden;
        }
        .btn::before {
            content: '';
            position: absolute;
            top: 0; left: -100%;
            width: 100%; height: 100%;
            background: linear-gradient(90deg, transparent, rgba(255,255,255,0.2), transparent);
            transition: left 0.5s;
        }
        .btn:hover::before { left: 100%; }
        .btn:hover { transform: translateY(-2px); box-shadow: var(--glow-primary); }
        .btn:active { transform: translateY(0); box-shadow: none; }
        .btn-outline { background: transparent; border: 1.5px solid var(--border); color: var(--text); }
        .btn-outline::before { display: none; }
        .btn-outline:hover { background: var(--hover); border-color: var(--primary); transform: translateY(-1px); box-shadow: none; }
        .btn-icon { width: 40px; height: 40px; padding: 0; justify-content: center; border-radius: 10px; }
        .btn-icon::before { display: none; }
        .btn-danger { background: var(--danger-light); color: var(--danger); border: 1px solid rgba(239,68,68,0.3); }
        .btn-danger::before { display: none; }
        .btn-danger:hover { background: var(--gradient-danger); color: white; box-shadow: 0 4px 15px rgba(239,68,68,0.4); }
        .btn-success { background: var(--gradient-success); color: white; }
        .btn-success:hover { box-shadow: var(--glow-success); }

        #batchBar { 
            display: none; 
            align-items: center; 
            gap: 12px; 
            background: var(--gradient); 
            color: white; 
            padding: 12px 20px; 
            border-radius: 12px; 
            animation: slideDown 0.3s ease-out;
            box-shadow: 0 4px 15px rgba(99, 102, 241, 0.3);
        }
        @keyframes slideDown { from { transform: translateY(-10px); opacity: 0; } to { transform: translateY(0); opacity: 1; } }

        .file-list { list-style: none; padding: 0; margin: 0; }
        .file-item { 
            display: flex; 
            align-items: center; 
            padding: 16px 20px; 
            border-bottom: 1px solid var(--border); 
            cursor: pointer; 
            transition: all var(--transition-fast); 
            user-select: none;
            position: relative;
            background: transparent;
        }
        .file-item::before {
            content: '';
            position: absolute;
            inset: 0;
            background: var(--gradient-subtle);
            opacity: 0;
            transition: opacity var(--transition-fast);
            pointer-events: none;
        }
        .file-item:hover::before { opacity: 1; }
        .file-item::after {
            content: '';
            position: absolute;
            left: 0; bottom: 0;
            width: 0; height: 3px;
            background: var(--gradient-primary);
            transition: width var(--transition-normal);
            border-radius: 0 3px 0 0;
        }
        .file-item:hover::after { width: 100%; }
        .file-item:hover { transform: translateX(4px); }
        .file-item.selected { 
            background: rgba(99, 102, 241, 0.08); 
            border-left: 4px solid var(--primary); 
        }
        
        .file-check { 
            margin-right: 16px; 
            width: 20px;
            height: 20px;
            cursor: pointer; 
            accent-color: var(--primary);
            transition: transform var(--transition-fast);
        }
        .file-check:hover { transform: scale(1.2); }
        .file-icon { 
            font-size: 1.6rem; 
            width: 48px; 
            text-align: center; 
            color: var(--text-secondary); 
            transition: all var(--transition-normal); 
        }
        .file-item:hover .file-icon { transform: scale(1.15) rotate(3deg); color: var(--primary); }
        .file-icon.folder { color: var(--folder); }
        .file-item:hover .file-icon.folder { color: var(--warning); transform: scale(1.15); }
        .file-info { flex: 1; min-width: 0; margin-right: 12px; }
        .file-name { 
            font-weight: 600; 
            overflow: hidden; 
            text-overflow: ellipsis; 
            white-space: nowrap; 
            font-size: 0.95rem; 
            transition: color var(--transition-fast); 
        }
        .file-item:hover .file-name { color: var(--primary); }
        .file-meta { font-size: 0.8rem; color: var(--text-secondary); margin-top: 4px; display: flex; gap: 12px; }
        .file-actions { opacity: 0; transition: all var(--transition-fast); display: flex; gap: 8px; }
        .file-item:focus-within .file-actions, .file-item:hover .file-actions { opacity: 1; transform: translateX(0); }
        
        .grid-view .file-list { 
            display: grid; 
            grid-template-columns: repeat(auto-fill, minmax(160px, 1fr)); 
            gap: 18px; 
            padding: 20px; 
        }
        .grid-view .file-item { 
            flex-direction: column; 
            text-align: center; 
            height: 190px; 
            justify-content: center; 
            border-radius: 18px; 
            border: 2px solid var(--border); 
            padding: 18px; 
            position: relative;
            transition: all var(--transition-normal);
            background: var(--card);
        }
        .grid-view .file-item::before { border-radius: 18px; }
        .grid-view .file-item::after { display: none; }
        .grid-view .file-item:hover { 
            transform: translateY(-8px) scale(1.02); 
            box-shadow: var(--shadow-lg), var(--glow-primary); 
            border-color: var(--primary);
        }
        .grid-view .file-check { position: absolute; top: 12px; left: 12px; z-index: 2; }
        .grid-view .file-icon { font-size: 3.5rem; margin-bottom: 14px; width: auto; }
        .grid-view .file-item:hover .file-icon { transform: scale(1.1); }
        .grid-view .file-info { margin: 0; width: 100%; }
        .grid-view .file-actions { display: none; } 
        .grid-view .file-item img.preview { 
            width: 100%; 
            height: 90px; 
            object-fit: cover; 
            border-radius: 12px; 
            margin-bottom: 10px; 
            transition: transform var(--transition-normal); 
        }
        .grid-view .file-item:hover img.preview { transform: scale(1.05); }

        .overlay { 
            position: fixed; 
            inset: 0; 
            background: rgba(0,0,0,0.6); 
            z-index: 2000; 
            display: none; 
            justify-content: center; 
            align-items: center; 
            backdrop-filter: blur(12px);
            -webkit-backdrop-filter: blur(12px);
            animation: fadeIn 0.25s cubic-bezier(0.4, 0, 0.2, 1);
        }
        @keyframes fadeIn { from { opacity: 0; } to { opacity: 1; } }
        
        .modal { 
            background: var(--glass-bg); 
            backdrop-filter: blur(20px);
            -webkit-backdrop-filter: blur(20px);
            padding: 32px; 
            border-radius: 28px; 
            width: 90%; 
            max-width: 440px; 
            max-height: 85vh; 
            overflow-y: auto; 
            position: relative; 
            box-shadow: var(--shadow-xl), 0 0 80px rgba(99, 102, 241, 0.2); 
            display: flex; 
            flex-direction: column;
            animation: modalSlide 0.35s cubic-bezier(0.34, 1.56, 0.64, 1);
            border: 1px solid var(--glass-border);
        }
        @keyframes modalSlide { 
            from { transform: translateY(-30px) scale(0.9); opacity: 0; } 
            to { transform: translateY(0) scale(1); opacity: 1; } 
        }
        .modal h3 {
            margin-top: 0;
            padding-bottom: 18px;
            border-bottom: 1px solid var(--border);
            display: flex;
            align-items: center;
            gap: 12px;
            font-size: 1.2rem;
            background: var(--gradient-primary);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            background-clip: text;
        }
        .modal.large { max-width: 980px; width: 95%; height: 85vh; }
        
        .context-menu { 
            position: fixed; 
            background: var(--glass-bg); 
            backdrop-filter: blur(16px);
            -webkit-backdrop-filter: blur(16px);
            border: 1px solid var(--glass-border); 
            border-radius: 16px; 
            box-shadow: var(--shadow-lg); 
            z-index: 1000; 
            display: none; 
            overflow: hidden; 
            min-width: 200px;
            padding: 8px;
            animation: contextPop 0.2s cubic-bezier(0.34, 1.56, 0.64, 1);
        }
        @keyframes contextPop { from { transform: scale(0.9) translateY(-8px); opacity: 0; } to { transform: scale(1) translateY(0); opacity: 1; } }
        .ctx-item { 
            padding: 12px 16px; 
            cursor: pointer; 
            display: flex; 
            align-items: center; 
            gap: 12px; 
            font-size: 0.9rem; 
            border-radius: 10px;
            transition: all var(--transition-fast); 
        }
        .ctx-item:hover { background: var(--hover); transform: translateX(4px); }
        .ctx-item i { width: 18px; text-align: center; color: var(--text-secondary); }
        .ctx-item:hover i { color: var(--primary); }
        .ctx-item.danger { color: var(--danger); }
        .ctx-item.danger:hover { background: rgba(239, 68, 68, 0.1); }
        .ctx-item.danger i { color: var(--danger); }

        .editor-container { flex: 1; position: relative; overflow: hidden; border: 1px solid var(--border); border-radius: 12px; margin-top: 12px; display: flex; }
        .editor-area { width: 100%; height: 100%; padding: 18px; background: var(--bg); color: var(--text); font-family: 'JetBrains Mono', 'Consolas', monospace; resize: none; border: none; box-sizing: border-box; line-height: 1.6; font-size: 14px; outline: none; }
        .markdown-body { overflow-y: auto; line-height: 1.7; padding: 18px; }
        .markdown-body pre { background: #1e293b; color: #e2e8f0; padding: 1rem; border-radius: 8px; overflow-x: auto; }
        
        .stats-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 20px; margin-top: 16px; }
        .stat-card { 
            background: var(--gradient-subtle); 
            padding: 22px; 
            border-radius: 16px; 
            border: 1px solid var(--glass-border); 
            text-align: center;
            transition: all var(--transition-normal);
            position: relative;
            overflow: hidden;
        }
        .stat-card::before {
            content: '';
            position: absolute;
            inset: 0;
            background: var(--gradient-primary);
            opacity: 0;
            transition: opacity var(--transition-normal);
        }
        .stat-card:hover { transform: translateY(-4px); box-shadow: var(--shadow-md); }
        .stat-card:hover::before { opacity: 0.05; }
        .stat-value { font-size: 1.8rem; font-weight: 700; background: var(--gradient); -webkit-background-clip: text; -webkit-text-fill-color: transparent; background-clip: text; margin: 8px 0; position: relative; }
        .stat-label { font-size: 0.85rem; color: var(--text-secondary); font-weight: 500; position: relative; }

        #toast-container { position: fixed; bottom: 36px; left: 50%; transform: translateX(-50%); z-index: 3000; display: flex; flex-direction: column; gap: 14px; }
        .toast { 
            background: var(--glass-bg); 
            backdrop-filter: blur(16px);
            -webkit-backdrop-filter: blur(16px);
            color: var(--text); 
            padding: 16px 28px; 
            border-radius: 16px; 
            font-size: 0.95rem; 
            font-weight: 500;
            animation: toastSlide 0.35s cubic-bezier(0.34, 1.56, 0.64, 1); 
            display: flex; 
            align-items: center; 
            gap: 12px;
            box-shadow: var(--shadow-lg);
            border: 1px solid var(--glass-border);
        }
        .toast.success { 
            background: linear-gradient(135deg, #10b981 0%, #059669 100%); 
            color: white;
            border: none;
        }
        .toast.error { 
            background: linear-gradient(135deg, #ef4444 0%, #dc2626 100%); 
            color: white;
            border: none;
        }
        .toast.warning { 
            background: linear-gradient(135deg, #f59e0b 0%, #d97706 100%); 
            color: white;
            border: none;
        }
        .toast.info { 
            background: var(--gradient); 
            color: white;
            border: none;
        }
        @keyframes toastSlide { from { transform: translateY(40px) scale(0.9); opacity: 0; } to { transform: translateY(0) scale(1); opacity: 1; } }
        
        #drop-zone { 
            position: fixed; 
            inset: 0; 
            background: var(--gradient); 
            z-index: 9999; 
            display: none; 
            flex-direction: column; 
            justify-content: center; 
            align-items: center; 
            color: white; 
            font-size: 1.8rem; 
            font-weight: 600;
            animation: dropZonePulse 2s ease-in-out infinite;
        }
        @keyframes dropZonePulse {
            0%, 100% { opacity: 0.95; }
            50% { opacity: 1; }
        }
        #drop-zone i { font-size: 4rem; margin-bottom: 20px; animation: dropIconBounce 1s ease-in-out infinite; }
        @keyframes dropIconBounce {
            0%, 100% { transform: translateY(0); }
            50% { transform: translateY(-10px); }
        }
        
        .disk-bar { height: 8px; background: var(--border); border-radius: 4px; overflow: hidden; margin-top: 6px; }
        .disk-fill { height: 100%; background: linear-gradient(90deg, var(--success), #34d399); width: 0%; transition: width 0.6s ease-out; }

        @media (max-width: 600px) {
            .file-actions { opacity: 1; }
            .btn span { display: none; }
            .container { padding: 12px; }
            header { flex-direction: column; gap: 10px; align-items: stretch; }
            header h1 { font-size: 1.2rem; text-align: center; }
            header nav { justify-content: center; flex-wrap: wrap; }
            .toolbar { gap: 8px; }
            .search-box { min-width: 100%; }
            .btn { padding: 8px 12px; font-size: 0.85rem; height: 40px; }
            .btn-icon { width: 36px; height: 36px; }
            .modal { padding: 20px; width: 95%; }
            .file-item { padding: 12px 14px; }
            .grid-view .file-list { grid-template-columns: repeat(auto-fill, minmax(110px, 1fr)); gap: 10px; }
            .grid-view .file-item { height: 140px; }
            #toast-container { bottom: 70px; width: 90%; }
            .toast { padding: 12px 20px; font-size: 0.85rem; }
        }
        
        /* 로딩 상태 애니메이션 */
        .btn.loading { pointer-events: none; opacity: 0.7; }
        .btn.loading::after {
            content: '';
            width: 14px;
            height: 14px;
            border: 2px solid transparent;
            border-top-color: currentColor;
            border-radius: 50%;
            animation: spin 0.8s linear infinite;
            margin-left: 8px;
        }
        @keyframes spin { to { transform: rotate(360deg); } }
        
        /* 빈 폴더 상태 개선 */
        .empty-state {
            padding: 60px 40px;
            text-align: center;
            color: var(--text-secondary);
        }
        .empty-state i { font-size: 4rem; opacity: 0.3; margin-bottom: 16px; }
        .empty-state p { margin: 8px 0; }
        .empty-state .subtitle { font-size: 0.85rem; opacity: 0.7; }
        
        /* 그리드 뷰 접근성 개선 - 호버 시 액션 표시 */
        .grid-view .file-item:hover .file-actions,
        .grid-view .file-item:focus-within .file-actions {
            display: flex;
            position: absolute;
            bottom: 8px;
            left: 50%;
            transform: translateX(-50%);
            background: rgba(0,0,0,0.7);
            padding: 6px 10px;
            border-radius: 8px;
            gap: 8px;
        }
        
        /* v5: Breadcrumb 네비게이션 */
        .breadcrumb {
            display: flex;
            align-items: center;
            gap: 6px;
            padding: 10px 0;
            font-size: 0.9rem;
            flex-wrap: wrap;
        }
        .breadcrumb a {
            color: var(--primary);
            text-decoration: none;
            padding: 4px 8px;
            border-radius: 6px;
            transition: background 0.2s;
        }
        .breadcrumb a:hover {
            background: var(--hover);
        }
        .breadcrumb .separator {
            color: var(--text-secondary);
            font-size: 0.8rem;
        }
        .breadcrumb .current {
            color: var(--text);
            font-weight: 500;
        }
        
        /* v5: 파일 목록 키보드 포커스 */
        .file-item.keyboard-focused {
            outline: 2px solid var(--primary);
            outline-offset: -2px;
        }
        
        /* v7.0: 툴팁 스타일 */
        [data-tooltip] {
            position: relative;
        }
        [data-tooltip]::after {
            content: attr(data-tooltip);
            position: absolute;
            bottom: 100%;
            left: 50%;
            transform: translateX(-50%) translateY(-4px);
            background: rgba(15, 23, 42, 0.95);
            color: white;
            padding: 8px 12px;
            border-radius: 8px;
            font-size: 0.8rem;
            font-weight: 500;
            white-space: nowrap;
            opacity: 0;
            visibility: hidden;
            transition: all 0.2s ease;
            pointer-events: none;
            z-index: 1000;
            box-shadow: 0 4px 12px rgba(0,0,0,0.2);
        }
        [data-tooltip]:hover::after,
        [data-tooltip]:focus::after {
            opacity: 1;
            visibility: visible;
            transform: translateX(-50%) translateY(-8px);
        }
        
        /* v7.0: 드롭다운 메뉴 */
        .dropdown {
            position: relative;
            display: inline-block;
        }
        .dropdown-menu {
            position: absolute;
            top: 100%;
            right: 0;
            background: var(--card);
            border: 1px solid var(--border);
            border-radius: 12px;
            box-shadow: 0 10px 40px rgba(0,0,0,0.15);
            min-width: 200px;
            padding: 8px;
            opacity: 0;
            visibility: hidden;
            transform: translateY(-10px);
            transition: all 0.2s ease;
            z-index: 1000;
        }
        .dropdown.open .dropdown-menu {
            opacity: 1;
            visibility: visible;
            transform: translateY(4px);
        }
        .dropdown-item {
            display: flex;
            align-items: center;
            gap: 10px;
            padding: 10px 14px;
            border-radius: 8px;
            cursor: pointer;
            color: var(--text);
            font-size: 0.9rem;
            transition: background 0.15s;
            text-decoration: none;
        }
        .dropdown-item:hover {
            background: var(--hover);
        }
        .dropdown-item i {
            width: 18px;
            text-align: center;
            color: var(--text-secondary);
        }
        .dropdown-divider {
            height: 1px;
            background: var(--border);
            margin: 6px 0;
        }
        
        /* v7.0: 헤더 버튼 그룹 */
        .header-actions {
            display: flex;
            align-items: center;
            gap: 6px;
        }
        .header-group {
            display: flex;
            align-items: center;
            gap: 4px;
            padding: 4px;
            background: var(--bg);
            border-radius: 12px;
            border: 1px solid var(--border);
        }
        
        /* v7.0: 모바일 하단 액션바 */
        @media (max-width: 600px) {
            .mobile-bottom-bar {
                position: fixed;
                bottom: 0;
                left: 0;
                right: 0;
                background: var(--card);
                border-top: 1px solid var(--border);
                padding: 8px 16px;
                display: flex;
                justify-content: space-around;
                gap: 8px;
                z-index: 100;
                box-shadow: 0 -4px 20px rgba(0,0,0,0.1);
            }
            .mobile-bottom-bar .btn {
                flex: 1;
                justify-content: center;
                min-height: 44px;
            }
            body { padding-bottom: 120px; }
            .dropdown-menu { 
                right: auto; 
                left: 50%; 
                transform: translateX(-50%) translateY(-10px); 
            }
            .dropdown.open .dropdown-menu {
                transform: translateX(-50%) translateY(4px);
            }
        }
        
        /* v7.0: 파일 타입 아이콘 색상 */
        .file-icon.image { color: #ec4899; }
        .file-icon.video { color: #8b5cf6; }
        .file-icon.audio { color: #06b6d4; }
        .file-icon.document { color: #3b82f6; }
        .file-icon.archive { color: #84cc16; }
        .file-icon.code { color: #f97316; }
        
        /* v7.0: 스켈레톤 로딩 */
        .skeleton {
            background: linear-gradient(90deg, var(--border) 25%, var(--hover) 50%, var(--border) 75%);
            background-size: 200% 100%;
            animation: skeleton-loading 1.5s infinite;
            border-radius: 8px;
        }
        @keyframes skeleton-loading {
            0% { background-position: 200% 0; }
            100% { background-position: -200% 0; }
        }
    </style>
</head>
<body>
    <div id="drop-zone" aria-hidden="true"><i class="fa-solid fa-cloud-arrow-up" style="font-size:4rem; margin-bottom:20px;"></i>폴더나 파일을 여기에 놓으세요</div>
    <div id="toast-container" aria-live="polite"></div>
    
    <div id="ctxMenu" class="context-menu" aria-hidden="true">
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('download')"><i class="fa-solid fa-download"></i> 다운로드</div>
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('rename')"><i class="fa-solid fa-pen"></i> 이름 변경</div>
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('info')"><i class="fa-solid fa-circle-info"></i> 상세 정보</div>
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('bookmark')"><i class="fa-solid fa-star"></i> 북마크 추가</div>
        <!-- v7.0: 새로운 메뉴 항목 -->
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('tag')"><i class="fa-solid fa-tag"></i> 태그 추가</div>
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('memo')"><i class="fa-solid fa-note-sticky"></i> 메모</div>
        <div class="ctx-item" id="ctxFavorite" role="button" tabindex="0" onclick="handleCtx('favorite')" style="display:none"><i class="fa-solid fa-folder-heart"></i> 즐겨찾기</div>
        {% if role == 'admin' %}
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('share')"><i class="fa-solid fa-link"></i> 공유 링크</div>
        <div class="ctx-item" id="ctxUnzip" role="button" tabindex="0" onclick="handleCtx('unzip')" style="display:none"><i class="fa-solid fa-box-open"></i> 압축 해제</div>
        <div class="ctx-item" id="ctxEncrypt" role="button" tabindex="0" onclick="handleCtx('encrypt')"><i class="fa-solid fa-lock"></i> 암호화</div>
        <div class="ctx-item" id="ctxDecrypt" role="button" tabindex="0" onclick="handleCtx('decrypt')" style="display:none"><i class="fa-solid fa-unlock"></i> 복호화</div>
        <div class="ctx-item" role="button" tabindex="0" onclick="handleCtx('trash')"><i class="fa-solid fa-trash-can"></i> 휴지통으로</div>
        {% endif %}
        <div class="ctx-item danger" role="button" tabindex="0" onclick="handleCtx('delete')"><i class="fa-solid fa-trash"></i> 영구 삭제</div>
    </div>

    <div class="container">
        {% if not logged_in %}
            <div style="height:80vh; display:flex; justify-content:center; align-items:center;">
                <form method="post" class="card login-card" style="
                    padding: 44px; 
                    width: 100%; 
                    max-width: 360px; 
                    text-align: center; 
                    animation: modalSlide 0.5s ease-out;
                    background: var(--card);
                    backdrop-filter: blur(20px);
                    -webkit-backdrop-filter: blur(20px);
                    border: 1px solid var(--border);
                    position: relative;
                    overflow: hidden;
                ">
                    <!-- Animated gradient border -->
                    <div style="
                        position: absolute;
                        top: -2px; left: -2px; right: -2px; bottom: -2px;
                        background: var(--gradient);
                        border-radius: 18px;
                        z-index: -1;
                        opacity: 0.5;
                        animation: pulse 3s ease-in-out infinite;
                    "></div>
                    <style>
                        @keyframes pulse { 0%, 100% { opacity: 0.3; } 50% { opacity: 0.6; } }
                        @keyframes iconFloat { 0%, 100% { transform: translateY(0); } 50% { transform: translateY(-5px); } }
                    </style>
                    
                    <div class="login-icon" style="
                        width: 80px; 
                        height: 80px; 
                        background: var(--gradient-primary); 
                        border-radius: 50%; 
                        display: flex; 
                        align-items: center; 
                        justify-content: center; 
                        margin: 0 auto 24px;
                        box-shadow: var(--glow-primary);
                        animation: iconFloat 3s ease-in-out infinite;
                    ">
                        <i class="fa-solid fa-share-nodes" style="font-size: 2rem; color: white;"></i>
                    </div>
                    <h1 style="color: var(--text); margin-top: 0; font-size: 1.6rem; margin-bottom: 8px; font-weight: 700;">WebShare Pro</h1>
                    <p style="color: var(--text-secondary); font-size: 0.9rem; margin-bottom: 28px;">안전한 파일 공유 시스템</p>
                    <label for="password" class="sr-only" style="position:absolute;width:1px;height:1px;overflow:hidden;clip:rect(0,0,0,0);">비밀번호</label>
                    <div class="input-group" style="position: relative; margin-bottom: 20px;">
                        <i class="fa-solid fa-lock" style="position: absolute; left: 16px; top: 50%; transform: translateY(-50%); color: var(--text-secondary); transition: color 0.2s;"></i>
                        <input type="password" id="password" name="password" placeholder="비밀번호 입력" required 
                               style="
                                   width: 100%; 
                                   padding: 16px 48px 16px 48px; 
                                   border-radius: 14px; 
                                   border: 2px solid var(--border); 
                                   background: var(--input-bg); 
                                   color: var(--text); 
                                   font-size: 1rem; 
                                   transition: all 0.25s;
                                   box-sizing: border-box;
                               "
                               onfocus="this.style.borderColor='var(--primary)'; this.style.boxShadow='0 0 0 4px rgba(99,102,241,0.1)'; this.previousElementSibling.style.color='var(--primary)';"
                               onblur="this.style.borderColor='var(--border)'; this.style.boxShadow='none'; this.previousElementSibling.style.color='var(--text-secondary)';">
                        <button type="button" onclick="togglePasswordVisibility()" class="pw-toggle" style="position: absolute; right: 14px; top: 50%; transform: translateY(-50%); background: none; border: none; cursor: pointer; color: var(--text-secondary); padding: 6px; transition: color 0.2s;" aria-label="비밀번호 표시">
                            <i id="pwToggleIcon" class="fa-solid fa-eye"></i>
                        </button>
                    </div>
                    <input type="hidden" name="csrf_token" value="{{ csrf_token() }}">
                    <button type="submit" class="btn" style="width: 100%; justify-content: center; padding: 16px; font-size: 1.05rem; font-weight: 600;">
                        <i class="fa-solid fa-arrow-right-to-bracket"></i> 접속하기
                    </button>
                    {% if error %}<p style="color: var(--danger); font-size: 0.9rem; margin-top: 16px; background: var(--danger-light); padding: 12px; border-radius: 10px; border: 1px solid rgba(239,68,68,0.2);" role="alert"><i class="fa-solid fa-exclamation-circle"></i> {{ error }}</p>{% endif %}
                </form>
            </div>
        {% else %}
            <header>
                <h1 style="margin:0; color:var(--primary); cursor:pointer; font-size:1.5rem" onclick="location.href='/'" tabindex="0" role="link"><i class="fa-solid fa-folder-tree"></i> WebShare</h1>
                <nav class="header-actions" aria-label="메인 메뉴">
                    <!-- 역할 배지 -->
                    <span style="background:rgba(79,70,229,0.1); color:var(--primary); padding:6px 12px; border-radius:20px; font-size:0.8rem; font-weight:bold;">
                        {{ '👑 관리자' if role == 'admin' else '👤 게스트' }}
                    </span>
                    
                    <!-- v7.0: 빠른 접근 그룹 -->
                    <div class="header-group">
                        <button class="btn btn-outline btn-icon" onclick="openModal('recentModal'); loadRecentFiles()" data-tooltip="최근 파일"><i class="fa-solid fa-clock-rotate-left"></i></button>
                        <button class="btn btn-outline btn-icon" onclick="openModal('bookmarkModal'); loadBookmarks()" data-tooltip="북마크"><i class="fa-solid fa-star"></i></button>
                        <button class="btn btn-outline btn-icon" onclick="openModal('clipModal'); loadClipboard()" data-tooltip="클립보드"><i class="fa-regular fa-clipboard"></i></button>
                    </div>
                    
                    {% if role == 'admin' %}
                    <!-- v7.0: 관리 드롭다운 -->
                    <div class="dropdown" id="adminDropdown">
                        <button class="btn btn-outline btn-icon" onclick="toggleDropdown('adminDropdown')" data-tooltip="관리 메뉴"><i class="fa-solid fa-gear"></i></button>
                        <div class="dropdown-menu">
                            <div class="dropdown-item" onclick="openModal('trashModal'); loadTrash(); closeDropdowns()">
                                <i class="fa-solid fa-trash-can"></i> 휴지통
                            </div>
                            <div class="dropdown-item" onclick="openModal('shareListModal'); loadShareLinks(); closeDropdowns()">
                                <i class="fa-solid fa-link"></i> 공유 링크
                            </div>
                            <div class="dropdown-item" onclick="openModal('sessionsModal'); loadActiveSessions(); closeDropdowns()">
                                <i class="fa-solid fa-users"></i> 접속자 현황
                            </div>
                            <div class="dropdown-item" onclick="openUserManagement(); closeDropdowns()">
                                <i class="fa-solid fa-users-gear"></i> 사용자 관리
                            </div>
                            <div class="dropdown-divider"></div>
                            <div class="dropdown-item" onclick="openModal('accessDashboardModal'); loadAccessDashboard(); closeDropdowns()">
                                <i class="fa-solid fa-chart-bar"></i> 접속 대시보드
                            </div>
                        </div>
                    </div>
                    {% endif %}
                    
                    <!-- v7.0: 설정 그룹 -->
                    <div class="header-group">
                        <button class="btn btn-outline btn-icon" onclick="openModal('statsModal'); fetchStats()" data-tooltip="서버 상태"><i class="fa-solid fa-chart-line"></i></button>
                        <button class="btn btn-outline btn-icon" onclick="toggleLanguage()" data-tooltip="한/영 전환"><i class="fa-solid fa-globe"></i></button>
                        <button class="btn btn-outline btn-icon" onclick="toggleTheme()" data-tooltip="테마 변경"><i class="fa-solid fa-moon"></i></button>
                        <button class="btn btn-outline btn-icon" onclick="openModal('helpModal')" data-tooltip="도움말"><i class="fa-solid fa-circle-question"></i></button>
                    </div>
                    
                    <!-- 로그아웃 -->
                    <a href="/logout" class="btn btn-danger btn-icon" data-tooltip="로그아웃" style="display:flex;align-items:center;text-decoration:none"><i class="fa-solid fa-power-off"></i></a>
                </nav>
            </header>

            <!-- v5: Breadcrumb 네비게이션 -->
            {% if current_path %}
            <nav class="breadcrumb" aria-label="폴더 경로">
                <a href="/"><i class="fa-solid fa-home"></i></a>
                <span class="separator">/</span>
                {% set path_parts = current_path.split('/') %}
                {% for i in range(path_parts | length) %}
                    {% if i < path_parts | length - 1 %}
                        <a href="/browse/{{ path_parts[:i+1] | join('/') }}">{{ path_parts[i] }}</a>
                        <span class="separator">/</span>
                    {% else %}
                        <span class="current">{{ path_parts[i] }}</span>
                    {% endif %}
                {% endfor %}
            </nav>
            {% endif %}

            <div class="toolbar" role="toolbar" aria-label="파일 도구">
                <div class="search-box">
                    <i class="fa-solid fa-magnifying-glass" aria-hidden="true"></i>
                    <label for="searchInput" class="sr-only" style="position:absolute;width:1px;height:1px;overflow:hidden;clip:rect(0,0,0,0);">검색</label>
                    <input type="text" id="searchInput" placeholder="파일 검색..." onkeyup="filterFiles()" aria-label="파일 검색" autocomplete="off">
                </div>
                
                <select id="sortOrder" class="sort-select" onchange="sortFiles()" aria-label="정렬 방식">
                    <option value="name">이름순</option>
                    <option value="size">크기순</option>
                    <option value="date">날짜순</option>
                </select>

                <div id="batchBar" role="region" aria-live="polite">
                    <span id="batchCount">0개 선택됨</span>
                    <button class="btn-icon" style="border:1px solid rgba(255,255,255,0.3); background:rgba(255,255,255,0.2); color:white" onclick="batchDownload()" title="일괄 다운로드" aria-label="일괄 다운로드"><i class="fa-solid fa-file-zipper"></i></button>
                    {% if can_modify %}
                    <button class="btn-icon" style="border:1px solid rgba(255,255,255,0.3); background:rgba(255,255,255,0.2); color:white" onclick="batchDelete()" title="일괄 삭제" aria-label="일괄 삭제"><i class="fa-solid fa-trash"></i></button>
                    {% endif %}
                </div>

                <div style="display:flex; gap:8px;">
                    <button class="btn btn-outline" onclick="toggleView()" title="뷰 전환" aria-label="뷰 전환"><i id="viewIcon" class="fa-solid fa-list"></i></button>
                    {% if current_path %}
                    <a href="/zip/{{ current_path }}" class="btn btn-outline" title="현재 폴더 압축 다운로드" aria-label="ZIP 다운로드" style="text-decoration:none;display:flex;align-items:center;gap:5px"><i class="fa-solid fa-file-zipper"></i> ZIP</a>
                    {% endif %}
                    {% if can_modify %}
                    <button class="btn" onclick="document.getElementById('fileInput').click()"><span>업로드</span> <i class="fa-solid fa-upload"></i></button>
                    <button class="btn btn-outline" onclick="openModal('mkdirModal')" aria-label="폴더 생성"><i class="fa-solid fa-folder-plus"></i></button>
                    {% endif %}
                </div>
            </div>
            <input type="file" id="fileInput" multiple style="display:none" onchange="handleFileSelect(this.files)">

            <main id="fileContainer" class="card" role="main">
                <ul class="file-list" id="fileList" aria-label="파일 목록">
                    {% if current_path %}
                    {% set parent_path = '/'.join(current_path.split('/')[:-1]) %}
                    {% set parent_link = '/' if parent_path == '' else '/browse/' + parent_path %}
                    <li class="file-item parent-folder" tabindex="0" role="link" onclick="location.href='{{ parent_link }}'" onkeydown="if(event.key==='Enter') location.href='{{ parent_link }}'">
                        <div class="file-icon folder"><i class="fa-solid fa-turn-up"></i></div>
                        <div class="file-info"><div class="file-name">.. (상위 폴더)</div></div>
                    </li>
                    {% endif %}
                    
                    {% for item in items %}
                    <li class="file-item data-item" 
                        tabindex="0"
                        role="listitem"
                        data-path="{{ item.rel_path }}" 
                        data-name="{{ item.name }}" 
                        data-type="{{ item.type }}" 
                        data-size="{{ item.raw_size }}" 
                        data-date="{{ item.raw_mtime }}"
                        data-ext="{{ item.ext }}"
                        oncontextmenu="openCtx(event, '{{ item.rel_path }}', '{{ item.name }}', '{{ item.type }}')"
                        onkeydown="if(event.key==='Enter') handleItemClick('{{ item.rel_path }}', '{{ item.type }}', {{ 'true' if item.is_dir else 'false' }}, '{{ item.ext }}')">
                        
                        <input type="checkbox" class="file-check" value="{{ item.name }}" onclick="event.stopPropagation(); toggleBatch(this)" aria-label="{{ item.name }} 선택">
                        
                        <div class="file-icon {{ 'folder' if item.is_dir else '' }}" aria-hidden="true">
                            {% if item.is_dir %}<i class="fa-solid fa-folder"></i>
                            {% elif item.type == 'image' %}<i class="fa-solid fa-image"></i>
                            {% elif item.type == 'video' %}<i class="fa-solid fa-film"></i>
                            {% elif item.type == 'audio' %}<i class="fa-solid fa-music"></i>
                            {% elif item.type == 'text' %}<i class="fa-solid fa-file-code"></i>
                            {% elif item.type == 'archive' %}<i class="fa-solid fa-file-zipper"></i>
                            {% elif item.ext == '.pdf' %}<i class="fa-solid fa-file-pdf"></i>
                            {% else %}<i class="fa-solid fa-file"></i>{% endif %}
                        </div>
                        
                        {% if item.type == 'image' %}<img src="/download/{{ item.rel_path }}" class="preview" style="display:none;" loading="lazy" alt="{{ item.name }}">{% endif %}
                        
                        <div class="file-info" onclick="handleItemClick('{{ item.rel_path }}', '{{ item.type }}', {{ 'true' if item.is_dir else 'false' }}, '{{ item.ext }}')">
                            <div class="file-name">{{ item.name }}</div>
                            <div class="file-meta">{{ item.size }} • {{ item.mod_time }}</div>
                        </div>
                        
                        <div class="file-actions">
                            {% if item.type == 'text' %}
                            <button class="btn-icon btn-outline" onclick="event.stopPropagation(); openEditor('{{ item.rel_path }}', '{{ item.name }}', '{{ item.ext }}')" aria-label="편집"><i class="fa-solid fa-pen"></i></button>
                            {% endif %}
                            <button class="btn-icon btn-outline" onclick="event.stopPropagation(); downloadItem('{{ item.rel_path }}')" aria-label="다운로드"><i class="fa-solid fa-download"></i></button>
                            {% if can_modify and not item.is_dir %}
                            <button class="btn-icon btn-danger" onclick="event.stopPropagation(); deleteItem('{{ item.rel_path }}')" aria-label="삭제"><i class="fa-solid fa-trash"></i></button>
                            {% endif %}
                        </div>
                    </li>
                    {% endfor %}
                    {% if not items %}
                    <div class="empty-state">
                        <i class="fa-solid fa-folder-open"></i>
                        <p>폴더가 비어있습니다</p>
                        <p class="subtitle">파일을 드래그하거나 업로드 버튼을 클릭하세요</p>
                    </div>
                    {% endif %}
                </ul>
            </main>

            <div class="disk-info" style="margin-top:20px; font-size:0.8rem; opacity:0.8;" role="status">
                <div style="display:flex; justify-content:space-between;">
                    <span><i class="fa-solid fa-hard-drive"></i> 저장소 상태</span>
                    <span id="diskText">계산 중...</span>
                </div>
                <div class="disk-bar" aria-hidden="true"><div id="diskFill" class="disk-fill"></div></div>
            </div>
            
            {% if can_modify %}
            <div style="text-align:center; margin-top:20px; font-size:0.8rem; opacity:0.6;">
                <i class="fa-solid fa-circle-info"></i> 폴더나 파일을 화면에 드래그하여 업로드하세요.
            </div>
            {% endif %}
        {% endif %}
    </div>

    <!-- Modals -->
    
    <!-- v5.1: 최근 파일 모달 -->
    <div id="recentModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:500px;">
            <h3><i class="fa-solid fa-clock-rotate-left"></i> 최근 파일</h3>
            <div id="recentList" style="max-height:400px; overflow-y:auto;"></div>
            <div style="text-align:right; margin-top:15px">
                <button class="btn" onclick="closeModal('recentModal')">닫기</button>
            </div>
        </div>
    </div>
    
    <!-- v5.1: 접속자 모니터링 모달 -->
    <div id="sessionsModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:550px;">
            <h3><i class="fa-solid fa-users"></i> 접속자 현황 <span id="sessionCount" style="font-size:0.8rem; opacity:0.7;"></span></h3>
            <div id="sessionsList" style="max-height:300px; overflow-y:auto;"></div>
            <div style="text-align:right; margin-top:15px">
                <button class="btn btn-outline" onclick="loadActiveSessions()">새로고침</button>
                <button class="btn" onclick="closeModal('sessionsModal')">닫기</button>
            </div>
        </div>
    </div>

    <div id="statsModal" class="overlay" role="dialog" aria-modal="true" aria-labelledby="statsTitle">
        <div class="modal">
            <h3 id="statsTitle"><i class="fa-solid fa-chart-line"></i> 서버 상태</h3>
            <div class="stats-grid">
                <div class="stat-card">
                    <div class="stat-value" id="st_uptime">-</div>
                    <div class="stat-label">가동 시간</div>
                </div>
                <div class="stat-card">
                    <div class="stat-value" id="st_req">-</div>
                    <div class="stat-label">총 요청</div>
                </div>
                <div class="stat-card">
                    <div class="stat-value" id="st_sent">-</div>
                    <div class="stat-label">보낸 데이터</div>
                </div>
                <div class="stat-card">
                    <div class="stat-value" id="st_recv">-</div>
                    <div class="stat-label">받은 데이터</div>
                </div>
            </div>
            <div style="text-align:right; margin-top:20px">
                <button class="btn" onclick="closeModal('statsModal')">닫기</button>
            </div>
        </div>
    </div>

    <div id="helpModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:500px;">
            <h3 style="margin-top:0;"><i class="fa-solid fa-book"></i> 사용 가이드</h3>
            <div style="line-height:1.7; color:var(--text)">
                <p><b>📁 파일/폴더 업로드</b><br>- 드래그 앤 드롭으로 <b>폴더째 업로드</b> 가능<br>- '업로드' 버튼으로 파일 여러 개 선택</p>
                <p><b>👁️ 미리보기 지원</b><br>- 이미지, 동영상, 오디오, <b>PDF</b>, 텍스트/코드</p>
                <p><b>✏️ 코드 뷰어</b><br>- 구문 강조 및 Markdown 미리보기</p>
                
                <hr style="border:none; border-top:1px solid var(--border); margin:16px 0;">
                
                <p style="margin-bottom:8px;"><b>⌨️ 키보드 단축키</b></p>
                <div style="display:grid; grid-template-columns:auto 1fr; gap:6px 16px; font-size:0.9rem;">
                    <kbd style="background:var(--hover); padding:4px 8px; border-radius:4px; border:1px solid var(--border);">Ctrl+U</kbd>
                    <span>파일 업로드</span>
                    <kbd style="background:var(--hover); padding:4px 8px; border-radius:4px; border:1px solid var(--border);">Ctrl+N</kbd>
                    <span>새 폴더 생성</span>
                    <kbd style="background:var(--hover); padding:4px 8px; border-radius:4px; border:1px solid var(--border);">Ctrl+A</kbd>
                    <span>전체 선택</span>
                    <kbd style="background:var(--hover); padding:4px 8px; border-radius:4px; border:1px solid var(--border);">Delete</kbd>
                    <span>선택 항목 삭제</span>
                    <kbd style="background:var(--hover); padding:4px 8px; border-radius:4px; border:1px solid var(--border);">F2</kbd>
                    <span>이름 변경</span>
                    <kbd style="background:var(--hover); padding:4px 8px; border-radius:4px; border:1px solid var(--border);">Escape</kbd>
                    <span>모달 닫기</span>
                </div>
            </div>
            <div style="text-align:right; margin-top:20px"><button class="btn" onclick="closeModal('helpModal')">닫기</button></div>
        </div>
    </div>

    <div id="clipModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal">
            <h3><i class="fa-regular fa-clipboard"></i> 공유 클립보드</h3>
            <label for="clipText" class="sr-only">클립보드 내용</label>
            <textarea id="clipText" style="width:100%; height:150px; padding:10px; border:1px solid var(--border); border-radius:8px; resize:none; background:var(--bg); color:var(--text); box-sizing:border-box;"></textarea>
            <div style="margin-top:10px; text-align:right; display:flex; gap:5px; justify-content:flex-end;">
                <button class="btn btn-outline" onclick="loadClipboard()">새로고침</button>
                <button class="btn" onclick="saveClipboard()">저장하기</button>
                <button class="btn btn-outline" onclick="closeModal('clipModal')">닫기</button>
            </div>
        </div>
    </div>

    <div id="editorModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal large">
            <h3 style="display:flex; justify-content:space-between; align-items:center; margin-top:0;">
                <span><i class="fa-solid fa-file-lines"></i> <span id="editorTitle"></span></span>
                <div style="display:flex; gap:10px; align-items:center">
                    <button id="previewToggle" class="btn-outline" style="font-size:0.8rem; padding:4px 8px; border-radius:4px; display:none" onclick="toggleMarkdownPreview()">미리보기</button>
                    <button class="btn-icon" style="border:none" onclick="closeModal('editorModal')" aria-label="닫기"><i class="fa-solid fa-xmark"></i></button>
                </div>
            </h3>
            <div class="editor-container">
                <textarea id="editorContent" class="editor-area" spellcheck="false" aria-label="코드 편집 영역"></textarea>
                <div id="codePreview" class="editor-area markdown-body" style="display:none; overflow-y:auto;" aria-label="미리보기 영역" tabindex="0"></div>
                <div id="mediaContainer" style="display:none; width:100%; height:100%; justify-content:center; align-items:center;"></div>
            </div>
            <div style="margin-top:15px; text-align:right;">
                <button class="btn btn-outline" onclick="closeModal('editorModal')">닫기</button>
                <button id="saveBtn" class="btn" onclick="saveFileContent()">저장</button>
            </div>
        </div>
    </div>

    <div id="mkdirModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal">
            <h3>새 폴더 생성</h3>
            <label for="newFolderInput" class="sr-only">폴더 이름</label>
            <input type="text" id="newFolderInput" placeholder="폴더 이름" style="width:100%; padding:10px; border:1px solid var(--border); border-radius:6px; box-sizing:border-box; background:var(--bg); color:var(--text);">
            <div style="margin-top:15px; text-align:right; gap:5px; display:flex; justify-content:flex-end">
                <button class="btn btn-outline" onclick="closeModal('mkdirModal')">취소</button>
                <button class="btn" onclick="createFolder()">생성</button>
            </div>
        </div>
    </div>
    
    <div id="progressModal" class="overlay" role="alertdialog" aria-modal="true">
        <div class="modal" style="text-align:center;">
            <h3><i class="fa-solid fa-cloud-arrow-up"></i> 업로드 중...</h3>
            <div id="progressFileInfo" style="font-size:0.9rem; margin-bottom:10px; color:var(--text); opacity:0.8;"></div>
            <div style="background:var(--border); height:8px; border-radius:4px; overflow:hidden; margin:15px 0;">
                <div id="progressBar" style="width:0%; height:100%; background:linear-gradient(90deg, var(--primary), #818cf8); transition:width 0.2s;" role="progressbar" aria-valuenow="0" aria-valuemin="0" aria-valuemax="100"></div>
            </div>
            <div id="progressText" style="font-size:1.2rem; font-weight:bold; color:var(--primary);">0%</div>
            <div id="progressStats" style="font-size:0.85rem; margin-top:10px; color:var(--text); opacity:0.7;"></div>
        </div>
    </div>

    <!-- 파일 정보 모달 -->
    <div id="fileInfoModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal">
            <h3><i class="fa-solid fa-circle-info"></i> 파일 정보</h3>
            <div id="fileInfoContent" style="line-height:1.8;"></div>
            <div style="margin-top:15px; text-align:right;">
                <button class="btn" onclick="closeModal('fileInfoModal')">닫기</button>
            </div>
        </div>
    </div>

    <!-- 북마크 모달 -->
    <div id="bookmarkModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal">
            <h3><i class="fa-solid fa-star"></i> 북마크</h3>
            <div id="bookmarkList" style="max-height:300px; overflow-y:auto;"></div>
            <div style="margin-top:15px; text-align:right;">
                <button class="btn" onclick="closeModal('bookmarkModal')">닫기</button>
            </div>
        </div>
    </div>

    <!-- 휴지통 모달 -->
    <div id="trashModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:500px;">
            <h3><i class="fa-solid fa-trash-can"></i> 휴지통</h3>
            <div id="trashList" style="max-height:300px; overflow-y:auto;"></div>
            <div style="margin-top:15px; text-align:right; display:flex; gap:5px; justify-content:flex-end;">
                <button class="btn btn-danger" onclick="emptyTrash()">휴지통 비우기</button>
                <button class="btn btn-outline" onclick="closeModal('trashModal')">닫기</button>
            </div>
        </div>
    </div>

    <!-- 공유 링크 목록 모달 -->
    <div id="shareListModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:600px;">
            <h3><i class="fa-solid fa-link"></i> 공유 링크 관리</h3>
            <div id="shareList" style="max-height:300px; overflow-y:auto;"></div>
            <div style="margin-top:15px; text-align:right;">
                <button class="btn" onclick="closeModal('shareListModal')">닫기</button>
            </div>
        </div>
    </div>

    <!-- 공유 링크 생성 모달 -->
    <div id="createShareModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal">
            <h3><i class="fa-solid fa-link"></i> 공유 링크 생성</h3>
            <p id="sharePathDisplay" style="word-break:break-all; color:var(--text); opacity:0.8;"></p>
            <label for="shareHours">유효 시간:</label>
            <select id="shareHours" style="width:100%; padding:8px; border:1px solid var(--border); border-radius:6px; background:var(--bg); color:var(--text); margin-top:5px;">
                <option value="1">1시간</option>
                <option value="6">6시간</option>
                <option value="24" selected>24시간</option>
                <option value="72">3일</option>
                <option value="168">7일</option>
            </select>
            <div id="generatedLink" style="margin-top:15px; display:none;">
                <label>생성된 링크:</label>
                <input type="text" id="shareLinkInput" readonly style="width:100%; padding:8px; border:1px solid var(--border); border-radius:6px; background:var(--bg); color:var(--text); margin-top:5px;">
                <button class="btn btn-outline" onclick="copyShareLink()" style="margin-top:10px;width:100%;"><i class="fa-solid fa-copy"></i> 복사</button>
            </div>
            <div style="margin-top:15px; text-align:right; display:flex; gap:5px; justify-content:flex-end;">
                <button class="btn btn-outline" onclick="closeModal('createShareModal')">취소</button>
                <button class="btn" id="createShareBtn" onclick="createShareLink()">생성</button>
            </div>
        </div>
    </div>

    <!-- v6.0: 비디오 플레이어 모달 -->
    <div id="videoPlayerModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal large" style="padding:0; background:#000;">
            <div style="position:relative; width:100%; height:100%;">
                <button onclick="closeModal('videoPlayerModal')" style="position:absolute; top:15px; right:15px; z-index:10; background:rgba(0,0,0,0.5); border:none; color:white; width:40px; height:40px; border-radius:50%; cursor:pointer; font-size:1.2rem;">
                    <i class="fa-solid fa-times"></i>
                </button>
                <video id="videoPlayer" controls style="width:100%; height:100%; object-fit:contain;">
                    Your browser does not support video playback.
                </video>
            </div>
        </div>
    </div>

    <!-- v6.0: 오디오 플레이어 모달 -->
    <div id="audioPlayerModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:500px;">
            <h3><i class="fa-solid fa-music"></i> 오디오 플레이어</h3>
            <div id="audioNowPlaying" style="text-align:center; margin:15px 0; font-weight:600; color:var(--primary);"></div>
            <audio id="audioPlayer" controls style="width:100%; margin-bottom:15px;"></audio>
            <div id="audioPlaylist" style="max-height:250px; overflow-y:auto; border:1px solid var(--border); border-radius:10px;"></div>
            <div style="margin-top:15px; display:flex; gap:8px; justify-content:center;">
                <button class="btn btn-outline btn-icon" onclick="audioPlayPrev()"><i class="fa-solid fa-backward-step"></i></button>
                <button class="btn btn-outline btn-icon" onclick="audioPlayNext()"><i class="fa-solid fa-forward-step"></i></button>
                <button class="btn" onclick="closeModal('audioPlayerModal')">닫기</button>
            </div>
        </div>
    </div>

    <!-- v6.0: 이미지 갤러리 모달 -->
    <div id="galleryModal" class="overlay" role="dialog" aria-modal="true" style="background:rgba(0,0,0,0.95);">
        <div style="position:relative; width:100%; height:100%; display:flex; align-items:center; justify-content:center;">
            <button onclick="closeModal('galleryModal')" style="position:absolute; top:20px; right:20px; z-index:10; background:rgba(255,255,255,0.1); border:none; color:white; width:50px; height:50px; border-radius:50%; cursor:pointer; font-size:1.5rem;">
                <i class="fa-solid fa-times"></i>
            </button>
            <button onclick="galleryPrev()" style="position:absolute; left:20px; top:50%; transform:translateY(-50%); background:rgba(255,255,255,0.1); border:none; color:white; width:60px; height:60px; border-radius:50%; cursor:pointer; font-size:1.5rem;">
                <i class="fa-solid fa-chevron-left"></i>
            </button>
            <img id="galleryImage" src="" style="max-width:90%; max-height:90%; object-fit:contain; border-radius:8px;">
            <button onclick="galleryNext()" style="position:absolute; right:20px; top:50%; transform:translateY(-50%); background:rgba(255,255,255,0.1); border:none; color:white; width:60px; height:60px; border-radius:50%; cursor:pointer; font-size:1.5rem;">
                <i class="fa-solid fa-chevron-right"></i>
            </button>
            <div id="galleryInfo" style="position:absolute; bottom:30px; color:white; text-align:center; font-size:0.9rem; opacity:0.8;"></div>
        </div>
    </div>

    <!-- v6.0: 사용자 관리 모달 -->
    {% if role == 'admin' %}
    <div id="userManageModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:600px;">
            <h3><i class="fa-solid fa-users-gear"></i> 사용자 관리</h3>
            <div style="margin-bottom:15px;">
                <button class="btn" onclick="showAddUserForm()"><i class="fa-solid fa-user-plus"></i> 새 사용자</button>
            </div>
            <div id="userFormArea" style="display:none; background:var(--hover); padding:15px; border-radius:10px; margin-bottom:15px;">
                <div style="display:grid; grid-template-columns:1fr 1fr; gap:10px;">
                    <input type="text" id="newUsername" placeholder="사용자명" style="padding:10px; border:1px solid var(--border); border-radius:8px; background:var(--card); color:var(--text);">
                    <input type="password" id="newPassword" placeholder="비밀번호" style="padding:10px; border:1px solid var(--border); border-radius:8px; background:var(--card); color:var(--text);">
                    <select id="newRole" style="padding:10px; border:1px solid var(--border); border-radius:8px; background:var(--card); color:var(--text);">
                        <option value="user">일반 사용자</option>
                        <option value="admin">관리자</option>
                    </select>
                    <input type="number" id="newQuota" placeholder="용량 제한 (MB)" value="1024" style="padding:10px; border:1px solid var(--border); border-radius:8px; background:var(--card); color:var(--text);">
                </div>
                <div style="margin-top:10px; display:flex; gap:8px; justify-content:flex-end;">
                    <button class="btn btn-outline" onclick="hideAddUserForm()">취소</button>
                    <button class="btn" onclick="createUser()">생성</button>
                </div>
            </div>
            <div id="userList" style="max-height:300px; overflow-y:auto;"></div>
            <div style="margin-top:15px; text-align:right;">
                <button class="btn" onclick="closeModal('userManageModal')">닫기</button>
            </div>
        </div>
    </div>
    {% endif %}

    <!-- v7.0: 태그 추가 모달 -->
    <div id="tagModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:400px;">
            <h3><i class="fa-solid fa-tag"></i> 태그 추가</h3>
            <p id="tagTargetPath" style="font-size:0.85rem; opacity:0.7; margin-bottom:15px;"></p>
            <div style="display:flex; gap:10px; margin-bottom:15px;">
                <input type="text" id="tagInput" placeholder="태그 이름" style="flex:1; padding:10px; border:1px solid var(--border); border-radius:8px; background:var(--bg); color:var(--text);">
                <input type="color" id="tagColor" value="#6366f1" style="width:50px; height:40px; border:none; border-radius:8px; cursor:pointer;">
            </div>
            <div id="existingTags" style="display:flex; flex-wrap:wrap; gap:6px; margin-bottom:15px;"></div>
            <div style="display:flex; gap:8px; justify-content:flex-end;">
                <button class="btn btn-outline" onclick="closeModal('tagModal')">취소</button>
                <button class="btn" onclick="addTag()">추가</button>
            </div>
        </div>
    </div>

    <!-- v7.0: 메모 모달 -->
    <div id="memoModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:500px;">
            <h3><i class="fa-solid fa-note-sticky"></i> 파일 메모</h3>
            <p id="memoTargetPath" style="font-size:0.85rem; opacity:0.7; margin-bottom:10px;"></p>
            <textarea id="memoText" placeholder="메모를 입력하세요..." style="width:100%; height:150px; padding:12px; border:1px solid var(--border); border-radius:10px; background:var(--bg); color:var(--text); resize:none; font-family:inherit;"></textarea>
            <p id="memoUpdated" style="font-size:0.75rem; opacity:0.5; margin-top:5px;"></p>
            <div style="margin-top:15px; display:flex; gap:8px; justify-content:flex-end;">
                <button class="btn btn-danger" onclick="deleteMemo()">삭제</button>
                <button class="btn btn-outline" onclick="closeModal('memoModal')">취소</button>
                <button class="btn" onclick="saveMemo()">저장</button>
            </div>
        </div>
    </div>

    <!-- v7.0: 암호화 모달 -->
    <div id="encryptModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:400px;">
            <h3><i class="fa-solid fa-lock"></i> 파일 암호화</h3>
            <p id="encryptTargetPath" style="font-size:0.85rem; opacity:0.7; margin-bottom:15px;"></p>
            <div style="margin-bottom:15px;">
                <label style="font-size:0.85rem; display:block; margin-bottom:5px;">암호화 비밀번호:</label>
                <input type="password" id="encryptPassword" placeholder="비밀번호 (기본: 관리자 암호)" style="width:100%; padding:10px; border:1px solid var(--border); border-radius:8px; background:var(--bg); color:var(--text); box-sizing:border-box;">
            </div>
            <p style="font-size:0.8rem; opacity:0.6; background:var(--hover); padding:10px; border-radius:8px;">
                <i class="fa-solid fa-info-circle"></i> 비밀번호를 잊으면 파일을 복구할 수 없습니다.
            </p>
            <div style="margin-top:15px; display:flex; gap:8px; justify-content:flex-end;">
                <button class="btn btn-outline" onclick="closeModal('encryptModal')">취소</button>
                <button class="btn" onclick="encryptFile()">암호화</button>
            </div>
        </div>
    </div>

    <!-- v7.0: 복호화 모달 -->
    <div id="decryptModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:400px;">
            <h3><i class="fa-solid fa-unlock"></i> 파일 복호화</h3>
            <p id="decryptTargetPath" style="font-size:0.85rem; opacity:0.7; margin-bottom:15px;"></p>
            <div style="margin-bottom:15px;">
                <label style="font-size:0.85rem; display:block; margin-bottom:5px;">복호화 비밀번호:</label>
                <input type="password" id="decryptPassword" placeholder="암호화 시 사용한 비밀번호" style="width:100%; padding:10px; border:1px solid var(--border); border-radius:8px; background:var(--bg); color:var(--text); box-sizing:border-box;">
            </div>
            <div style="margin-top:15px; display:flex; gap:8px; justify-content:flex-end;">
                <button class="btn btn-outline" onclick="closeModal('decryptModal')">취소</button>
                <button class="btn" onclick="decryptFile()">복호화</button>
            </div>
        </div>
    </div>

    <!-- v7.0: 즐겨찾기 모달 -->
    <div id="favoritesModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal" style="max-width:500px;">
            <h3><i class="fa-solid fa-folder-heart"></i> 즐겨찾기 폴더</h3>
            <div id="favoritesList" style="max-height:350px; overflow-y:auto;"></div>
            <div style="margin-top:15px; text-align:right;">
                <button class="btn" onclick="closeModal('favoritesModal')">닫기</button>
            </div>
        </div>
    </div>

    <!-- v7.0: 접속 로그 대시보드 모달 -->
    <div id="accessDashboardModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal large" style="max-width:800px;">
            <h3><i class="fa-solid fa-chart-line"></i> 접속 대시보드</h3>
            <div style="display:grid; grid-template-columns:1fr 1fr; gap:20px; margin-bottom:20px;">
                <div style="background:var(--bg); padding:15px; border-radius:12px; border:1px solid var(--border);">
                    <h4 style="margin:0 0 10px 0; font-size:0.9rem; opacity:0.7;">활동별 통계</h4>
                    <div id="actionStats"></div>
                </div>
                <div style="background:var(--bg); padding:15px; border-radius:12px; border:1px solid var(--border);">
                    <h4 style="margin:0 0 10px 0; font-size:0.9rem; opacity:0.7;">차단된 IP</h4>
                    <div id="blockedIpsList"></div>
                </div>
            </div>
            <div style="background:var(--bg); padding:15px; border-radius:12px; border:1px solid var(--border);">
                <h4 style="margin:0 0 10px 0; font-size:0.9rem; opacity:0.7;">최근 접속 기록</h4>
                <div id="recentLogs" style="max-height:200px; overflow-y:auto;"></div>
            </div>
            <div style="margin-top:15px; display:flex; gap:8px; justify-content:flex-end;">
                <button class="btn btn-outline" onclick="loadAccessDashboard()">새로고침</button>
                <button class="btn" onclick="closeModal('accessDashboardModal')">닫기</button>
            </div>
        </div>
    </div>

    <!-- v7.0: 문서 미리보기 모달 -->
    <div id="docPreviewModal" class="overlay" role="dialog" aria-modal="true">
        <div class="modal large" style="max-width:900px; max-height:90vh;">
            <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:15px;">
                <h3 id="docPreviewTitle" style="margin:0;"><i class="fa-solid fa-file-alt"></i> 문서 미리보기</h3>
                <button class="btn btn-outline" onclick="closeModal('docPreviewModal')" style="font-size:1.2rem; padding:5px 12px;">&times;</button>
            </div>
            <div id="docPreviewContent" style="max-height:calc(90vh - 150px); overflow-y:auto; background:var(--bg); padding:20px; border-radius:12px; border:1px solid var(--border);"></div>
            <div style="margin-top:15px; display:flex; gap:8px; justify-content:flex-end;">
                <button class="btn btn-outline" onclick="downloadCurrentDoc()"><i class="fa-solid fa-download"></i> 다운로드</button>
                <button class="btn" onclick="closeModal('docPreviewModal')">닫기</button>
            </div>
        </div>
    </div>

    <script>
        const currentPath = "{{ current_path }}";
        const canModify = {{ 'true' if can_modify else 'false' }};
        let selectedFiles = new Set();
        
        // Utility: Format bytes to human readable size
        function formatSize(bytes) {
            if (bytes < 1024) return bytes + ' B';
            if (bytes < 1024 * 1024) return (bytes / 1024).toFixed(1) + ' KB';
            if (bytes < 1024 * 1024 * 1024) return (bytes / 1024 / 1024).toFixed(1) + ' MB';
            return (bytes / 1024 / 1024 / 1024).toFixed(2) + ' GB';
        }
        
        // Utility: Escape HTML to prevent XSS
        function escapeHtml(text) {
            const div = document.createElement('div');
            div.textContent = text;
            return div.innerHTML;
        }
        
        // v5: 키보드 탐색 인덱스
        let currentKeyboardIndex = -1;
        
        // v5: 파일 목록 키보드 탐색
        function navigateFileList(direction) {
            const items = document.querySelectorAll('.file-item.data-item');
            if(items.length === 0) return;
            
            // 이전 포커스 제거
            if(currentKeyboardIndex >= 0 && items[currentKeyboardIndex]) {
                items[currentKeyboardIndex].classList.remove('keyboard-focused');
            }
            
            // 새 인덱스 계산
            currentKeyboardIndex += direction;
            if(currentKeyboardIndex < 0) currentKeyboardIndex = items.length - 1;
            if(currentKeyboardIndex >= items.length) currentKeyboardIndex = 0;
            
            // 새 포커스 적용
            items[currentKeyboardIndex].classList.add('keyboard-focused');
            items[currentKeyboardIndex].scrollIntoView({ block: 'nearest', behavior: 'smooth' });
        }
        
        // Utility: Toggle password visibility
        function togglePasswordVisibility() {
            const pwInput = document.getElementById('password');
            const icon = document.getElementById('pwToggleIcon');
            if (pwInput.type === 'password') {
                pwInput.type = 'text';
                icon.className = 'fa-solid fa-eye-slash';
            } else {
                pwInput.type = 'password';
                icon.className = 'fa-solid fa-eye';
            }
        }
        
        document.addEventListener('DOMContentLoaded', () => {
             // v7.1: Global Fetch Interceptor for CSRF
            const originalFetch = window.fetch;
            window.fetch = function(url, options) {
                if (options && options.method && ['POST', 'PUT', 'DELETE'].includes(options.method.toUpperCase())) {
                    options.headers = options.headers || {};
                    const token = document.querySelector('meta[name="csrf-token"]').content;
                    if (options.headers instanceof Headers) {
                        options.headers.append('X-CSRF-Token', token);
                    } else {
                        options.headers['X-CSRF-Token'] = token;
                    }
                }
                return originalFetch(url, options);
            };

            fetchDiskInfo();
            document.addEventListener('keydown', (e) => {
                // Escape: 모든 모달 닫기
                if(e.key === "Escape") {
                    document.querySelectorAll('.overlay').forEach(el => el.style.display = 'none');
                }
                
                // 입력 필드에 포커스 중이면 단축키 무시
                if(e.target.tagName === 'INPUT' || e.target.tagName === 'TEXTAREA') return;
                
                // Ctrl+U: 업로드
                if(e.ctrlKey && e.key === 'u' && canModify) {
                    e.preventDefault();
                    document.getElementById('fileInput').click();
                }
                
                // Ctrl+N: 새 폴더
                if(e.ctrlKey && e.key === 'n' && canModify) {
                    e.preventDefault();
                    openModal('mkdirModal');
                    document.getElementById('newFolderInput').focus();
                }
                
                // Delete: 선택된 파일 삭제
                if(e.key === 'Delete' && selectedFiles.size > 0 && canModify) {
                    e.preventDefault();
                    batchDelete();
                }
                
                // Ctrl+A: 모든 파일 선택
                if(e.ctrlKey && e.key === 'a') {
                    e.preventDefault();
                    document.querySelectorAll('.file-check').forEach(c => {
                        if(!c.checked) {
                            c.checked = true;
                            toggleBatch(c);
                        }
                    });
                }
                
                // F2: 선택된 항목 이름 변경
                if(e.key === 'F2' && selectedFiles.size === 1) {
                    e.preventDefault();
                    const fileName = Array.from(selectedFiles)[0];
                    const newName = prompt("새 이름:", fileName);
                    if(newName && newName !== fileName) {
                        fetch('/rename/' + currentPath, {
                            method:'POST', 
                            headers:{'Content-Type':'application/json'}, 
                            body:JSON.stringify({old_name: fileName, new_name: newName})
                        }).then(r=>r.json()).then(d => { 
                            if(d.success) location.reload(); 
                            else showToast(d.error, 'error'); 
                        });
                    }
                }
                
                // v5: 방향키로 파일 목록 탐색
                if(e.key === 'ArrowDown' || e.key === 'ArrowUp') {
                    e.preventDefault();
                    navigateFileList(e.key === 'ArrowDown' ? 1 : -1);
                }
                
                // v5: Enter로 선택된 항목 열기
                if(e.key === 'Enter' && currentKeyboardIndex >= 0) {
                    e.preventDefault();
                    const items = document.querySelectorAll('.file-item.data-item');
                    if(items[currentKeyboardIndex]) {
                        items[currentKeyboardIndex].querySelector('.file-info').click();
                    }
                }
            });
            
            // 단축키 힌트 표시
            console.log('📌 키보드 단축키: Ctrl+U(업로드), Ctrl+N(새폴더), Delete(삭제), Ctrl+A(전체선택), F2(이름변경)');
        });
        
        // v7.0: 드롭다운 메뉴 토글
        function toggleDropdown(id) {
            const dropdown = document.getElementById(id);
            const isOpen = dropdown.classList.contains('open');
            closeDropdowns();
            if (!isOpen) dropdown.classList.add('open');
        }
        
        function closeDropdowns() {
            document.querySelectorAll('.dropdown.open').forEach(d => d.classList.remove('open'));
        }
        
        // 드롭다운 외부 클릭 시 닫기
        document.addEventListener('click', e => {
            if (!e.target.closest('.dropdown')) closeDropdowns();
        });

        function fetchStats() {
            fetch('/metrics').then(r=>r.json()).then(d => {
                document.getElementById('st_uptime').innerText = d.uptime;
                document.getElementById('st_req').innerText = d.requests.toLocaleString();
                document.getElementById('st_sent').innerText = d.sent;
                document.getElementById('st_recv').innerText = d.recv;
            });
        }

        function toggleBatch(checkbox) {
            const row = checkbox.closest('.file-item');
            if (checkbox.checked) {
                selectedFiles.add(checkbox.value);
                row.classList.add('selected');
            } else {
                selectedFiles.delete(checkbox.value);
                row.classList.remove('selected');
            }
            updateBatchUI();
        }

        function updateBatchUI() {
            const bar = document.getElementById('batchBar');
            const cnt = document.getElementById('batchCount');
            if (selectedFiles.size > 0) {
                bar.style.display = 'flex';
                cnt.innerText = selectedFiles.size + '개 선택됨';
            } else {
                bar.style.display = 'none';
            }
        }

        function batchDownload() {
            if (selectedFiles.size === 0) return;
            const files = Array.from(selectedFiles);
            const form = document.createElement('form');
            form.method = 'POST';
            form.action = '/batch_download/' + currentPath;
            const input = document.createElement('input');
            input.type = 'hidden';
            input.name = 'files';
            input.value = JSON.stringify(files);
            form.appendChild(input);
            document.body.appendChild(form);
            form.submit();
            document.body.removeChild(form);
            document.querySelectorAll('.file-check').forEach(c => { c.checked = false; toggleBatch(c); });
        }

        function batchDelete() {
            if (selectedFiles.size === 0) return;
            if (!confirm(selectedFiles.size + "개 항목을 삭제하시겠습니까?")) return;
            fetch('/batch_delete/' + currentPath, {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({files: Array.from(selectedFiles)})
            }).then(r => r.json()).then(d => {
                if (d.success) location.reload(); else alert(d.error);
            });
        }

        let editPath = '';
        let isMarkdown = false;
        
        function handleItemClick(path, type, isDir, ext) {
            if(isDir) location.href = '/browse/' + path;
            // v6.0: 비디오는 스트리밍 플레이어로 재생
            else if (type === 'video') {
                playVideo(path);
            }
            // v6.0: 오디오는 현재 폴더 플레이리스트로 재생 (단일 파일은 기존 방식)
            else if (type === 'audio') {
                openEditor(path, path.split('/').pop(), ext, true);
            }
            else if (type === 'image' || ext.toLowerCase() === '.pdf') {
                openEditor(path, path.split('/').pop(), ext, true);
            }
            else if (type === 'text') {
                openEditor(path, path.split('/').pop(), ext, false);
            }
            // v7.0: 문서 미리보기 (Word, Excel, PowerPoint, CSV, JSON)
            else if (['.docx', '.xlsx', '.xls', '.pptx', '.csv', '.json'].includes(ext.toLowerCase())) {
                openDocumentPreview(path, path.split('/').pop());
            }
            else location.href = '/download/' + path;
        }

        function openEditor(path, name, ext, readOnly) {
            editPath = path;
            ext = ext.toLowerCase();
            isMarkdown = ext === '.md';
            
            document.getElementById('editorTitle').innerText = name;
            
            const editor = document.getElementById('editorContent');
            const preview = document.getElementById('codePreview');
            const media = document.getElementById('mediaContainer');
            const saveBtn = document.getElementById('saveBtn');
            const toggleBtn = document.getElementById('previewToggle');
            
            editor.style.display = 'none';
            preview.style.display = 'none';
            media.style.display = 'none';
            toggleBtn.style.display = 'none';
            saveBtn.style.display = readOnly ? 'none' : 'inline-block';

            if (readOnly) {
                media.style.display = 'flex';
                const url = '/download/' + path;
                if(ext === '.pdf') {
                    media.innerHTML = `<iframe src="${url}" style="width:100%; height:100%; border:none;" title="PDF Preview"></iframe>`;
                } else if (['.mp4', '.webm', '.ogg'].includes(ext)) {
                    media.innerHTML = `<video controls autoplay style="max-width:100%; max-height:100%"><source src="${url}"></video>`;
                } else if (['.mp3', '.wav', '.ogg'].includes(ext)) {
                    media.innerHTML = `<audio controls autoplay><source src="${url}"></audio>`;
                } else if (['.jpg', '.jpeg', '.png', '.gif', '.webp'].includes(ext)) {
                    media.innerHTML = `<img src="${url}" style="max-width:100%; max-height:100%; object-fit:contain">`;
                }
                openModal('editorModal');
            } else {
                editor.style.display = 'block';
                toggleBtn.style.display = 'inline-block';
                toggleBtn.innerText = "미리보기";
                editor.value = "Loading...";
                
                fetch('/get_content/' + path).then(r=>r.json()).then(d => {
                    if(d.error) { alert(d.error); return; }
                    editor.value = d.content;
                    openModal('editorModal');
                });
            }
        }

        function toggleMarkdownPreview() {
            const editor = document.getElementById('editorContent');
            const preview = document.getElementById('codePreview');
            const btn = document.getElementById('previewToggle');
            
            if(editor.style.display !== 'none') {
                if(isMarkdown) {
                    preview.innerHTML = marked.parse(editor.value);
                } else {
                    const safeContent = editor.value.replace(/&/g, "&amp;").replace(/</g, "&lt;").replace(/>/g, "&gt;");
                    const ext = editPath.split('.').pop();
                    preview.innerHTML = `<pre><code class="language-${ext}">${safeContent}</code></pre>`;
                    hljs.highlightElement(preview.querySelector('code'));
                }
                editor.style.display = 'none';
                preview.style.display = 'block';
                btn.innerText = "편집하기";
            } else {
                preview.style.display = 'none';
                editor.style.display = 'block';
                btn.innerText = "미리보기";
            }
        }

        function saveFileContent() {
            const content = document.getElementById('editorContent').value;
            fetch('/save_content/' + editPath, {
                method:'POST', headers:{'Content-Type':'application/json'},
                body:JSON.stringify({content: content})
            }).then(r=>r.json()).then(d => {
                if(d.success) { showToast('저장되었습니다.'); closeModal('editorModal'); }
                else alert(d.error);
            });
        }

        function handleFileSelect(files) { if(files.length > 0) uploadFiles(files); }

        function uploadFiles(files) {
            openModal('progressModal');
            const fd = new FormData();
            let totalSize = 0;
            
            for(let i=0; i<files.length; i++) {
                const file = files[i];
                const path = file.webkitRelativePath || file.name;
                fd.append('file', file);
                fd.append('paths', path);
                totalSize += file.size;
            }
            
            // 파일 정보 표시 (전역 formatSize 함수 사용)
            document.getElementById('progressFileInfo').innerText = 
                `${files.length}개 파일 (${formatSize(totalSize)})`;
            
            const xhr = new XMLHttpRequest();
            const startTime = Date.now();
            
            xhr.open('POST', '/upload/' + currentPath);
            // v7.1: CSRF Token for XHR
            xhr.setRequestHeader('X-CSRF-Token', document.querySelector('meta[name="csrf-token"]').content);
            
            xhr.upload.onprogress = e => {
                if(e.lengthComputable) {
                    const p = Math.round((e.loaded/e.total)*100);
                    const elapsed = (Date.now() - startTime) / 1000;
                    const speed = e.loaded / elapsed;
                    const remaining = (e.total - e.loaded) / speed;
                    
                    document.getElementById('progressBar').style.width = p+'%';
                    document.getElementById('progressBar').setAttribute('aria-valuenow', p);
                    document.getElementById('progressText').innerText = p+'%';
                    
                    // 속도와 예상 시간 표시
                    const speedStr = formatSize(speed) + '/s';
                    const remainStr = remaining > 60 
                        ? Math.ceil(remaining / 60) + '분 남음'
                        : Math.ceil(remaining) + '초 남음';
                    document.getElementById('progressStats').innerText = 
                        `${speedStr} • ${formatSize(e.loaded)} / ${formatSize(e.total)} • ${remainStr}`;
                }
            };
            xhr.onload = () => {
                showToast('업로드 완료!', 'success');
                setTimeout(() => location.reload(), 500);
            };
            xhr.onerror = () => { 
                showToast('업로드 실패', 'error'); 
                closeModal('progressModal');
            };
            xhr.send(fd);
        }

        const dropZone = document.getElementById('drop-zone');
        window.addEventListener('dragenter', e => { if(canModify && e.dataTransfer.types.includes('Files')) dropZone.style.display='flex'; });
        dropZone.addEventListener('dragleave', e => dropZone.style.display='none');
        dropZone.addEventListener('drop', e => {
            e.preventDefault(); dropZone.style.display='none';
            if(canModify) uploadFiles(e.dataTransfer.files);
        });
        window.addEventListener('dragover', e => e.preventDefault());

        function showToast(msg, type = 'info') {
            const icons = {
                success: '<i class="fa-solid fa-check-circle"></i>',
                error: '<i class="fa-solid fa-exclamation-circle"></i>',
                warning: '<i class="fa-solid fa-exclamation-triangle"></i>',
                info: '<i class="fa-solid fa-info-circle"></i>'
            };
            const t = document.createElement('div'); 
            t.className = 'toast ' + type; 
            t.innerHTML = (icons[type] || '') + ' ' + msg; 
            t.setAttribute('role', 'alert');
            document.getElementById('toast-container').appendChild(t);
            setTimeout(() => t.remove(), 3500);
        }
        function sortFiles() {
            const list = document.getElementById('fileList');
            const items = Array.from(list.querySelectorAll('.data-item'));
            const type = document.getElementById('sortOrder').value;
            const parent = list.querySelector('.parent-folder');
            items.sort((a, b) => {
                const isDirA = a.querySelector('.file-icon').classList.contains('folder');
                const isDirB = b.querySelector('.file-icon').classList.contains('folder');
                if (isDirA !== isDirB) return isDirA ? -1 : 1;
                if (type === 'name') return a.getAttribute('data-name').localeCompare(b.getAttribute('data-name'));
                if (type === 'size') return parseInt(b.getAttribute('data-size')) - parseInt(a.getAttribute('data-size'));
                if (type === 'date') return parseFloat(b.getAttribute('data-date')) - parseFloat(a.getAttribute('data-date'));
                return 0;
            });
            list.innerHTML = '';
            if(parent) list.appendChild(parent);
            items.forEach(item => list.appendChild(item));
        }
        function fetchDiskInfo() {
            fetch('/disk_info').then(r=>r.json()).then(d => {
                if(d.error) return;
                document.getElementById('diskText').innerText = `${d.used} / ${d.total} (${d.percent}%)`;
                document.getElementById('diskFill').style.width = d.percent + '%';
            });
        }
        function downloadItem(path) { location.href = '/download/' + path; }
        function deleteItem(path) {
            if(!confirm('정말 삭제하시겠습니까?')) return;
            fetch('/delete/' + path, {method:'POST'})
                .then(r=>r.json()).then(d=>{ if(d.success) location.reload(); else alert(d.error); });
        }
        function createFolder() {
            const name = document.getElementById('newFolderInput').value;
            if(!name) return;
            fetch('/mkdir/' + currentPath, {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({name: name})})
            .then(r=>r.json()).then(d => { if(d.success) location.reload(); else alert(d.error); });
        }
        
        let ctxTarget = null;
        document.addEventListener('click', () => document.getElementById('ctxMenu').style.display='none');
        function openCtx(e, path, name, type) {
            e.preventDefault();
            ctxTarget = {path, name, type};
            const unzipBtn = document.getElementById('ctxUnzip');
            if(unzipBtn) unzipBtn.style.display = (type === 'archive') ? 'flex' : 'none';
            
            // v7.0: 암호화/복호화 버튼 표시 로직
            const encryptBtn = document.getElementById('ctxEncrypt');
            const decryptBtn = document.getElementById('ctxDecrypt');
            const favoriteBtn = document.getElementById('ctxFavorite');
            
            if(encryptBtn && decryptBtn) {
                const isEncrypted = name.endsWith('.enc');
                const isFolder = (type === 'folder');
                encryptBtn.style.display = (!isEncrypted && !isFolder) ? 'flex' : 'none';
                decryptBtn.style.display = isEncrypted ? 'flex' : 'none';
            }
            
            // v7.0: 즐겨찾기는 폴더만
            if(favoriteBtn) {
                favoriteBtn.style.display = (type === 'folder') ? 'flex' : 'none';
            }
            
            const menu = document.getElementById('ctxMenu');
            menu.style.display = 'block';
            menu.style.left = e.pageX + 'px';
            menu.style.top = e.pageY + 'px';
        }
        function handleCtx(action) {
            if(!ctxTarget) return;
            if(action === 'download') downloadItem(ctxTarget.path);
            if(action === 'delete') {
                if(!confirm('영구적으로 삭제하시겠습니까? (복구 불가)')) return;
                deleteItem(ctxTarget.path);
            }
            if(action === 'unzip') {
                if(!confirm('압축 해제?')) return;
                fetch('/unzip/' + ctxTarget.path, {method:'POST'}).then(r=>r.json()).then(d=>{ 
                    if(d.success) { showToast('압축 해제 완료', 'success'); location.reload(); }
                    else showToast(d.error, 'error'); 
                });
            }
            if(action === 'rename') {
                const newName = prompt("새 이름:", ctxTarget.name);
                if(newName && newName !== ctxTarget.name) {
                    fetch('/rename/' + currentPath, {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({old_name: ctxTarget.name, new_name: newName})})
                    .then(r=>r.json()).then(d=>{ if(d.success) location.reload(); else showToast(d.error, 'error'); });
                }
            }
            if(action === 'info') {
                showFileInfo(ctxTarget.path);
            }
            if(action === 'bookmark') {
                addBookmark(ctxTarget.path, ctxTarget.name);
            }
            if(action === 'share') {
                openShareModal(ctxTarget.path);
            }
            if(action === 'trash') {
                if(!confirm('휴지통으로 이동하시겠습니까?')) return;
                fetch('/trash', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({path: ctxTarget.path})})
                .then(r=>r.json()).then(d=>{ 
                    if(d.success) { showToast('휴지통으로 이동됨', 'success'); location.reload(); }
                    else showToast(d.error, 'error'); 
                });
            }
            // v7.0: 태그 추가
            if(action === 'tag') {
                openTagModal(ctxTarget.path, ctxTarget.name);
            }
            // v7.0: 메모
            if(action === 'memo') {
                openMemoModal(ctxTarget.path, ctxTarget.name);
            }
            // v7.0: 즐겨찾기
            if(action === 'favorite') {
                addFavorite(ctxTarget.path, ctxTarget.name);
            }
            // v7.0: 암호화
            if(action === 'encrypt') {
                document.getElementById('encryptTargetPath').textContent = '대상: ' + ctxTarget.name;
                document.getElementById('encryptPassword').value = '';
                openModal('encryptModal');
            }
            // v7.0: 복호화
            if(action === 'decrypt') {
                document.getElementById('decryptTargetPath').textContent = '대상: ' + ctxTarget.name;
                document.getElementById('decryptPassword').value = '';
                openModal('decryptModal');
            }
        }
        
        // 파일 정보 표시
        function showFileInfo(path) {
            fetch('/file_info/' + path).then(r=>r.json()).then(d => {
                if(d.error) { showToast(d.error, 'error'); return; }
                
                // 전역 formatSize 함수 사용
                let html = `
                    <p><strong>이름:</strong> ${d.name}</p>
                    <p><strong>경로:</strong> ${d.path}</p>
                    <p><strong>타입:</strong> ${d.is_dir ? '폴더' : '파일'}</p>
                    <p><strong>크기:</strong> ${formatSize(d.size)}</p>
                    <p><strong>생성:</strong> ${new Date(d.created).toLocaleString()}</p>
                    <p><strong>수정:</strong> ${new Date(d.modified).toLocaleString()}</p>
                `;
                
                if(!d.is_dir) {
                    html += `<p><strong>MIME:</strong> ${d.mime_type || '-'}</p>`;
                    if(d.md5) html += `<p><strong>MD5:</strong> <code style="font-size:0.8rem;">${d.md5}</code></p>`;
                } else {
                    html += `<p><strong>파일:</strong> ${d.file_count || 0}개</p>`;
                    html += `<p><strong>폴더:</strong> ${d.folder_count || 0}개</p>`;
                }
                
                document.getElementById('fileInfoContent').innerHTML = html;
                openModal('fileInfoModal');
            });
        }
        
        // 북마크 관련
        function addBookmark(path, name) {
            fetch('/bookmarks', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({path, name})})
            .then(r=>r.json()).then(d => {
                if(d.success) showToast('북마크 추가됨', 'success');
                else showToast(d.error, 'warning');
            });
        }
        
        function loadBookmarks() {
            fetch('/bookmarks').then(r=>r.json()).then(d => {
                const list = document.getElementById('bookmarkList');
                if(!d.bookmarks || d.bookmarks.length === 0) {
                    list.innerHTML = '<p style="text-align:center; opacity:0.6;">북마크가 없습니다.</p>';
                    return;
                }
                list.innerHTML = d.bookmarks.map(b => `
                    <div style="display:flex; align-items:center; padding:8px; border-bottom:1px solid var(--border);">
                        <i class="fa-solid fa-star" style="color:var(--folder); margin-right:10px;"></i>
                        <a href="/browse/${escapeHtml(b.path)}" style="flex:1; color:var(--text); text-decoration:none;">${escapeHtml(b.name)}</a>
                        <button class="btn-icon btn-danger" onclick="removeBookmark('${escapeHtml(b.path)}')" style="border:none;background:transparent;"><i class="fa-solid fa-xmark"></i></button>
                    </div>
                `).join('');
            });
        }
        
        function removeBookmark(path) {
            fetch('/bookmarks', {method:'DELETE', headers:{'Content-Type':'application/json'}, body:JSON.stringify({path})})
            .then(r=>r.json()).then(d => {
                if(d.success) { showToast('북마크 삭제됨', 'success'); loadBookmarks(); }
            });
        }
        
        // 휴지통 관련
        function loadTrash() {
            fetch('/trash/list').then(r=>r.json()).then(d => {
                const list = document.getElementById('trashList');
                if(!d.items || d.items.length === 0) {
                    list.innerHTML = '<p style="text-align:center; opacity:0.6;">휴지통이 비어있습니다.</p>';
                    return;
                }
                list.innerHTML = d.items.map(item => `
                    <div style="display:flex; align-items:center; padding:8px; border-bottom:1px solid var(--border);">
                        <i class="fa-solid ${item.is_dir ? 'fa-folder' : 'fa-file'}" style="margin-right:10px; color:var(--text); opacity:0.5;"></i>
                        <div style="flex:1;">
                            <div>${escapeHtml(item.original_name)}</div>
                            <div style="font-size:0.75rem; opacity:0.6;">${new Date(item.deleted_at).toLocaleString()}</div>
                        </div>
                        <button class="btn btn-outline" style="font-size:0.75rem; padding:4px 8px;" onclick="restoreFromTrash('${escapeHtml(item.name)}')"><i class="fa-solid fa-undo"></i></button>
                    </div>
                `).join('');
            });
        }
        
        function restoreFromTrash(name) {
            fetch('/trash/restore', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({name})})
            .then(r=>r.json()).then(d => {
                if(d.success) { showToast('복원됨', 'success'); loadTrash(); }
                else showToast(d.error, 'error');
            });
        }
        
        function emptyTrash() {
            if(!confirm('휴지통을 비우시겠습니까? (모든 항목 영구 삭제)')) return;
            fetch('/trash/empty', {method:'POST'}).then(r=>r.json()).then(d => {
                if(d.success) { showToast('휴지통 비움', 'success'); loadTrash(); }
                else showToast(d.error, 'error');
            });
        }
        
        // 공유 링크 관련
        let currentSharePath = '';
        
        function openShareModal(path) {
            currentSharePath = path;
            document.getElementById('sharePathDisplay').innerText = '대상: ' + path;
            document.getElementById('generatedLink').style.display = 'none';
            document.getElementById('createShareBtn').disabled = false;
            openModal('createShareModal');
        }
        
        function createShareLink() {
            const hours = parseInt(document.getElementById('shareHours').value);
            fetch('/share/create', {
                method:'POST', 
                headers:{'Content-Type':'application/json'}, 
                body:JSON.stringify({path: currentSharePath, hours})
            }).then(r=>r.json()).then(d => {
                if(d.success) {
                    const fullLink = window.location.origin + d.link;
                    document.getElementById('shareLinkInput').value = fullLink;
                    document.getElementById('generatedLink').style.display = 'block';
                    document.getElementById('createShareBtn').disabled = true;
                    showToast('공유 링크 생성됨', 'success');
                } else {
                    showToast(d.error, 'error');
                }
            });
        }
        
        function copyShareLink() {
            const input = document.getElementById('shareLinkInput');
            input.select();
            document.execCommand('copy');
            showToast('클립보드에 복사되었습니다', 'success');
        }
        
        function loadShareLinks() {
            fetch('/share/list').then(r=>r.json()).then(d => {
                const list = document.getElementById('shareList');
                if(!d.links || d.links.length === 0) {
                    list.innerHTML = '<p style="text-align:center; opacity:0.6;">활성 공유 링크가 없습니다.</p>';
                    return;
                }
                list.innerHTML = d.links.map(link => `
                    <div style="display:flex; align-items:center; padding:8px; border-bottom:1px solid var(--border);">
                        <i class="fa-solid fa-link" style="margin-right:10px; color:var(--primary);"></i>
                        <div style="flex:1; min-width:0;">
                            <div style="overflow:hidden; text-overflow:ellipsis; white-space:nowrap;">${escapeHtml(link.path)}</div>
                            <div style="font-size:0.75rem; opacity:0.6;">만료: ${new Date(link.expires).toLocaleString()}</div>
                        </div>
                        <button class="btn btn-outline" style="font-size:0.75rem; padding:4px 8px; margin-right:5px;" onclick="navigator.clipboard.writeText(window.location.origin + '/share/${escapeHtml(link.token)}'); showToast('복사됨','success');"><i class="fa-solid fa-copy"></i></button>
                        <button class="btn-icon btn-danger" style="border:none;background:transparent;" onclick="deleteShareLink('${escapeHtml(link.token)}')"><i class="fa-solid fa-xmark"></i></button>
                    </div>
                `).join('');
            });
        }
        
        function deleteShareLink(token) {
            fetch('/share/delete/' + token, {method:'POST'}).then(r=>r.json()).then(d => {
                if(d.success) { showToast('링크 삭제됨', 'success'); loadShareLinks(); }
            });
        }
        
        function loadClipboard() { fetch('/clipboard').then(r=>r.json()).then(d => document.getElementById('clipText').value = d.content); }
        function saveClipboard() { fetch('/clipboard', {method:'POST', headers:{'Content-Type':'application/json'}, body:JSON.stringify({content: document.getElementById('clipText').value})}).then(()=> { showToast('저장됨', 'success'); closeModal('clipModal'); }); }
        function toggleTheme() {
            const html = document.documentElement;
            const isDark = html.getAttribute('data-theme') === 'dark';
            html.setAttribute('data-theme', isDark ? 'light' : 'dark');
            localStorage.setItem('theme', isDark ? 'light' : 'dark');
        }
        function toggleView() {
            const list = document.getElementById('fileList');
            const icon = document.getElementById('viewIcon');
            if(list.parentElement.classList.contains('grid-view')) {
                list.parentElement.classList.remove('grid-view');
                icon.className = 'fa-solid fa-list';
                localStorage.setItem('view', 'list');
            } else {
                list.parentElement.classList.add('grid-view');
                icon.className = 'fa-solid fa-border-all';
                localStorage.setItem('view', 'grid');
                document.querySelectorAll('.preview').forEach(img => img.style.display = 'block');
            }
        }
        function filterFiles() {
            const q = document.getElementById('searchInput').value.toLowerCase();
            document.querySelectorAll('.data-item').forEach(item => {
                const name = item.getAttribute('data-name').toLowerCase();
                item.style.display = name.includes(q) ? 'flex' : 'none';
            });
        }
        function openModal(id) { 
            document.getElementById(id).style.display = 'flex'; 
        }
        function closeModal(id) { document.getElementById(id).style.display = 'none'; }
        
        const savedTheme = localStorage.getItem('theme');
        if(savedTheme) document.documentElement.setAttribute('data-theme', savedTheme);
        else if (window.matchMedia && window.matchMedia('(prefers-color-scheme: dark)').matches) document.documentElement.setAttribute('data-theme', 'dark');
        if(localStorage.getItem('view') === 'grid') toggleView();
        
        // v5.1: 언어 전환
        function toggleLanguage() {
            const currentLang = localStorage.getItem('lang') || 'ko';
            const newLang = currentLang === 'ko' ? 'en' : 'ko';
            fetch('/set_language/' + newLang).then(r => r.json()).then(d => {
                if(d.success) {
                    localStorage.setItem('lang', newLang);
                    showToast(newLang === 'ko' ? '한국어로 변경됨' : 'Changed to English', 'success');
                    setTimeout(() => location.reload(), 500);
                }
            });
        }
        
        // v5.1: 최근 파일 로드
        function loadRecentFiles() {
            fetch('/recent_files').then(r => r.json()).then(d => {
                const list = document.getElementById('recentList');
                if(!d.files || d.files.length === 0) {
                    list.innerHTML = '<p style="text-align:center; opacity:0.6; padding:20px;">최근 파일이 없습니다</p>';
                    return;
                }
                list.innerHTML = d.files.map(f => `
                    <div style="display:flex; align-items:center; padding:10px; border-bottom:1px solid var(--border);">
                        <i class="fa-solid ${f.type === 'folder' ? 'fa-folder' : 'fa-file'}" style="margin-right:12px; color:${f.type === 'folder' ? 'var(--folder)' : 'var(--text-secondary)'};"></i>
                        <div style="flex:1; min-width:0;">
                            <a href="/browse/${escapeHtml(f.path)}" style="color:var(--text); text-decoration:none; display:block; overflow:hidden; text-overflow:ellipsis; white-space:nowrap;">${escapeHtml(f.name)}</a>
                            <div style="font-size:0.75rem; opacity:0.6;">${new Date(f.accessed).toLocaleString()}</div>
                        </div>
                    </div>
                `).join('');
            });
        }
        
        // v5.1: 접속자 목록 로드
        function loadActiveSessions() {
            fetch('/active_sessions').then(r => r.json()).then(d => {
                document.getElementById('sessionCount').textContent = `(${d.count}명 접속 중)`;
                const list = document.getElementById('sessionsList');
                if(d.sessions.length === 0) {
                    list.innerHTML = '<p style="text-align:center; opacity:0.6; padding:20px;">현재 접속자가 없습니다</p>';
                    return;
                }
                list.innerHTML = d.sessions.map(s => `
                    <div style="display:flex; align-items:center; padding:10px; border-bottom:1px solid var(--border);">
                        <i class="fa-solid fa-user-circle" style="font-size:1.5rem; margin-right:12px; color:${s.role === 'admin' ? 'var(--primary)' : 'var(--text-secondary)'};"></i>
                        <div style="flex:1;">
                            <div style="font-weight:500;">${escapeHtml(s.ip)}</div>
                            <div style="font-size:0.75rem; opacity:0.7;">${s.role === 'admin' ? '👑 관리자' : '👤 게스트'} · ${s.idle_minutes}분 전 활동</div>
                        </div>
                    </div>
                `).join('');
            });
        }
        
        // v5.1: 디스크 상태 체크
        function checkDiskStatus() {
            fetch('/disk_status').then(r => r.json()).then(d => {
                if(d.warning) {
                    showToast(`⚠️ 디스크 용량 경고! ${d.percent}% 사용 중 (잔여: ${d.free})`, 'error');
                }
            }).catch(() => {});
        }
        
        // v5.1: 드래그앤드롭 파일 이동
        function initFileDragDrop() {
            document.querySelectorAll('.file-item.data-item').forEach(item => {
                // 파일은 드래그 가능
                if(!item.querySelector('.fa-folder')) {
                    item.setAttribute('draggable', 'true');
                    item.addEventListener('dragstart', (e) => {
                        e.dataTransfer.setData('text/plain', item.getAttribute('data-name'));
                        e.dataTransfer.effectAllowed = 'move';
                        item.style.opacity = '0.5';
                    });
                    item.addEventListener('dragend', () => {
                        item.style.opacity = '1';
                    });
                }
                
                // 폴더는 드롭 대상
                if(item.querySelector('.fa-folder')) {
                    item.addEventListener('dragover', (e) => {
                        e.preventDefault();
                        e.dataTransfer.dropEffect = 'move';
                        item.style.background = 'var(--primary)';
                        item.style.opacity = '0.8';
                    });
                    item.addEventListener('dragleave', () => {
                        item.style.background = '';
                        item.style.opacity = '1';
                    });
                    item.addEventListener('drop', (e) => {
                        e.preventDefault();
                        item.style.background = '';
                        item.style.opacity = '1';
                        const sourceFile = e.dataTransfer.getData('text/plain');
                        const destFolder = item.getAttribute('data-name');
                        if(sourceFile && destFolder && canModify) {
                            const srcPath = currentPath ? currentPath + '/' + sourceFile : sourceFile;
                            const dstPath = currentPath ? currentPath + '/' + destFolder : destFolder;
                            fetch('/move', {
                                method: 'POST',
                                headers: {'Content-Type': 'application/json'},
                                body: JSON.stringify({source: srcPath, destination: dstPath})
                            }).then(r => r.json()).then(d => {
                                if(d.success) {
                                    showToast(`${sourceFile} → ${destFolder}로 이동됨`, 'success');
                                    location.reload();
                                } else {
                                    showToast(d.error || '이동 실패', 'error');
                                }
                            });
                        }
                    });
                }
            });
        }
        
        // 페이지 로드 시 초기화
        document.addEventListener('DOMContentLoaded', () => {
            initFileDragDrop();
            checkDiskStatus();
            // v7.0: 메타데이터 로드는 서버 시작 시 처리됨
        });
        
        // ==========================================
        // v7.0: 태그 관리
        // ==========================================
        let currentTagPath = '';
        
        function openTagModal(path, name) {
            currentTagPath = path;
            document.getElementById('tagTargetPath').textContent = '대상: ' + name;
            document.getElementById('tagInput').value = '';
            loadExistingTags(path);
            openModal('tagModal');
        }
        
        function loadExistingTags(path) {
            fetch('/api/tags?path=' + encodeURIComponent(path)).then(r => r.json()).then(d => {
                const container = document.getElementById('existingTags');
                if(!d.tags || d.tags.length === 0) {
                    container.innerHTML = '<span style="opacity:0.5; font-size:0.85rem;">태그 없음</span>';
                    return;
                }
                container.innerHTML = d.tags.map(t => `
                    <span style="background:${t.color}; color:white; padding:4px 10px; border-radius:12px; font-size:0.8rem; display:inline-flex; align-items:center; gap:5px;">
                        ${escapeHtml(t.tag)}
                        <i class="fa-solid fa-xmark" style="cursor:pointer;" onclick="removeTag('${escapeHtml(path)}', '${escapeHtml(t.tag)}')"></i>
                    </span>
                `).join('');
            });
        }
        
        function addTag() {
            const tag = document.getElementById('tagInput').value.trim();
            const color = document.getElementById('tagColor').value;
            if(!tag) { showToast('태그 이름을 입력하세요', 'warning'); return; }
            
            fetch('/api/tags', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({path: currentTagPath, tag, color})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('태그 추가됨', 'success');
                    document.getElementById('tagInput').value = '';
                    loadExistingTags(currentTagPath);
                } else {
                    showToast(d.error || '태그 추가 실패', 'error');
                }
            });
        }
        
        function removeTag(path, tag) {
            fetch('/api/tags', {
                method: 'DELETE',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({path, tag})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('태그 삭제됨', 'success');
                    loadExistingTags(path);
                }
            });
        }
        
        // ==========================================
        // v7.0: 메모 관리
        // ==========================================
        let currentMemoPath = '';
        
        function openMemoModal(path, name) {
            currentMemoPath = path;
            document.getElementById('memoTargetPath').textContent = '대상: ' + name;
            document.getElementById('memoText').value = '';
            document.getElementById('memoUpdated').textContent = '';
            
            fetch('/api/memo/' + encodeURIComponent(path)).then(r => r.json()).then(d => {
                document.getElementById('memoText').value = d.memo || '';
                if(d.updated) {
                    document.getElementById('memoUpdated').textContent = '마지막 수정: ' + new Date(d.updated).toLocaleString();
                }
            });
            openModal('memoModal');
        }
        
        function saveMemo() {
            const memo = document.getElementById('memoText').value;
            fetch('/api/memo/' + encodeURIComponent(currentMemoPath), {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({memo})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('메모 저장됨', 'success');
                    closeModal('memoModal');
                } else {
                    showToast('메모 저장 실패', 'error');
                }
            });
        }
        
        function deleteMemo() {
            if(!confirm('메모를 삭제하시겠습니까?')) return;
            fetch('/api/memo/' + encodeURIComponent(currentMemoPath), {
                method: 'DELETE'
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('메모 삭제됨', 'success');
                    closeModal('memoModal');
                }
            });
        }
        
        // ==========================================
        // v7.0: 암호화/복호화
        // ==========================================
        function encryptFile() {
            const password = document.getElementById('encryptPassword').value || '';
            fetch('/encrypt/' + encodeURIComponent(ctxTarget.path), {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({password})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('파일 암호화됨: ' + d.new_path, 'success');
                    closeModal('encryptModal');
                    setTimeout(() => location.reload(), 500);
                } else {
                    showToast(d.error || '암호화 실패', 'error');
                }
            });
        }
        
        function decryptFile() {
            const password = document.getElementById('decryptPassword').value || '';
            fetch('/decrypt/' + encodeURIComponent(ctxTarget.path), {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({password})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('파일 복호화됨: ' + d.new_path, 'success');
                    closeModal('decryptModal');
                    setTimeout(() => location.reload(), 500);
                } else {
                    showToast(d.error || '복호화 실패', 'error');
                }
            });
        }
        
        // ==========================================
        // v7.0: 즐겨찾기 관리
        // ==========================================
        function addFavorite(path, name) {
            fetch('/api/favorites', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({path, name})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('즐겨찾기 추가됨', 'success');
                } else {
                    showToast(d.error || '즐겨찾기 추가 실패', 'warning');
                }
            });
        }
        
        function loadFavorites() {
            fetch('/api/favorites').then(r => r.json()).then(d => {
                const list = document.getElementById('favoritesList');
                if(!d.favorites || d.favorites.length === 0) {
                    list.innerHTML = '<p style="text-align:center; opacity:0.6;">즐겨찾기가 없습니다.</p>';
                    return;
                }
                list.innerHTML = d.favorites.map(f => `
                    <div style="display:flex; align-items:center; padding:10px; border-bottom:1px solid var(--border);">
                        <i class="fa-solid fa-folder-heart" style="color:var(--danger); margin-right:12px;"></i>
                        <a href="/browse/${escapeHtml(f.path)}" style="flex:1; color:var(--text); text-decoration:none;">${escapeHtml(f.name)}</a>
                        <button class="btn-icon btn-danger" onclick="removeFavorite('${escapeHtml(f.path)}')" style="border:none;background:transparent;">
                            <i class="fa-solid fa-xmark"></i>
                        </button>
                    </div>
                `).join('');
            });
        }
        
        function removeFavorite(path) {
            fetch('/api/favorites', {
                method: 'DELETE',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({path})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('즐겨찾기 삭제됨', 'success');
                    loadFavorites();
                }
            });
        }
        
        // ==========================================
        // v7.0: 접속 대시보드
        // ==========================================
        function loadAccessDashboard() {
            fetch('/api/access_dashboard').then(r => r.json()).then(d => {
                // 활동별 통계
                const actionStats = document.getElementById('actionStats');
                if(d.action_stats && Object.keys(d.action_stats).length > 0) {
                    actionStats.innerHTML = Object.entries(d.action_stats).map(([action, count]) => `
                        <div style="display:flex; justify-content:space-between; padding:4px 0; font-size:0.85rem;">
                            <span>${escapeHtml(action)}</span>
                            <span style="font-weight:bold;">${count}</span>
                        </div>
                    `).join('');
                } else {
                    actionStats.innerHTML = '<span style="opacity:0.5;">데이터 없음</span>';
                }
                
                // 차단 IP
                const blockedList = document.getElementById('blockedIpsList');
                if(d.blocked_ips && d.blocked_ips.length > 0) {
                    blockedList.innerHTML = d.blocked_ips.map(b => `
                        <div style="display:flex; justify-content:space-between; align-items:center; padding:6px 0; font-size:0.85rem; border-bottom:1px solid var(--border);">
                            <span><i class="fa-solid fa-ban" style="color:var(--danger); margin-right:5px;"></i>${escapeHtml(b.ip)}</span>
                            <button class="btn btn-outline" style="font-size:0.7rem; padding:2px 8px;" onclick="unblockIp('${escapeHtml(b.ip)}')">해제</button>
                        </div>
                    `).join('');
                } else {
                    blockedList.innerHTML = '<span style="opacity:0.5; color:var(--success);">차단된 IP 없음</span>';
                }
                
                // 최근 로그
                const recentLogs = document.getElementById('recentLogs');
                if(d.recent_logs && d.recent_logs.length > 0) {
                    recentLogs.innerHTML = d.recent_logs.map(log => `
                        <div style="display:flex; gap:10px; padding:6px 0; font-size:0.8rem; border-bottom:1px solid var(--border);">
                            <span style="width:60px; opacity:0.6;">${new Date(log.time).toLocaleTimeString()}</span>
                            <span style="width:100px;">${escapeHtml(log.ip)}</span>
                            <span style="flex:1;">${escapeHtml(log.action)}</span>
                        </div>
                    `).join('');
                } else {
                    recentLogs.innerHTML = '<span style="opacity:0.5;">로그 없음</span>';
                }
            });
        }
        
        function unblockIp(ip) {
            fetch('/api/unblock/' + encodeURIComponent(ip), {method: 'POST'})
            .then(r => r.json()).then(d => {
                if(d.success) {
                    showToast('IP 차단 해제됨', 'success');
                    loadAccessDashboard();
                } else {
                    showToast(d.error || '차단 해제 실패', 'error');
                }
            });
        }
        
        // ==========================================
        // v7.0: 문서 미리보기
        // ==========================================
        let currentDocPath = '';
        
        function openDocumentPreview(path, filename) {
            currentDocPath = path;
            document.getElementById('docPreviewTitle').innerHTML = '<i class="fa-solid fa-file-alt"></i> ' + escapeHtml(filename);
            document.getElementById('docPreviewContent').innerHTML = '<div style="text-align:center; padding:40px;"><i class="fa-solid fa-spinner fa-spin" style="font-size:2rem;"></i><br>로딩 중...</div>';
            openModal('docPreviewModal');
            
            fetch('/preview/' + encodeURIComponent(path)).then(r => r.json()).then(d => {
                if(d.success) {
                    if(d.type === 'html') {
                        document.getElementById('docPreviewContent').innerHTML = d.content;
                    } else {
                        document.getElementById('docPreviewContent').innerText = d.content;
                    }
                } else {
                    document.getElementById('docPreviewContent').innerHTML = '<div style="color:var(--danger); text-align:center; padding:20px;"><i class="fa-solid fa-exclamation-circle"></i> ' + (d.error || '미리보기 실패') + '</div>';
                }
            }).catch(e => {
                document.getElementById('docPreviewContent').innerHTML = '<div style="color:var(--danger); text-align:center; padding:20px;"><i class="fa-solid fa-exclamation-circle"></i> 네트워크 오류</div>';
            });
        }
        
        function downloadCurrentDoc() {
            if(currentDocPath) {
                location.href = '/download/' + currentDocPath;
            }
        }

        // ==========================================
        // v6.0: 비디오 플레이어
        // ==========================================
        function playVideo(path) {
            const video = document.getElementById('videoPlayer');
            video.src = '/stream/' + path;
            openModal('videoPlayerModal');
            video.play();
        }
        
        // ==========================================
        // v6.0: 오디오 플레이어
        // ==========================================
        let audioTracks = [];
        let currentTrackIndex = 0;
        
        function openAudioPlayer(folderPath) {
            fetch('/playlist/' + folderPath).then(r => r.json()).then(d => {
                if(d.error || d.count === 0) {
                    showToast('오디오 파일이 없습니다', 'warning');
                    return;
                }
                audioTracks = d.tracks;
                currentTrackIndex = 0;
                renderAudioPlaylist();
                playAudioTrack(0);
                openModal('audioPlayerModal');
            });
        }
        
        function renderAudioPlaylist() {
            const list = document.getElementById('audioPlaylist');
            list.innerHTML = audioTracks.map((t, i) => `
                <div class="audio-track" onclick="playAudioTrack(${i})" style="
                    padding:12px 15px; cursor:pointer; display:flex; align-items:center; gap:10px;
                    border-bottom:1px solid var(--border); transition: background 0.2s;
                    ${i === currentTrackIndex ? 'background:var(--hover); color:var(--primary);' : ''}
                " onmouseover="this.style.background='var(--hover)'" onmouseout="this.style.background='${i === currentTrackIndex ? 'var(--hover)' : ''}'">
                    <i class="fa-solid ${i === currentTrackIndex ? 'fa-volume-high' : 'fa-music'}" style="width:20px;"></i>
                    <span style="flex:1; overflow:hidden; text-overflow:ellipsis; white-space:nowrap;">${escapeHtml(t.name)}</span>
                </div>
            `).join('');
        }
        
        function playAudioTrack(index) {
            if(index < 0 || index >= audioTracks.length) return;
            currentTrackIndex = index;
            const track = audioTracks[index];
            const audio = document.getElementById('audioPlayer');
            audio.src = track.stream_url;
            audio.play();
            document.getElementById('audioNowPlaying').textContent = track.name;
            renderAudioPlaylist();
        }
        
        function audioPlayPrev() {
            playAudioTrack((currentTrackIndex - 1 + audioTracks.length) % audioTracks.length);
        }
        
        function audioPlayNext() {
            playAudioTrack((currentTrackIndex + 1) % audioTracks.length);
        }
        
        // 오디오 자동 재생: 다음 트랙
        document.addEventListener('DOMContentLoaded', () => {
            const audio = document.getElementById('audioPlayer');
            if(audio) {
                audio.addEventListener('ended', audioPlayNext);
            }
        });
        
        // ==========================================
        // v6.0: 이미지 갤러리
        // ==========================================
        let galleryImages = [];
        let currentGalleryIndex = 0;
        
        function openGallery(folderPath, startIndex = 0) {
            fetch('/gallery/' + folderPath).then(r => r.json()).then(d => {
                if(d.error || d.count === 0) {
                    showToast('이미지가 없습니다', 'warning');
                    return;
                }
                galleryImages = d.images;
                currentGalleryIndex = startIndex;
                showGalleryImage();
                openModal('galleryModal');
            });
        }
        
        function showGalleryImage() {
            if(galleryImages.length === 0) return;
            const img = galleryImages[currentGalleryIndex];
            document.getElementById('galleryImage').src = img.url;
            document.getElementById('galleryInfo').textContent = `${img.name} (${currentGalleryIndex + 1} / ${galleryImages.length})`;
        }
        
        function galleryPrev() {
            currentGalleryIndex = (currentGalleryIndex - 1 + galleryImages.length) % galleryImages.length;
            showGalleryImage();
        }
        
        function galleryNext() {
            currentGalleryIndex = (currentGalleryIndex + 1) % galleryImages.length;
            showGalleryImage();
        }
        
        // 갤러리 키보드 탐색
        document.addEventListener('keydown', (e) => {
            if(document.getElementById('galleryModal').style.display === 'flex') {
                if(e.key === 'ArrowLeft') galleryPrev();
                if(e.key === 'ArrowRight') galleryNext();
                if(e.key === 'Escape') closeModal('galleryModal');
            }
        });
        
        // ==========================================
        // v6.0: 사용자 관리
        // ==========================================
        function showAddUserForm() {
            document.getElementById('userFormArea').style.display = 'block';
        }
        
        function hideAddUserForm() {
            document.getElementById('userFormArea').style.display = 'none';
            document.getElementById('newUsername').value = '';
            document.getElementById('newPassword').value = '';
        }
        
        function loadUsers() {
            fetch('/api/users').then(r => r.json()).then(d => {
                const list = document.getElementById('userList');
                const users = Object.entries(d.users || {});
                if(users.length === 0) {
                    list.innerHTML = '<p style="text-align:center; opacity:0.6; padding:20px;">등록된 사용자가 없습니다</p>';
                    return;
                }
                list.innerHTML = users.map(([username, info]) => `
                    <div style="display:flex; align-items:center; padding:12px; border-bottom:1px solid var(--border);">
                        <i class="fa-solid ${info.role === 'admin' ? 'fa-user-shield' : 'fa-user'}" style="font-size:1.3rem; margin-right:12px; color:${info.role === 'admin' ? 'var(--primary)' : 'var(--text-secondary)'};"></i>
                        <div style="flex:1;">
                            <div style="font-weight:500;">${escapeHtml(username)}</div>
                            <div style="font-size:0.75rem; opacity:0.7;">
                                ${info.role === 'admin' ? '👑 관리자' : '👤 사용자'} · 
                                ${info.quota_mb > 0 ? info.usage_mb + ' / ' + info.quota_mb + ' MB' : '무제한'}
                            </div>
                        </div>
                        <button class="btn btn-outline btn-icon" onclick="deleteUser('${username}')" title="삭제">
                            <i class="fa-solid fa-trash"></i>
                        </button>
                    </div>
                `).join('');
            });
        }
        
        function createUser() {
            const username = document.getElementById('newUsername').value.trim();
            const password = document.getElementById('newPassword').value;
            const role = document.getElementById('newRole').value;
            const quota = parseInt(document.getElementById('newQuota').value) || 1024;
            
            if(!username || !password) {
                showToast('사용자명과 비밀번호를 입력하세요', 'warning');
                return;
            }
            
            fetch('/api/users', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({username, password, role, quota_mb: quota})
            }).then(r => r.json()).then(d => {
                if(d.success) {
                    showToast(`사용자 '${username}' 생성됨`, 'success');
                    hideAddUserForm();
                    loadUsers();
                } else {
                    showToast(d.error || '생성 실패', 'error');
                }
            });
        }
        
        function deleteUser(username) {
            if(!confirm(`'${username}' 사용자를 삭제하시겠습니까?`)) return;
            fetch('/api/users/' + username, {method: 'DELETE'})
                .then(r => r.json()).then(d => {
                    if(d.success) {
                        showToast(`'${username}' 삭제됨`, 'success');
                        loadUsers();
                    } else {
                        showToast(d.error || '삭제 실패', 'error');
                    }
                });
        }
        
        function openUserManagement() {
            openModal('userManageModal');
            loadUsers();
        }
        
        // ==========================================
        // v6.0: 청크 업로드
        // ==========================================
        async function uploadLargeFile(file, targetPath) {
            const CHUNK_SIZE = 5 * 1024 * 1024; // 5MB
            const totalChunks = Math.ceil(file.size / CHUNK_SIZE);
            
            // 세션 시작
            const initRes = await fetch('/upload/chunk/init', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({
                    filename: file.name,
                    total_size: file.size,
                    path: targetPath
                })
            }).then(r => r.json());
            
            if(!initRes.success) {
                showToast('업로드 시작 실패: ' + initRes.error, 'error');
                return false;
            }
            
            const sessionId = initRes.session_id;
            
            // 청크 업로드
            for(let i = 0; i < totalChunks; i++) {
                const start = i * CHUNK_SIZE;
                const end = Math.min(start + CHUNK_SIZE, file.size);
                const chunk = file.slice(start, end);
                
                const formData = new FormData();
                formData.append('index', i);
                formData.append('chunk', chunk);
                
                const chunkRes = await fetch('/upload/chunk/' + sessionId, {
                    method: 'POST',
                    body: formData
                }).then(r => r.json());
                
                if(!chunkRes.success) {
                    showToast('청크 업로드 실패', 'error');
                    await fetch('/upload/chunk/' + sessionId + '/cancel', {method: 'POST'});
                    return false;
                }
                
                // 진행률 업데이트
                const progress = Math.round(((i + 1) / totalChunks) * 100);
                document.getElementById('uploadProgressText').textContent = `${file.name}: ${progress}%`;
                document.getElementById('uploadProgressBar').style.width = progress + '%';
            }
            
            // 완료
            const completeRes = await fetch('/upload/chunk/' + sessionId + '/complete', {
                method: 'POST'
            }).then(r => r.json());
            
            if(completeRes.success) {
                showToast(`${file.name} 업로드 완료`, 'success');
                return true;
            } else {
                showToast('파일 병합 실패: ' + completeRes.error, 'error');
                return false;
            }
        }
    </script>
</body>
</html>
"""

# ==========================================
# 4. Flask 웹 서버 로직
# ==========================================
app = Flask(__name__)
app.secret_key = os.urandom(24)
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024 * 1024  # 10GB 제한
app.jinja_env.globals['csrf_token'] = generate_csrf_token  # 템플릿에서 csrf_token() 사용 가능

clipboard_store = ""
login_block = {} 

def get_real_ip():
    if request.headers.get('X-Forwarded-For'):
        return request.headers.get('X-Forwarded-For').split(',')[0]
    return request.remote_addr

@app.before_request
def before_request():
    g.start = time.time()
    with _stats_lock:
        STATS['requests'] += 1
        STATS['active_connections'] += 1
    
    # v5.1: IP 화이트리스트 체크 (로그인 페이지 제외)
    client_ip = get_real_ip()
    if request.endpoint and request.endpoint not in ['index', 'static']:
        if not check_ip_whitelist(client_ip):
            logger.add(f"차단된 IP: {client_ip}", "WARN")
            return jsonify({'error': get_text('ip_blocked')}), 403

    # v7.1: CSRF 토큰 검증
    if request.method == "POST":
        # 로그인 등 일부 예외 처리 필요 시 추가
        token = session.pop('_csrf_token', None)
        if not token or token != request.form.get('csrf_token'):
            # AJAX 요청 헤더 확인
            if not token or token != request.headers.get('X-CSRF-Token'):
                # [주의] 세션 초기화나 재발급 로직이 필요할 수 있음
                # 여기서는 검증 실패 시 403
                logger.add(f"CSRF 검증 실패: {client_ip}", "WARN")
                return jsonify({'error': 'CSRF token missing or incorrect'}), 403
    
    # 세션 타임아웃 검사
    if session.get('logged_in'):
        last_active = session.get('last_active')
        if last_active:
            timeout = conf.get('session_timeout') or SESSION_TIMEOUT_MINUTES
            if datetime.now().timestamp() - last_active > timeout * 60:
                session.clear()
                logger.add(f"세션 만료: {client_ip}")
        session['last_active'] = datetime.now().timestamp()
        
        # v5.1: 활성 세션 추적
        sid = session.get('_id', id(session))
        ACTIVE_SESSIONS[sid] = {
            'ip': client_ip,
            'role': session.get('role', 'guest'),
            'login_time': session.get('login_time', datetime.now()),
            'last_active': datetime.now()
        }

@app.after_request
def after_request(response):
    """응답 후 처리 (스레드 안전)"""
    with _stats_lock:
        if response.content_length:
            STATS['bytes_sent'] += response.content_length
        STATS['active_connections'] = max(0, STATS['active_connections'] - 1)
    return response

def login_required(role_req='guest'):
    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            if not session.get('logged_in'):
                return jsonify({'error': '로그인이 필요합니다.'}), 401
            if role_req == 'admin' and session.get('role') != 'admin':
                return jsonify({'error': '관리자 권한이 필요합니다.'}), 403
            return f(*args, **kwargs)
        return decorated_function
    return decorator

def check_security():
    """v7.0: IP 차단 상태 확인 (새로운 LOGIN_ATTEMPTS 시스템 사용)"""
    ip = get_real_ip()
    is_blocked, remaining = check_ip_blocked(ip)
    return not is_blocked

@app.route('/', defaults={'path': ''}, methods=['GET', 'POST'])
@app.route('/browse/<path:path>', methods=['GET', 'POST'])
def index(path):
    if not check_security():
        return render_template_string(HTML_TEMPLATE, logged_in=False, error="보안 차단됨: 잠시 후 다시 시도하세요.")

    if request.method == 'POST':
        pw = request.form.get('password')
        ip = get_real_ip()
        
        # v7.0: IP 차단 상태 먼저 확인
        is_blocked, remaining = check_ip_blocked(ip)
        if is_blocked:
            return render_template_string(HTML_TEMPLATE, logged_in=False, 
                error=f"IP 차단됨: {remaining}분 후 다시 시도하세요.")
        
        # v4: verify_password 사용 (해시 + 평문 호환)
        if verify_password(conf.get('admin_pw'), pw):
            session['logged_in'] = True
            session['role'] = 'admin'
            session['last_active'] = datetime.now().timestamp()
            record_login_attempt(ip, True)  # v7.0: 성공 기록
            logger.add(f"관리자 로그인: {ip}")
            log_access(ip, 'login', 'admin')
            return redirect(url_for('index', path=path))
        elif verify_password(conf.get('guest_pw'), pw):
            session['logged_in'] = True
            session['role'] = 'guest'
            session['last_active'] = datetime.now().timestamp()
            record_login_attempt(ip, True)  # v7.0: 성공 기록
            logger.add(f"게스트 로그인: {ip}")
            log_access(ip, 'login', 'guest')
            return redirect(url_for('index', path=path))
        else:
            record_login_attempt(ip, False)  # v7.0: 실패 기록
            attempts = LOGIN_ATTEMPTS.get(ip, {}).get('attempts', 0)
            remaining_attempts = MAX_LOGIN_ATTEMPTS - attempts
            log_access(ip, 'login_failed', f'remaining: {remaining_attempts}')
            
            if remaining_attempts <= 0:
                return render_template_string(HTML_TEMPLATE, logged_in=False, 
                    error=f"로그인 시도 초과. {LOGIN_BLOCK_MINUTES}분간 차단됩니다.")
            return render_template_string(HTML_TEMPLATE, logged_in=False, 
                error=f"비밀번호가 올바르지 않습니다. (남은 시도: {remaining_attempts}회)")

    if not session.get('logged_in'):
        return render_template_string(HTML_TEMPLATE, logged_in=False)

    base_dir = conf.get('folder')
    abs_path = os.path.join(base_dir, path)
    
    try:
        if not os.path.abspath(abs_path).startswith(os.path.abspath(base_dir)):
            return abort(403)
    except Exception: return abort(403)

    if not os.path.exists(abs_path): return abort(404)

    items = []
    try:
        with os.scandir(abs_path) as entries:
            for entry in entries:
                f_type = 'file'
                if entry.is_dir(): f_type = 'folder'
                else:
                    if entry.name.lower().endswith(tuple(TEXT_EXTENSIONS)):
                        f_type = 'text'
                    else:
                        mt, _ = mimetypes.guess_type(entry.name)
                        if mt:
                            if mt.startswith('image'): f_type = 'image'
                            elif mt.startswith('video'): f_type = 'video'
                            elif mt.startswith('audio'): f_type = 'audio'
                            elif mt in ['application/zip', 'application/x-rar-compressed']: f_type = 'archive'

                stat = entry.stat()
                size_str = "-"
                raw_size = 0
                if not entry.is_dir():
                    raw_size = stat.st_size
                    if raw_size < 1024: size_str = f"{raw_size} B"
                    elif raw_size < 1024*1024: size_str = f"{raw_size/1024:.1f} KB"
                    else: size_str = f"{raw_size/(1024*1024):.1f} MB"

                items.append({
                    'name': entry.name,
                    'is_dir': entry.is_dir(),
                    'type': f_type,
                    'size': size_str,
                    'raw_size': raw_size,
                    'raw_mtime': stat.st_mtime,
                    'ext': os.path.splitext(entry.name)[1],
                    'mod_time': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M'),
                    'rel_path': os.path.relpath(entry.path, base_dir).replace('\\', '/')
                })
    except Exception as e:
        logger.add(f"탐색 오류: {e}", "ERROR")
    
    # 디버그: 찾은 파일 수 로깅
    logger.add(f"폴더 탐색: {path or '/'} ({len(items)}개 항목)")

    items.sort(key=lambda x: (not x['is_dir'], x['name'].lower()))
    can_modify = (session.get('role') == 'admin') or (conf.get('allow_guest_upload'))
    return render_template_string(HTML_TEMPLATE, logged_in=True, role=session.get('role'), 
                                  items=items, current_path=path, can_modify=can_modify)

@app.route('/metrics')
@login_required()
def metrics():
    uptime = datetime.now() - SERVER_START_TIME
    uptime_str = str(uptime).split('.')[0]
    
    # 스레드 안전하게 STATS 읽기
    with _stats_lock:
        stats_copy = STATS.copy()

    return jsonify({
        'uptime': uptime_str,
        'requests': stats_copy['requests'],
        'sent': fmt_bytes(stats_copy['bytes_sent']),
        'recv': fmt_bytes(stats_copy['bytes_received']),
        'active': stats_copy['active_connections']
    })

@app.route('/upload/<path:path>', methods=['POST'])
def upload_file(path):
    # 권한 체크 로직 통합
    if not (session.get('role')=='admin' or conf.get('allow_guest_upload')):
        return jsonify({'error':'권한 없음'}), 403
    
    target_dir = os.path.join(conf.get('folder'), path)
    files = request.files.getlist('file')
    paths = request.form.getlist('paths') 
    
    count = 0
    total_size = 0
    for i, file in enumerate(files):
        if file.filename:
            file.seek(0, os.SEEK_END)
            total_size += file.tell()
            file.seek(0)
            
            # [수정됨] safe_filename 사용으로 한글 지원
            safe_name = safe_filename(file.filename)

            if paths and len(paths) > i and '/' in paths[i]:
                rel_path = paths[i]
                if '..' in rel_path: continue
                # 경로 부분의 파일명도 안전하게 처리해야 함
                parts = rel_path.split('/')
                safe_parts = [safe_filename(p) for p in parts]
                save_path = os.path.join(target_dir, *safe_parts)
                os.makedirs(os.path.dirname(save_path), exist_ok=True)
                file.save(save_path)
            else:
                file.save(os.path.join(target_dir, safe_name))
            count += 1
    
    with _stats_lock:
        STATS['bytes_received'] += total_size
    logger.add(f"업로드: {count}개 항목 -> /{path}")
    return jsonify({'success': True})

@app.route('/batch_download/<path:path>', methods=['POST'])
def batch_download(path):
    if not session.get('logged_in'): return abort(401)
    base_dir = conf.get('folder')
    
    # 현재 디렉토리 경로 검증
    is_valid, current_dir, error = validate_path(base_dir, path)
    if not is_valid:
        return abort(403)
    
    try:
        data = json.loads(request.form.get('files'))
        mem_zip = io.BytesIO()
        with zipfile.ZipFile(mem_zip, 'w', zipfile.ZIP_DEFLATED) as zf:
            for item_name in data:
                # 각 항목 경로 검증
                item_rel = os.path.join(path, safe_filename(item_name)).replace('\\', '/')
                is_valid_item, item_path, _ = validate_path(base_dir, item_rel)
                if not is_valid_item:
                    continue  # 유효하지 않은 경로 건너뜀
                
                if os.path.isfile(item_path):
                    zf.write(item_path, item_name)
                elif os.path.isdir(item_path):
                    for root, dirs, files in os.walk(item_path):
                        for file in files:
                            abs_file = os.path.join(root, file)
                            rel_file = os.path.relpath(abs_file, current_dir)
                            zf.write(abs_file, rel_file)
        
        mem_zip.seek(0)
        return send_file(mem_zip, download_name=f"batch_download.zip", as_attachment=True)
    except Exception as e:
        logger.add(f"배치 다운로드 오류: {e}", "ERROR")
        return abort(500)

@app.route('/batch_delete/<path:path>', methods=['POST'])
@login_required('admin')
def batch_delete(path):
    base_dir = conf.get('folder')
    current_dir = os.path.join(base_dir, path)
    data = request.get_json()
    files = data.get('files', [])
    
    count = 0
    try:
        for item_name in files:
            # [수정됨] safe_filename 적용
            item_path = os.path.join(current_dir, safe_filename(item_name))
            if os.path.exists(item_path):
                if os.path.isfile(item_path): os.remove(item_path)
                else: shutil.rmtree(item_path)
                count += 1
        logger.add(f"일괄 삭제: {count}개 항목")
        return jsonify({'success': True})
    except Exception as e: return jsonify({'success': False, 'error': str(e)})

@app.route('/download/<path:filename>')
def download_file(filename):
    if not session.get('logged_in'): return abort(401)
    
    # 경로 검증
    is_valid, full_path, error = validate_path(conf.get('folder'), filename)
    if not is_valid:
        logger.add(f"다운로드 경로 검증 실패: {filename}", "WARN")
        return abort(403)
    
    if not os.path.exists(full_path):
        return abort(404)
    
    return send_from_directory(conf.get('folder'), filename)

@app.route('/mkdir/<path:path>', methods=['POST'])
def mkdir(path):
    if not (session.get('role')=='admin' or conf.get('allow_guest_upload')): return jsonify({'error':'권한 없음'}), 403
    try:
        data = request.get_json()
        # [수정됨] safe_filename 적용
        new_dir = os.path.join(conf.get('folder'), path, safe_filename(data['name']))
        os.makedirs(new_dir, exist_ok=True)
        logger.add(f"폴더 생성: {data['name']}")
        return jsonify({'success': True})
    except Exception as e: return jsonify({'success': False, 'error': str(e)})

@app.route('/delete/<path:path>', methods=['POST'])
@login_required('admin')
def delete_item(path):
    # 경로 검증
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid:
        return jsonify({'success': False, 'error': error}), 403
    
    try:
        if os.path.isfile(full_path): 
            os.remove(full_path)
        else: 
            shutil.rmtree(full_path)
        logger.add(f"삭제: {path}")
        return jsonify({'success': True})
    except Exception as e: 
        logger.add(f"삭제 오류: {e}", "ERROR")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/rename/<path:path>', methods=['POST'])
@login_required('admin')
def rename_item(path):
    data = request.get_json()
    base = os.path.join(conf.get('folder'), path)
    # [수정됨] safe_filename 적용
    old = os.path.join(base, safe_filename(data['old_name']))
    new = os.path.join(base, safe_filename(data['new_name']))
    try:
        os.rename(old, new)
        logger.add(f"이름변경: {data['old_name']} -> {data['new_name']}")
        return jsonify({'success': True})
    except Exception as e: return jsonify({'success': False, 'error': str(e)})

@app.route('/zip/<path:path>')
@login_required()
def download_zip(path):
    base_dir = conf.get('folder')
    target_dir = os.path.join(base_dir, path)
    if not os.path.exists(target_dir): return abort(404)
    
    mem_zip = io.BytesIO()
    with zipfile.ZipFile(mem_zip, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(target_dir):
            for file in files:
                file_path = os.path.join(root, file)
                zf.write(file_path, os.path.relpath(file_path, target_dir))
    mem_zip.seek(0)
    return send_file(mem_zip, download_name=f"{os.path.basename(target_dir)}.zip", as_attachment=True)

@app.route('/unzip/<path:path>', methods=['POST'])
@login_required('admin')
def unzip_file(path):
    """ZIP 파일 압축 해제 (Zip Slip 공격 방지 포함)"""
    zip_path = os.path.join(conf.get('folder'), path)
    extract_to = os.path.splitext(zip_path)[0]
    
    try:
        with zipfile.ZipFile(zip_path, 'r') as zf:
            # Zip Slip 공격 방지: 각 파일의 대상 경로 검증
            extract_to_abs = os.path.abspath(extract_to)
            for member in zf.namelist():
                # 경로 정규화
                member_path = os.path.normpath(os.path.join(extract_to, member))
                # 대상 디렉토리 외부로 탈출하는지 확인
                if not os.path.abspath(member_path).startswith(extract_to_abs + os.sep) and \
                   os.path.abspath(member_path) != extract_to_abs:
                    logger.add(f"Zip Slip 공격 감지: {member}", "WARN")
                    return jsonify({'success': False, 'error': f'보안 위협 감지: 잘못된 경로 "{member}"'}), 400
            
            # 안전하게 압축 해제
            zf.extractall(extract_to)
        logger.add(f"압축해제: {path}")
        return jsonify({'success': True})
    except zipfile.BadZipFile:
        return jsonify({'success': False, 'error': '잘못된 ZIP 파일입니다.'})
    except Exception as e: 
        logger.add(f"압축해제 오류: {e}", "ERROR")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/get_content/<path:path>')
@login_required()
def get_content(path):
    # 경로 검증
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid:
        return jsonify({'error': error}), 403
    
    try:
        with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
            return jsonify({'content': f.read()})
    except Exception as e: 
        logger.add(f"파일 읽기 오류: {e}", "ERROR")
        return jsonify({'error': str(e)})

@app.route('/save_content/<path:path>', methods=['POST'])
@login_required('admin')
def save_content(path):
    # 경로 검증
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid:
        return jsonify({'success': False, 'error': error}), 403
    
    try:
        content = request.get_json().get('content', '')
        with open(full_path, 'w', encoding='utf-8') as f:
            f.write(content)
        logger.add(f"파일수정: {path}")
        return jsonify({'success': True})
    except Exception as e: 
        logger.add(f"파일 저장 오류: {e}", "ERROR")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/disk_info')
@login_required()
def disk_info():
    try:
        t, u, f = shutil.disk_usage(conf.get('folder'))
        return jsonify({
            'total': f"{t/1024**3:.1f}GB", 
            'used': f"{u/1024**3:.1f}GB", 
            'percent': round((u/t)*100, 1)
        })
    except Exception as e:
        logger.add(f"디스크 정보 조회 오류: {e}", "ERROR")
        return jsonify({'error': '디스크 정보를 가져올 수 없습니다.'})

@app.route('/clipboard', methods=['GET', 'POST'])
def clipboard_handler():
    global clipboard_store
    if not session.get('logged_in'): return jsonify({'error':'Auth required'}), 401
    if request.method == 'POST':
        clipboard_store = request.get_json().get('content', '')
        return jsonify({'success': True})
    return jsonify({'content': clipboard_store})

# ==========================================
# v4 신규 API: 파일 관리 기능 확장
# ==========================================

@app.route('/copy', methods=['POST'])
@login_required('admin')
def copy_item():
    """파일/폴더 복사"""
    data = request.get_json()
    src_path = data.get('source', '')
    dst_path = data.get('destination', '')
    
    base_dir = conf.get('folder')
    is_valid_src, full_src, _ = validate_path(base_dir, src_path)
    is_valid_dst, full_dst, _ = validate_path(base_dir, dst_path)
    
    if not is_valid_src or not is_valid_dst:
        return jsonify({'success': False, 'error': '잘못된 경로입니다.'})
    
    if not os.path.exists(full_src):
        return jsonify({'success': False, 'error': '원본을 찾을 수 없습니다.'})
    
    try:
        if os.path.isdir(full_src):
            shutil.copytree(full_src, full_dst)
        else:
            os.makedirs(os.path.dirname(full_dst), exist_ok=True)
            shutil.copy2(full_src, full_dst)
        logger.add(f"복사: {src_path} -> {dst_path}")
        log_access(get_real_ip(), 'copy', f"{src_path} -> {dst_path}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/move', methods=['POST'])
@login_required('admin')
def move_item():
    """파일/폴더 이동"""
    data = request.get_json()
    src_path = data.get('source', '')
    dst_path = data.get('destination', '')
    
    base_dir = conf.get('folder')
    is_valid_src, full_src, _ = validate_path(base_dir, src_path)
    is_valid_dst, full_dst, _ = validate_path(base_dir, dst_path)
    
    if not is_valid_src or not is_valid_dst:
        return jsonify({'success': False, 'error': '잘못된 경로입니다.'})
    
    if not os.path.exists(full_src):
        return jsonify({'success': False, 'error': '원본을 찾을 수 없습니다.'})
    
    try:
        os.makedirs(os.path.dirname(full_dst), exist_ok=True)
        shutil.move(full_src, full_dst)
        logger.add(f"이동: {src_path} -> {dst_path}")
        log_access(get_real_ip(), 'move', f"{src_path} -> {dst_path}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/search')
@login_required()
def search_files():
    """서버 전체 파일 검색"""
    query = request.args.get('q', '').lower().strip()
    if not query or len(query) < 2:
        return jsonify({'results': [], 'error': '검색어는 2자 이상이어야 합니다.'})
    
    base_dir = conf.get('folder')
    results = []
    max_results = 100
    
    try:
        for root, dirs, files in os.walk(base_dir):
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            
            for name in files + dirs:
                if query in name.lower():
                    rel_path = os.path.relpath(os.path.join(root, name), base_dir).replace('\\', '/')
                    results.append({'name': name, 'path': rel_path, 'is_dir': name in dirs})
                    if len(results) >= max_results:
                        break
            if len(results) >= max_results:
                break
    except Exception as e:
        logger.add(f"검색 오류: {e}", "ERROR")
    
    return jsonify({'results': results, 'count': len(results)})

@app.route('/thumbnail/<path:filepath>')
@login_required()
def get_thumbnail(filepath):
    """이미지 썸네일 생성 (스레드 안전, LRU 캐시)"""
    is_valid, full_path, _ = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.exists(full_path):
        return abort(404)
    
    cache_key = f"{filepath}_{os.path.getmtime(full_path)}"
    
    # 캐시 확인 (스레드 안전)
    with _cache_lock:
        if cache_key in THUMBNAIL_CACHE:
            return send_file(io.BytesIO(THUMBNAIL_CACHE[cache_key]), mimetype='image/jpeg')
    
    try:
        img = Image.open(full_path)
        img.thumbnail((150, 150), Image.Resampling.LANCZOS)
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        buffer = io.BytesIO()
        img.save(buffer, format='JPEG', quality=70)
        buffer.seek(0)
        
        # 캐시 저장 (스레드 안전, LRU 방식 개선)
        with _cache_lock:
            if len(THUMBNAIL_CACHE) >= MAX_THUMBNAIL_CACHE:
                # OrderedDict: popitem(last=False)는 가장 먼저 들어온(오래된) 항목 제거
                THUMBNAIL_CACHE.popitem(last=False)
            THUMBNAIL_CACHE[cache_key] = buffer.getvalue()
            # 최신 항목을 끝으로 이동 (갱신)
            if cache_key in THUMBNAIL_CACHE:
                THUMBNAIL_CACHE.move_to_end(cache_key)
        
        buffer.seek(0)
        return send_file(buffer, mimetype='image/jpeg')
    except Exception as e:
        logger.add(f"썸네일 생성 실패: {e}", "ERROR")
        return abort(500)

@app.route('/versions/<path:filepath>')
@login_required('admin')
def list_versions(filepath):
    """파일 버전 목록 조회"""
    filename = os.path.basename(filepath)
    version_dir = os.path.join(conf.get('folder'), VERSION_FOLDER_NAME)
    if not os.path.exists(version_dir):
        return jsonify({'versions': []})
    versions = []
    for f in os.listdir(version_dir):
        if f.endswith(f'_{filename}'):
            full_path = os.path.join(version_dir, f)
            versions.append({'name': f, 'timestamp': f.rsplit('_', 1)[0], 'size': os.path.getsize(full_path)})
    versions.sort(key=lambda x: x['timestamp'], reverse=True)
    return jsonify({'versions': versions})

@app.route('/versions/restore', methods=['POST'])
@login_required('admin')
def restore_version():
    """파일 버전 복원"""
    data = request.get_json()
    version_name = data.get('version', '')
    target_path = data.get('target', '')
    version_dir = os.path.join(conf.get('folder'), VERSION_FOLDER_NAME)
    version_path = os.path.join(version_dir, safe_filename(version_name))
    is_valid, full_target, _ = validate_path(conf.get('folder'), target_path)
    if not os.path.exists(version_path) or not is_valid:
        return jsonify({'success': False, 'error': '파일을 찾을 수 없습니다.'})
    try:
        shutil.copy2(version_path, full_target)
        logger.add(f"버전 복원: {version_name} -> {target_path}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/access_log')
@login_required('admin')
def get_access_log():
    """접속 기록 조회"""
    return jsonify({'logs': ACCESS_LOG})

# ==========================================
# v7.0 신규 API 엔드포인트
# ==========================================

@app.route('/api/blocked_ips')
@login_required('admin')
def api_blocked_ips():
    """v7.0: 차단된 IP 목록 조회"""
    return jsonify({'blocked': get_blocked_ips()})

@app.route('/api/unblock/<ip>', methods=['POST'])
@login_required('admin')
def api_unblock_ip(ip):
    """v7.0: IP 차단 해제"""
    if unblock_ip(ip):
        return jsonify({'success': True})
    return jsonify({'success': False, 'error': '차단된 IP가 아닙니다.'})

@app.route('/encrypt/<path:filepath>', methods=['POST'])
@login_required('admin')
def api_encrypt_file(filepath):
    """v7.0: 파일 암호화"""
    data = request.get_json()
    password = data.get('password', conf.get('admin_pw'))  # 기본값: 관리자 비밀번호
    
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return jsonify({'success': False, 'error': '파일을 찾을 수 없습니다.'}), 404
    
    if full_path.endswith('.enc'):
        return jsonify({'success': False, 'error': '이미 암호화된 파일입니다.'})
    
    success, result = encrypt_file_aes(full_path, password)
    if success:
        return jsonify({'success': True, 'new_path': os.path.basename(result)})
    return jsonify({'success': False, 'error': result})

@app.route('/decrypt/<path:filepath>', methods=['POST'])
@login_required('admin')
def api_decrypt_file(filepath):
    """v7.0: 파일 복호화"""
    data = request.get_json()
    password = data.get('password', conf.get('admin_pw'))
    
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return jsonify({'success': False, 'error': '파일을 찾을 수 없습니다.'}), 404
    
    success, result = decrypt_file_aes(full_path, password)
    if success:
        return jsonify({'success': True, 'new_path': os.path.basename(result)})
    return jsonify({'success': False, 'error': result})

@app.route('/api/tags', methods=['GET', 'POST', 'DELETE'])
@login_required()
def api_file_tags():
    """v7.0: 파일 태그 관리"""
    global FILE_TAGS
    
    if request.method == 'GET':
        path = request.args.get('path', '')
        if path:
            return jsonify({'tags': FILE_TAGS.get(path, [])})
        return jsonify({'all_tags': FILE_TAGS})
    
    data = request.get_json()
    path = data.get('path', '')
    
    if request.method == 'POST':
        tag = data.get('tag', '')
        color = data.get('color', '#6366f1')  # 기본 색상: 보라색
        
        if not path or not tag:
            return jsonify({'success': False, 'error': '경로와 태그가 필요합니다.'})
        
        if path not in FILE_TAGS:
            FILE_TAGS[path] = []
        
        # 중복 태그 확인
        if any(t['tag'] == tag for t in FILE_TAGS[path]):
            return jsonify({'success': False, 'error': '이미 존재하는 태그입니다.'})
        
        FILE_TAGS[path].append({'tag': tag, 'color': color})
        save_metadata()
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        tag = data.get('tag', '')
        if path in FILE_TAGS:
            FILE_TAGS[path] = [t for t in FILE_TAGS[path] if t['tag'] != tag]
            if not FILE_TAGS[path]:
                del FILE_TAGS[path]
            save_metadata()
        return jsonify({'success': True})

@app.route('/api/favorites', methods=['GET', 'POST', 'DELETE'])
@login_required()
def api_favorites():
    """v7.0: 즐겨찾기 폴더 관리"""
    global FAVORITE_FOLDERS
    
    if request.method == 'GET':
        return jsonify({'favorites': FAVORITE_FOLDERS})
    
    data = request.get_json()
    path = data.get('path', '')
    name = data.get('name', os.path.basename(path) if path else '')
    
    if request.method == 'POST':
        if not path:
            return jsonify({'success': False, 'error': '경로가 필요합니다.'})
        
        # 중복 확인
        if any(f['path'] == path for f in FAVORITE_FOLDERS):
            return jsonify({'success': False, 'error': '이미 즐겨찾기에 추가되어 있습니다.'})
        
        FAVORITE_FOLDERS.append({
            'path': path, 
            'name': name, 
            'added': datetime.now().isoformat()
        })
        save_metadata()
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        FAVORITE_FOLDERS = [f for f in FAVORITE_FOLDERS if f['path'] != path]
        save_metadata()
        return jsonify({'success': True})

@app.route('/api/memo/<path:filepath>', methods=['GET', 'POST', 'DELETE'])
@login_required()
def api_file_memo(filepath):
    """v7.0: 파일 메모 관리"""
    global FILE_MEMOS
    
    if request.method == 'GET':
        memo = FILE_MEMOS.get(filepath, {})
        return jsonify({'memo': memo.get('memo', ''), 'updated': memo.get('updated', '')})
    
    if request.method == 'POST':
        data = request.get_json()
        memo_text = data.get('memo', '')
        
        FILE_MEMOS[filepath] = {
            'memo': memo_text,
            'updated': datetime.now().isoformat()
        }
        save_metadata()
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        if filepath in FILE_MEMOS:
            del FILE_MEMOS[filepath]
            save_metadata()
        return jsonify({'success': True})

@app.route('/video_thumbnail/<path:filepath>')
@login_required()
def api_video_thumbnail(filepath):
    """v7.0: 동영상 썸네일 반환"""
    is_valid, full_path, _ = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return abort(404)
    
    thumb_path = generate_video_thumbnail(full_path)
    if thumb_path and os.path.exists(thumb_path):
        return send_file(thumb_path, mimetype='image/jpeg')
    
    # 썸네일 생성 실패 시 기본 아이콘 반환 (없으면 404)
    return abort(404)

@app.route('/api/access_dashboard')
@login_required('admin')
def api_access_dashboard():
    """v7.0: 접속 대시보드 데이터"""
    # 시간별 접속 통계 계산
    hourly_stats = {}
    action_stats = {}
    ip_stats = {}
    
    for log in ACCESS_LOG:
        try:
            log_time = datetime.fromisoformat(log['time'])
            hour = log_time.strftime('%H:00')
            hourly_stats[hour] = hourly_stats.get(hour, 0) + 1
            
            action = log.get('action', 'unknown')
            action_stats[action] = action_stats.get(action, 0) + 1
            
            ip = log.get('ip', 'unknown')
            ip_stats[ip] = ip_stats.get(ip, 0) + 1
        except:
            continue
    
    # 최근 10개 접속 기록
    recent_logs = ACCESS_LOG[:10]
    
    # 현재 차단 IP
    blocked = get_blocked_ips()
    
    return jsonify({
        'hourly_stats': hourly_stats,
        'action_stats': action_stats,
        'ip_stats': ip_stats,
        'recent_logs': recent_logs,
        'blocked_ips': blocked,
        'total_logs': len(ACCESS_LOG)
    })

@app.route('/api/trash_settings', methods=['GET', 'POST'])
@login_required('admin')
def api_trash_settings():
    """v7.0: 휴지통 설정"""
    if request.method == 'GET':
        return jsonify({
            'auto_delete_days': conf.get('trash_auto_delete_days') or TRASH_AUTO_DELETE_DAYS
        })
    
    data = request.get_json()
    days = data.get('days', TRASH_AUTO_DELETE_DAYS)
    conf.set('trash_auto_delete_days', int(days))
    conf.save()
    return jsonify({'success': True})

@app.route('/api/cleanup_trash', methods=['POST'])
@login_required('admin')
def api_cleanup_trash():
    """v7.0: 휴지통 자동 정리 수동 실행"""
    deleted_count = auto_cleanup_trash()
    return jsonify({'success': True, 'deleted': deleted_count})

@app.route('/preview/<path:filepath>')
@login_required()
def api_document_preview(filepath):
    """v7.0: 문서 미리보기 (Word, Excel, PowerPoint)"""
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.isfile(full_path):
        return jsonify({'error': '파일을 찾을 수 없습니다.'}), 404
    
    ext = os.path.splitext(full_path)[1].lower()
    content = ""
    preview_type = "text"
    
    try:
        # Word (.docx)
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(full_path)
                paragraphs = []
                for para in doc.paragraphs[:100]:  # 최대 100 문단
                    if para.text.strip():
                        paragraphs.append(f"<p>{para.text}</p>")
                content = "\n".join(paragraphs) if paragraphs else "<p>문서가 비어있습니다.</p>"
                preview_type = "html"
            except ImportError:
                content = "python-docx 라이브러리가 필요합니다. pip install python-docx"
        
        # Excel (.xlsx)
        elif ext in ['.xlsx', '.xls']:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(full_path, read_only=True, data_only=True)
                sheet = wb.active
                rows = []
                for i, row in enumerate(sheet.iter_rows(max_row=50, values_only=True)):
                    if i >= 50: break
                    cells = "".join([f"<td>{cell if cell is not None else ''}</td>" for cell in row[:20]])
                    rows.append(f"<tr>{cells}</tr>")
                content = f"<table border='1' style='border-collapse:collapse; width:100%;'>{''.join(rows)}</table>"
                preview_type = "html"
                wb.close()
            except ImportError:
                content = "openpyxl 라이브러리가 필요합니다. pip install openpyxl"
        
        # PowerPoint (.pptx)
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(full_path)
                slides_content = []
                for i, slide in enumerate(prs.slides[:20]):
                    if i >= 20: break
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        slides_content.append(f"<div style='border:1px solid #ccc; padding:15px; margin:10px 0; border-radius:8px;'><strong>슬라이드 {i+1}</strong><br>{'<br>'.join(slide_text)}</div>")
                content = "".join(slides_content) if slides_content else "<p>프레젠테이션이 비어있습니다.</p>"
                preview_type = "html"
            except ImportError:
                content = "python-pptx 라이브러리가 필요합니다. pip install python-pptx"
        
        # CSV
        elif ext == '.csv':
            import csv
            with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = []
                for i, row in enumerate(reader):
                    if i >= 100: break
                    cells = "".join([f"<td>{cell}</td>" for cell in row[:20]])
                    rows.append(f"<tr>{cells}</tr>")
                content = f"<table border='1' style='border-collapse:collapse; width:100%;'>{''.join(rows)}</table>"
                preview_type = "html"
        
        # JSON
        elif ext == '.json':
            with open(full_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                content = f"<pre>{json.dumps(data, ensure_ascii=False, indent=2)[:10000]}</pre>"
                preview_type = "html"
        
        else:
            content = "지원하지 않는 파일 형식입니다."
        
        return jsonify({
            'success': True,
            'content': content,
            'type': preview_type,
            'filename': os.path.basename(full_path)
        })
        
    except Exception as e:
        logger.add(f"문서 미리보기 오류: {e}", "ERROR")
        return jsonify({'success': False, 'error': str(e)})



def set_autostart(enable: bool = True):
    """Windows 시작 시 자동 실행 설정"""
    if sys.platform != 'win32':
        return
    
    key_path = r"Software\Microsoft\Windows\CurrentVersion\Run"
    app_name = "WebSharePro"
    exe_path = sys.executable
    
    # 파이썬 스크립트로 실행 중인 경우 (개발 모드)
    if not getattr(sys, 'frozen', False):
        exe_path = f'"{sys.executable}" "{os.path.abspath(__file__)}"'
    else:
        exe_path = f'"{sys.executable}"'

    try:
        import winreg
        key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_ALL_ACCESS)
        try:
            if enable:
                winreg.SetValueEx(key, app_name, 0, winreg.REG_SZ, exe_path)
            else:
                try:
                    winreg.DeleteValue(key, app_name)
                except FileNotFoundError:
                    pass
        finally:
            winreg.CloseKey(key)
    except Exception as e:
        logger.add(f"자동 실행 설정 실패: {e}", "ERROR")

# ==========================================
# 새 기능: 공유 링크
# ==========================================
# secrets는 상단에서 import됨

@app.route('/share/create', methods=['POST'])
@login_required('admin')
def create_share_link():
    """v7.0: 임시 공유 링크 생성 (비밀번호, 다운로드 제한 지원)"""
    data = request.get_json()
    path = data.get('path', '')
    hours = data.get('hours', 24)  # 기본 24시간 유효
    password = data.get('password', '')  # v7.0: 공유 링크 비밀번호
    max_downloads = data.get('max_downloads', 0)  # v7.0: 최대 다운로드 횟수 (0=무제한)
    
    # 경로 검증
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid or not os.path.exists(full_path):
        return jsonify({'success': False, 'error': '유효하지 않은 경로입니다.'}), 400
    
    # 토큰 생성
    token = secrets.token_urlsafe(16)
    expires = datetime.now() + timedelta(hours=hours)
    
    SHARE_LINKS[token] = {
        'path': path,
        'expires': expires,
        'created_by': session.get('role', 'unknown'),
        'is_dir': os.path.isdir(full_path),
        'password_hash': hash_password(password) if password else None,  # v7.0
        'max_downloads': max_downloads,  # v7.0
        'download_count': 0,  # v7.0
        'created_at': datetime.now().isoformat()  # v7.0
    }
    
    features = []
    if password: features.append('비밀번호')
    if max_downloads > 0: features.append(f'최대 {max_downloads}회')
    feature_str = f" [{', '.join(features)}]" if features else ""
    
    logger.add(f"공유 링크 생성: {path} ({hours}시간){feature_str}")
    return jsonify({
        'success': True,
        'token': token,
        'expires': expires.isoformat(),
        'link': f"/share/{token}",
        'has_password': bool(password),
        'max_downloads': max_downloads
    })

@app.route('/share/<token>', methods=['GET', 'POST'])
def access_share_link(token):
    """v7.0: 공유 링크로 파일 접근 (비밀번호, 다운로드 제한 지원)"""
    if token not in SHARE_LINKS:
        return abort(404)
    
    share_info = SHARE_LINKS[token]
    
    # 만료 확인
    if datetime.now() > share_info['expires']:
        del SHARE_LINKS[token]
        return abort(410)  # Gone
    
    # v7.0: 다운로드 횟수 제한 확인
    max_downloads = share_info.get('max_downloads', 0)
    if max_downloads > 0 and share_info.get('download_count', 0) >= max_downloads:
        return render_template_string(SHARE_EXPIRED_TEMPLATE, 
            message="다운로드 횟수가 초과되었습니다.")
    
    # v7.0: 비밀번호 확인 (타이밍 공격 방지)
    password_hash = share_info.get('password_hash')
    if password_hash:
        if request.method == 'POST':
            entered_password = request.form.get('password', '')
            entered_hash = hash_password(entered_password)
            # 타이밍 공격 방지를 위한 안전한 비교
            if not secrets.compare_digest(entered_hash, password_hash):
                return render_template_string(SHARE_PASSWORD_TEMPLATE, 
                    token=token, error="비밀번호가 올바르지 않습니다.")
        else:
            # GET 요청 시 비밀번호 폼 표시
            return render_template_string(SHARE_PASSWORD_TEMPLATE, token=token, error=None)
    
    # 경로 검증
    is_valid, full_path, error = validate_path(conf.get('folder'), share_info['path'])
    if not is_valid or not os.path.exists(full_path):
        return abort(404)
    
    # v7.0: 다운로드 횟수 증가
    share_info['download_count'] = share_info.get('download_count', 0) + 1
    
    if share_info['is_dir']:
        # 폴더인 경우 ZIP으로 다운로드
        mem_zip = io.BytesIO()
        with zipfile.ZipFile(mem_zip, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(full_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    zf.write(file_path, os.path.relpath(file_path, full_path))
        mem_zip.seek(0)
        return send_file(mem_zip, download_name=f"{os.path.basename(full_path)}.zip", as_attachment=True)
    else:
        return send_from_directory(conf.get('folder'), share_info['path'])

@app.route('/share/list')
@login_required('admin')
def list_share_links():
    """활성 공유 링크 목록"""
    now = datetime.now()
    active_links = []
    expired_tokens = []
    
    # 만료된 링크 정리 및 락 적용
    with _share_links_lock:
        for token, info in SHARE_LINKS.items():
            if now > info['expires']:
                expired_tokens.append(token)
            else:
                active_links.append({
                    'token': token,
                    'path': info['path'],
                    'expires': info['expires'].isoformat(),
                    'is_dir': info['is_dir']
                })
        
        for token in expired_tokens:
            del SHARE_LINKS[token]
    
    return jsonify({'links': active_links})

@app.route('/share/delete/<token>', methods=['POST'])
@login_required('admin')
def delete_share_link(token):
    """공유 링크 삭제"""
    with _share_links_lock:
        if token in SHARE_LINKS:
            del SHARE_LINKS[token]
            return jsonify({'success': True})
    return jsonify({'success': False, 'error': '링크를 찾을 수 없습니다.'})

# ==========================================
# 새 기능: 파일 정보 상세
# ==========================================
@app.route('/file_info/<path:path>')
@login_required()
def get_file_info(path):
    """파일 상세 정보 조회"""
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid or not os.path.exists(full_path):
        return jsonify({'error': '파일을 찾을 수 없습니다.'}), 404
    
    stat = os.stat(full_path)
    info = {
        'name': os.path.basename(full_path),
        'path': path,
        'is_dir': os.path.isdir(full_path),
        'size': stat.st_size,
        'created': datetime.fromtimestamp(stat.st_ctime).isoformat(),
        'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
        'accessed': datetime.fromtimestamp(stat.st_atime).isoformat(),
    }
    
    if not info['is_dir']:
        # 파일 해시 계산 (청크 단위, 메모리 효율적)
        if stat.st_size < 10 * 1024 * 1024:  # 10MB 이하
            try:
                md5_hash = hashlib.md5()
                with open(full_path, 'rb') as f:
                    for chunk in iter(lambda: f.read(8192), b''):
                        md5_hash.update(chunk)
                info['md5'] = md5_hash.hexdigest()
            except Exception:
                pass
        
        # MIME 타입
        mime_type, _ = mimetypes.guess_type(full_path)
        info['mime_type'] = mime_type or 'application/octet-stream'
    else:
        # 폴더 내 파일/폴더 개수
        try:
            items = os.listdir(full_path)
            info['file_count'] = len([i for i in items if os.path.isfile(os.path.join(full_path, i))])
            info['folder_count'] = len([i for i in items if os.path.isdir(os.path.join(full_path, i))])
        except Exception:
            pass
    
    return jsonify(info)

# ==========================================
# 새 기능: 북마크
# ==========================================
@app.route('/bookmarks', methods=['GET', 'POST', 'DELETE'])
@login_required()
def handle_bookmarks():
    """북마크 관리"""
    global BOOKMARKS
    
    if request.method == 'GET':
        return jsonify({'bookmarks': BOOKMARKS})
    
    elif request.method == 'POST':
        data = request.get_json()
        path = data.get('path', '')
        name = data.get('name', os.path.basename(path))
        
        # 중복 확인
        if any(b['path'] == path for b in BOOKMARKS):
            return jsonify({'success': False, 'error': '이미 북마크되어 있습니다.'})
        
        BOOKMARKS.append({'path': path, 'name': name, 'added': datetime.now().isoformat()})
        return jsonify({'success': True})
    
    elif request.method == 'DELETE':
        data = request.get_json()
        path = data.get('path', '')
        BOOKMARKS = [b for b in BOOKMARKS if b['path'] != path]
        return jsonify({'success': True})

# ==========================================
# v5.1 신규 API: 확장 기능
# ==========================================

@app.route('/recent_files')
@login_required()
def get_recent_files():
    """v5.1: 최근 파일 목록"""
    return jsonify({'files': RECENT_FILES})

@app.route('/folder_size/<path:folder_path>')
@login_required()
def api_folder_size(folder_path):
    """v5.1: 폴더 크기 계산 (비동기)"""
    is_valid, full_path, error = validate_path(conf.get('folder'), folder_path)
    if not is_valid or not os.path.isdir(full_path):
        return jsonify({'error': '폴더를 찾을 수 없습니다.'}), 404
    
    size = get_folder_size(full_path)
    # 포맷팅
    if size < 1024:
        size_str = f"{size} B"
    elif size < 1024 * 1024:
        size_str = f"{size / 1024:.1f} KB"
    elif size < 1024 * 1024 * 1024:
        size_str = f"{size / 1024 / 1024:.1f} MB"
    else:
        size_str = f"{size / 1024 / 1024 / 1024:.2f} GB"
    
    return jsonify({'size': size, 'size_formatted': size_str})

@app.route('/active_sessions')
@login_required('admin')
def get_active_sessions():
    """v5.1: 활성 세션 목록 (접속자 모니터링)"""
    now = datetime.now()
    timeout = conf.get('session_timeout') or SESSION_TIMEOUT_MINUTES
    active = []
    
    for sid, info in ACTIVE_SESSIONS.items():
        # 타임아웃 체크
        last_active = info.get('last_active')
        if last_active:
            elapsed = (now - last_active).total_seconds() / 60
            if elapsed < timeout:
                active.append({
                    'ip': info.get('ip', 'unknown'),
                    'role': info.get('role', 'guest'),
                    'login_time': info.get('login_time', now).isoformat(),
                    'last_active': last_active.isoformat(),
                    'idle_minutes': round(elapsed, 1)
                })
    
    return jsonify({'sessions': active, 'count': len(active)})

@app.route('/set_language/<lang>')
def set_language(lang):
    """v5.1: 언어 변경"""
    if lang in I18N:
        session['language'] = lang
        return jsonify({'success': True, 'language': lang})
    return jsonify({'success': False, 'error': '지원하지 않는 언어입니다.'}), 400

# 주의: /move 라우트는 라인 2918의 move_item() 함수에서 이미 정의됨

@app.route('/disk_status')
@login_required()
def get_disk_status():
    """v5.1: 디스크 상태 및 경고"""
    try:
        t, u, f = shutil.disk_usage(conf.get('folder'))
        percent = round((u / t) * 100, 1)
        threshold = conf.get('disk_warning_threshold') or 90
        
        return jsonify({
            'total': f"{t / 1024**3:.1f}GB",
            'used': f"{u / 1024**3:.1f}GB",
            'free': f"{f / 1024**3:.1f}GB",
            'percent': percent,
            'warning': percent >= threshold,
            'threshold': threshold
        })
    except (OSError, IOError) as e:
        return jsonify({'error': str(e)}), 500

# ==========================================
# 새 기능: 휴지통 (Trash)
# ==========================================
@app.route('/trash', methods=['POST'])
@login_required('admin')
def move_to_trash():
    """파일을 휴지통으로 이동"""
    data = request.get_json()
    path = data.get('path', '')
    
    is_valid, full_path, error = validate_path(conf.get('folder'), path)
    if not is_valid or not os.path.exists(full_path):
        return jsonify({'success': False, 'error': '파일을 찾을 수 없습니다.'}), 404
    
    # 휴지통 폴더 생성
    trash_dir = os.path.join(conf.get('folder'), TRASH_FOLDER_NAME)
    os.makedirs(trash_dir, exist_ok=True)
    
    # 타임스탬프를 붙여 이동
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    base_name = os.path.basename(full_path)
    trash_name = f"{timestamp}_{base_name}"
    trash_path = os.path.join(trash_dir, trash_name)
    
    try:
        shutil.move(full_path, trash_path)
        logger.add(f"휴지통 이동: {path}")
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/trash/list')
@login_required('admin')
def list_trash():
    """휴지통 목록"""
    trash_dir = os.path.join(conf.get('folder'), TRASH_FOLDER_NAME)
    if not os.path.exists(trash_dir):
        return jsonify({'items': []})
    
    items = []
    for name in os.listdir(trash_dir):
        full_path = os.path.join(trash_dir, name)
        stat = os.stat(full_path)
        items.append({
            'name': name,
            'original_name': extract_original_name_from_trash(name),
            'is_dir': os.path.isdir(full_path),
            'size': stat.st_size,
            'deleted_at': datetime.fromtimestamp(stat.st_mtime).isoformat()
        })
    
    return jsonify({'items': items})

@app.route('/trash/restore', methods=['POST'])
@login_required('admin')
def restore_from_trash():
    """휴지통에서 복원"""
    data = request.get_json()
    name = data.get('name', '')
    
    trash_dir = os.path.join(conf.get('folder'), TRASH_FOLDER_NAME)
    trash_path = os.path.join(trash_dir, safe_filename(name))
    
    if not os.path.exists(trash_path):
        return jsonify({'success': False, 'error': '파일을 찾을 수 없습니다.'})
    
    # 원래 이름 추출 (정규식으로 타임스탬프 제거)
    original_name = extract_original_name_from_trash(name)
    restore_path = os.path.join(conf.get('folder'), original_name)
    
    # 동일 이름 파일 존재 시 이름 변경 (덮어쓰기 방지)
    if os.path.exists(restore_path):
        base, ext = os.path.splitext(original_name)
        counter = 1
        while os.path.exists(restore_path):
            restore_path = os.path.join(conf.get('folder'), f"{base}_복원{counter}{ext}")
            counter += 1
            
    try:
        shutil.move(trash_path, restore_path)
        restored_name = os.path.basename(restore_path)
        logger.add(f"휴지통 복원: {name} -> {restored_name}")
        return jsonify({'success': True, 'restored_name': restored_name})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/trash/empty', methods=['POST'])
@login_required('admin')
def empty_trash():
    """휴지통 비우기"""
    trash_dir = os.path.join(conf.get('folder'), TRASH_FOLDER_NAME)
    if os.path.exists(trash_dir):
        try:
            shutil.rmtree(trash_dir)
            logger.add("휴지통 비움")
            return jsonify({'success': True})
        except Exception as e:
            return jsonify({'success': False, 'error': str(e)})
    return jsonify({'success': True})

@app.route('/logout')
def logout():
    session.clear()
    return redirect('/')

# ==========================================
# v6.0: 미디어 스트리밍 (Media Streaming)
# ==========================================
@app.route('/stream/<path:filepath>')
@login_required()
def stream_media(filepath):
    """HTTP Range 요청을 지원하는 미디어 스트리밍"""
    is_valid, full_path, error = validate_path(conf.get('folder'), filepath)
    if not is_valid or not os.path.exists(full_path):
        return abort(404)
    
    file_size = os.path.getsize(full_path)
    mime_type, _ = mimetypes.guess_type(full_path)
    if not mime_type:
        mime_type = 'application/octet-stream'
    
    range_header = request.headers.get('Range')
    
    if range_header:
        # Range 요청 파싱
        byte_start = 0
        byte_end = file_size - 1
        
        match = re.match(r'bytes=(\d+)-(\d*)', range_header)
        if match:
            byte_start = int(match.group(1))
            if match.group(2):
                byte_end = int(match.group(2))
        
        byte_end = min(byte_end, file_size - 1)
        content_length = byte_end - byte_start + 1
        
        def generate():
            with open(full_path, 'rb') as f:
                f.seek(byte_start)
                remaining = content_length
                chunk_size = 1024 * 1024  # 1MB chunks
                while remaining > 0:
                    read_size = min(chunk_size, remaining)
                    data = f.read(read_size)
                    if not data:
                        break
                    remaining -= len(data)
                    yield data
        
        response = app.response_class(
            generate(),
            status=206,
            mimetype=mime_type,
            direct_passthrough=True
        )
        response.headers['Content-Range'] = f'bytes {byte_start}-{byte_end}/{file_size}'
        response.headers['Content-Length'] = content_length
        response.headers['Accept-Ranges'] = 'bytes'
        return response
    else:
        # 전체 파일 스트리밍
        def generate_full():
            with open(full_path, 'rb') as f:
                while True:
                    data = f.read(1024 * 1024)
                    if not data:
                        break
                    yield data
        
        response = app.response_class(
            generate_full(),
            mimetype=mime_type,
            direct_passthrough=True
        )
        response.headers['Content-Length'] = file_size
        response.headers['Accept-Ranges'] = 'bytes'
        return response

@app.route('/playlist/<path:folder_path>')
@login_required()
def get_playlist(folder_path):
    """폴더 내 오디오 파일 플레이리스트"""
    is_valid, full_path, error = validate_path(conf.get('folder'), folder_path)
    if not is_valid or not os.path.isdir(full_path):
        return jsonify({'error': '폴더를 찾을 수 없습니다.'}), 404
    
    audio_extensions = {'.mp3', '.wav', '.ogg', '.m4a', '.flac', '.aac', '.wma'}
    tracks = []
    
    for name in sorted(os.listdir(full_path)):
        ext = os.path.splitext(name)[1].lower()
        if ext in audio_extensions:
            rel_path = os.path.join(folder_path, name).replace('\\', '/')
            tracks.append({
                'name': name,
                'path': rel_path,
                'stream_url': f'/stream/{rel_path}'
            })
    
    return jsonify({'folder': folder_path, 'tracks': tracks, 'count': len(tracks)})

@app.route('/gallery/<path:folder_path>')
@login_required()
def get_gallery(folder_path):
    """폴더 내 이미지 갤러리"""
    is_valid, full_path, error = validate_path(conf.get('folder'), folder_path)
    if not is_valid or not os.path.isdir(full_path):
        return jsonify({'error': '폴더를 찾을 수 없습니다.'}), 404
    
    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg'}
    images = []
    
    for name in sorted(os.listdir(full_path)):
        ext = os.path.splitext(name)[1].lower()
        if ext in image_extensions:
            rel_path = os.path.join(folder_path, name).replace('\\', '/')
            images.append({
                'name': name,
                'path': rel_path,
                'url': f'/download/{rel_path}',
                'thumbnail': f'/thumbnail/{rel_path}'
            })
    
    return jsonify({'folder': folder_path, 'images': images, 'count': len(images)})

# ==========================================
# v6.0: 다중 사용자 관리 (Multi-User Management)
# ==========================================
USERS_FILE = "webshare_users.json"

def load_users():
    """사용자 목록 로드"""
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            pass
    # 기본 사용자 (기존 설정과 호환)
    return {
        'users': {
            '_legacy_admin': {
                'password_hash': conf.get('admin_pw', '1234'),
                'role': 'admin',
                'quota_mb': 0,
                'folders': ['*'],
                'created': datetime.now().isoformat()
            },
            '_legacy_guest': {
                'password_hash': conf.get('guest_pw', '0000'),
                'role': 'guest',
                'quota_mb': 0,
                'folders': ['*'],
                'created': datetime.now().isoformat()
            }
        }
    }

def save_users(users_data):
    """사용자 목록 저장"""
    try:
        with open(USERS_FILE, 'w', encoding='utf-8') as f:
            json.dump(users_data, f, indent=2, ensure_ascii=False)
        return True
    except IOError as e:
        logger.add(f"사용자 저장 실패: {e}", "ERROR")
        return False

def get_user_usage(username):
    """사용자 업로드 용량 계산"""
    user_folder = os.path.join(conf.get('folder'), f'_user_{username}')
    if os.path.exists(user_folder):
        return get_folder_size(user_folder)
    return 0

@app.route('/api/users', methods=['GET', 'POST'])
@login_required('admin')
def manage_users():
    """사용자 목록 조회 및 생성"""
    users_data = load_users()
    
    if request.method == 'GET':
        # 비밀번호 해시 제외하고 반환
        safe_users = {}
        for username, info in users_data.get('users', {}).items():
            if not username.startswith('_legacy_'):
                safe_users[username] = {
                    'role': info.get('role', 'user'),
                    'quota_mb': info.get('quota_mb', 0),
                    'folders': info.get('folders', []),
                    'created': info.get('created', ''),
                    'usage_mb': round(get_user_usage(username) / 1024 / 1024, 2)
                }
        return jsonify({'users': safe_users})
    
    elif request.method == 'POST':
        data = request.get_json()
        username = data.get('username', '').strip()
        password = data.get('password', '')
        role = data.get('role', 'user')
        quota_mb = data.get('quota_mb', 1024)
        folders = data.get('folders', [])
        
        if not username or not password:
            return jsonify({'success': False, 'error': '사용자명과 비밀번호가 필요합니다.'}), 400
        
        # v7.1: 사용자명 유효성 검증
        if not re.match(r'^[a-zA-Z0-9_\-]{3,20}$', username):
            return jsonify({'success': False, 
                'error': '사용자명은 3-20자의 영문, 숫자, _, -만 사용 가능합니다.'}), 400
        
        if username in users_data.get('users', {}):
            return jsonify({'success': False, 'error': '이미 존재하는 사용자입니다.'}), 400
        
        # 사용자 생성
        users_data.setdefault('users', {})[username] = {
            'password_hash': hash_password(password),
            'role': role,
            'quota_mb': quota_mb,
            'folders': folders if folders else [f'/_user_{username}'],
            'created': datetime.now().isoformat()
        }
        
        # 사용자 폴더 생성
        user_folder = os.path.join(conf.get('folder'), f'_user_{username}')
        os.makedirs(user_folder, exist_ok=True)
        
        if save_users(users_data):
            logger.add(f"사용자 생성: {username}")
            return jsonify({'success': True})
        return jsonify({'success': False, 'error': '저장 실패'}), 500

@app.route('/api/users/<username>', methods=['GET', 'PUT', 'DELETE'])
@login_required('admin')
def manage_single_user(username):
    """개별 사용자 관리"""
    users_data = load_users()
    users = users_data.get('users', {})
    
    if username not in users:
        return jsonify({'error': '사용자를 찾을 수 없습니다.'}), 404
    
    if request.method == 'GET':
        info = users[username]
        return jsonify({
            'username': username,
            'role': info.get('role', 'user'),
            'quota_mb': info.get('quota_mb', 0),
            'folders': info.get('folders', []),
            'created': info.get('created', ''),
            'usage_mb': round(get_user_usage(username) / 1024 / 1024, 2)
        })
    
    elif request.method == 'PUT':
        data = request.get_json()
        
        if 'password' in data and data['password']:
            users[username]['password_hash'] = hash_password(data['password'])
        if 'role' in data:
            users[username]['role'] = data['role']
        if 'quota_mb' in data:
            users[username]['quota_mb'] = data['quota_mb']
        if 'folders' in data:
            users[username]['folders'] = data['folders']
        
        if save_users(users_data):
            logger.add(f"사용자 수정: {username}")
            return jsonify({'success': True})
        return jsonify({'success': False, 'error': '저장 실패'}), 500
    
    elif request.method == 'DELETE':
        if username.startswith('_legacy_'):
            return jsonify({'success': False, 'error': '기본 사용자는 삭제할 수 없습니다.'}), 400
        
        del users[username]
        if save_users(users_data):
            logger.add(f"사용자 삭제: {username}")
            return jsonify({'success': True})
        return jsonify({'success': False, 'error': '저장 실패'}), 500

# ==========================================
# v6.0: 청크 업로드 (Chunk Upload)
# ==========================================
UPLOAD_SESSIONS = {}  # {session_id: {'filename': str, 'path': str, 'chunks': [], 'total_size': int}}

@app.route('/upload/chunk/init', methods=['POST'])
def init_chunk_upload():
    """청크 업로드 세션 시작"""
    if not (session.get('role') == 'admin' or conf.get('allow_guest_upload')):
        return jsonify({'error': '권한 없음'}), 403
    
    data = request.get_json()
    filename = safe_filename(data.get('filename', 'unnamed'))
    total_size = data.get('total_size', 0)
    target_path = data.get('path', '')
    
    # v7.1: 경로 검증 추가 (Path Traversal 방지)
    is_valid, full_target_path, error = validate_path(conf.get('folder'), target_path)
    if not is_valid:
        return jsonify({'error': '잘못된 경로입니다.'}), 400
    
    session_id = secrets.token_urlsafe(16)
    with _upload_session_lock:
        UPLOAD_SESSIONS[session_id] = {
            'filename': filename,
            'path': target_path, # 상대 경로 저장
        'chunks': [],
        'total_size': total_size,
        'received': 0,
        'created': datetime.now()
    }
    
    return jsonify({
        'success': True,
        'session_id': session_id,
        'chunk_size': 5 * 1024 * 1024  # 5MB 권장 청크 크기
    })

@app.route('/upload/chunk/<session_id>', methods=['POST'])
def upload_chunk(session_id):
    """청크 데이터 업로드"""
    if session_id not in UPLOAD_SESSIONS:
        return jsonify({'error': '세션을 찾을 수 없습니다.'}), 404
    
    with _upload_session_lock:
        upload_info = UPLOAD_SESSIONS[session_id]
        
    chunk_index = request.form.get('index', type=int)
    chunk_file = request.files.get('chunk')
    
    if chunk_file is None:
        return jsonify({'error': '청크 데이터가 없습니다.'}), 400
    
    # 임시 파일에 청크 저장
    temp_dir = os.path.join(conf.get('folder'), '.webshare_uploads', session_id)
    os.makedirs(temp_dir, exist_ok=True)
    
    chunk_path = os.path.join(temp_dir, f'chunk_{chunk_index:05d}')
    chunk_file.save(chunk_path)
    
    chunk_size = os.path.getsize(chunk_path)
    with _upload_session_lock:
        upload_info['received'] += chunk_size
        upload_info['chunks'].append(chunk_index)
        
        progress = round((upload_info['received'] / upload_info['total_size']) * 100, 1) if upload_info['total_size'] > 0 else 0
    
    return jsonify({
        'success': True,
        'chunk_index': chunk_index,
        'received': upload_info['received'],
        'progress': progress
    })

@app.route('/upload/chunk/<session_id>/complete', methods=['POST'])
def complete_chunk_upload(session_id):
    """청크 업로드 완료 및 파일 병합"""
    if session_id not in UPLOAD_SESSIONS:
        return jsonify({'error': '세션을 찾을 수 없습니다.'}), 404
    
    with _upload_session_lock:
        upload_info = UPLOAD_SESSIONS[session_id]
        
    temp_dir = os.path.join(conf.get('folder'), '.webshare_uploads', session_id)
    
    # 최종 파일 경로
    # v7.1: 경로 재검증
    is_valid, full_dir, _ = validate_path(conf.get('folder'), upload_info['path'])
    if not is_valid:
        return jsonify({'error': '잘못된 경로'}), 400
        
    target_dir = full_dir # 이미 thread-local full_dir 사용 가능하지만 validate_path가 리턴한 값 사용
    os.makedirs(target_dir, exist_ok=True) # full_dir은 절대경로여야 함. validate_path 리턴값 확인 필요 -> full_path 리턴함
    
    final_path = os.path.join(target_dir, upload_info['filename'])
    
    # v7.1: 최종 쓰기 경로도 검증 (이중 체크)
    is_valid_file, _, _ = validate_path(conf.get('folder'), os.path.join(upload_info['path'], upload_info['filename']))
    if not is_valid_file:
         return jsonify({'error': '잘못된 파일 경로'}), 400
    
    try:
        # 청크 병합
        with open(final_path, 'wb') as outfile:
            for chunk_index in sorted(upload_info['chunks']):
                chunk_path = os.path.join(temp_dir, f'chunk_{chunk_index:05d}')
                if os.path.exists(chunk_path):
                    with open(chunk_path, 'rb') as infile:
                        outfile.write(infile.read())
        
        # 임시 폴더 정리
        # 임시 폴더 정리
        shutil.rmtree(temp_dir, ignore_errors=True)
        with _upload_session_lock:
            if session_id in UPLOAD_SESSIONS:
                del UPLOAD_SESSIONS[session_id]
        
        logger.add(f"청크 업로드 완료: {upload_info['filename']}")
        STATS['bytes_received'] += os.path.getsize(final_path)
        
        return jsonify({'success': True, 'filename': upload_info['filename']})
    except Exception as e:
        logger.add(f"청크 병합 오류: {e}", "ERROR")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/upload/chunk/<session_id>/cancel', methods=['POST'])
def cancel_chunk_upload(session_id):
    """청크 업로드 취소"""
    if session_id in UPLOAD_SESSIONS:
        temp_dir = os.path.join(conf.get('folder'), '.webshare_uploads', session_id)
        shutil.rmtree(temp_dir, ignore_errors=True)
        with _upload_session_lock:
            if session_id in UPLOAD_SESSIONS:
                del UPLOAD_SESSIONS[session_id]
    return jsonify({'success': True})

def cleanup_expired_upload_sessions():
    """만료된 청크 업로드 세션 정리 (1시간 이상 된 세션 삭제)"""
    now = datetime.now()
    expired_sessions = []
    max_age_hours = 1  # 1시간 이상 된 세션 정리
    
    max_age_hours = 1  # 1시간 이상 된 세션 정리
    
    with _upload_session_lock:
        current_sessions = list(UPLOAD_SESSIONS.items())
        
    for session_id, info in current_sessions:
        created = info.get('created')
        if created:
            age_hours = (now - created).total_seconds() / 3600
            if age_hours >= max_age_hours:
                expired_sessions.append(session_id)
    
    for session_id in expired_sessions:
        try:
            temp_dir = os.path.join(conf.get('folder'), '.webshare_uploads', session_id)
            shutil.rmtree(temp_dir, ignore_errors=True)
            with _upload_session_lock:
                if session_id in UPLOAD_SESSIONS:
                    del UPLOAD_SESSIONS[session_id]
            logger.add(f"만료 업로드 세션 정리: {session_id}")
        except Exception as e:
            logger.add(f"세션 정리 오류: {e}", "ERROR")
    
    return len(expired_sessions)


# ==========================================
# 5. 서버 스레드 관리 (Aggressive Shutdown)
# ==========================================
class ServerThread(threading.Thread):
    def __init__(self, use_https=False):
        threading.Thread.__init__(self)
        self.server = None
        self.daemon = True
        self.use_https = use_https
        self.port = int(conf.get('port'))

    def run(self):
        try:
            log = logging.getLogger('werkzeug')
            log.setLevel(logging.ERROR)
            
            ssl_ctx = None
            proto = "http"
            if self.use_https:
                try:
                    # cryptography 라이브러리가 없으면 에러 발생 가능
                    ssl_ctx = 'adhoc' 
                    proto = "https"
                except Exception as e:
                    logger.add(f"HTTPS(adhoc) 설정 실패: {e}\nHTTP로 전환합니다.", "ERROR")
                    self.use_https = False
                    ssl_ctx = None
                    proto = "http"

            import werkzeug.serving
            if hasattr(werkzeug.serving, 'make_server'):
                # Werkzeug 서버 생성
                self.server = make_server(
                    conf.get('display_host'), 
                    self.port, 
                    app,
                    threaded=True,
                    ssl_context=ssl_ctx
                )
            else:
                logger.add("Werkzeug 버전 호환성 경고: make_server를 찾을 수 없습니다.", "WARN")
                return

            logger.add(f"서버 시작: {proto}://{conf.get('display_host')}:{self.port}")
            
            # serve_forever 실행 (shutdown 시 socket error가 날 수 있으므로 예외 처리)
            try:
                self.server.serve_forever()
            except OSError:
                pass # 서버 소켓이 강제 종료되면 발생하는 정상적인 현상
            except Exception as e:
                logger.add(f"서버 실행 중 오류: {e}", "ERROR")
                
        except OSError as e:
            if e.errno == 98 or e.errno == 10048: # Address already in use
                logger.add(f"포트 {self.port}가 이미 사용 중입니다.", "ERROR")
            else:
                logger.add(f"서버 시작 오류: {e}", "ERROR")
        except Exception as e:
            logger.add(f"서버 치명적 오류: {e}", "ERROR")

    def shutdown(self):
        if self.server:
            try:
                logger.add("서버 종료 신호 전송 중...")
                
                # [강력한 종료 로직]
                # 1. 종료 플래그 설정 (모든 가능성 고려)
                if hasattr(self.server, '_BaseServer__shutdown_request'):
                    self.server._BaseServer__shutdown_request = True
                if hasattr(self.server, '_shutdown_request'):
                    self.server._shutdown_request = True
                    
                # 2. 소켓 강제 종료 (블로킹 해제 핵심)
                if hasattr(self.server, 'socket') and self.server.socket:
                    try:
                        import socket
                        self.server.socket.shutdown(socket.SHUT_RDWR)
                    except Exception: pass
                    try:
                        self.server.socket.close()
                    except Exception: pass
                
                # 3. 공식 shutdown 호출 (타임아웃 적용?)
                # serve_forever가 루프를 돌고 있다면, 위 소켓 close로 인해 이미 에러가 나거나
                # 플래그 체크로 종료되었을 것임.
                try:
                    self.server.shutdown()
                except Exception: pass
                
                try:
                    self.server.server_close()
                except Exception: pass
                
            except Exception as e:
                logger.add(f"서버 종료 중 예외 (무시됨): {e}", "WARN")
            
            logger.add("서버가 중지되었습니다.")

server_thread = None

# ==========================================
# 6. Modern GUI (PyQt6 with Tkinter fallback)
# ==========================================

if PYQT6_AVAILABLE:
    # ==========================================
    # PyQt6 Modern GUI Implementation
    # ==========================================
    
    STYLESHEET = """
    QMainWindow, QWidget {
        background-color: #0f172a;
        color: #f1f5f9;
        font-family: 'Segoe UI', 'Malgun Gothic', sans-serif;
        font-size: 13px;
    }
    
    QTabWidget::pane {
        border: 1px solid #334155;
        border-radius: 12px;
        background-color: #1e293b;
        padding: 8px;
    }
    
    QTabBar::tab {
        background-color: transparent;
        color: #94a3b8;
        padding: 14px 28px;
        margin-right: 6px;
        border-top-left-radius: 10px;
        border-top-right-radius: 10px;
        font-weight: 500;
    }
    
    QTabBar::tab:selected {
        background-color: #334155;
        color: #f1f5f9;
        font-weight: 600;
    }
    
    QTabBar::tab:hover:!selected {
        background-color: rgba(51, 65, 85, 0.5);
        color: #e2e8f0;
    }
    
    QPushButton {
        background-color: qlineargradient(x1:0, y1:0, x2:1, y2:1, stop:0 #6366f1, stop:1 #8b5cf6);
        color: white;
        border: none;
        padding: 14px 28px;
        border-radius: 12px;
        font-weight: 600;
        font-size: 13px;
    }
    
    QPushButton:hover {
        background-color: qlineargradient(x1:0, y1:0, x2:1, y2:1, stop:0 #818cf8, stop:1 #a78bfa);
    }
    
    QPushButton:pressed {
        background-color: qlineargradient(x1:0, y1:0, x2:1, y2:1, stop:0 #4f46e5, stop:1 #7c3aed);
    }
    
    QPushButton:disabled {
        background-color: #475569;
        color: #64748b;
    }
    
    QPushButton#stopBtn {
        background-color: qlineargradient(x1:0, y1:0, x2:1, y2:1, stop:0 #ef4444, stop:1 #dc2626);
    }
    
    QPushButton#stopBtn:hover {
        background-color: qlineargradient(x1:0, y1:0, x2:1, y2:1, stop:0 #f87171, stop:1 #ef4444);
    }
    
    QPushButton#outlineBtn {
        background-color: transparent;
        border: 2px solid #475569;
        color: #f1f5f9;
        border-radius: 12px;
    }
    
    QPushButton#outlineBtn:hover {
        background-color: rgba(51, 65, 85, 0.6);
        border-color: #818cf8;
    }
    
    QPushButton#outlineBtn:pressed {
        background-color: #334155;
    }
    
    QLineEdit, QComboBox {
        background-color: #1e293b;
        border: 2px solid #475569;
        border-radius: 10px;
        padding: 14px 16px;
        min-height: 22px;
        color: #f1f5f9;
        font-size: 13px;
        selection-background-color: #6366f1;
    }
    
    QComboBox {
        min-height: 24px;
        padding-right: 32px;
    }
    
    QComboBox QAbstractItemView {
        background-color: #1e293b;
        border: 1px solid #475569;
        border-radius: 8px;
        selection-background-color: #4f46e5;
        padding: 6px;
    }
    
    QLineEdit:focus, QComboBox:focus {
        border-color: #818cf8;
        background-color: #1e293b;
    }
    
    QLineEdit:hover, QComboBox:hover {
        border-color: #64748b;
    }
    
    QComboBox::drop-down {
        border: none;
        padding-right: 12px;
        width: 24px;
    }
    
    QComboBox::down-arrow {
        image: none;
        border: none;
    }
    
    QTextEdit {
        background-color: #0f172a;
        border: 2px solid #334155;
        border-radius: 10px;
        padding: 12px;
        color: #94a3b8;
        font-family: 'Cascadia Code', 'Consolas', 'Courier New', monospace;
        font-size: 12px;
        selection-background-color: #6366f1;
    }
    
    QTextEdit:focus {
        border-color: #475569;
    }
    
    QGroupBox {
        border: 2px solid #334155;
        border-radius: 12px;
        margin-top: 16px;
        padding: 20px 16px 16px 16px;
        font-weight: 600;
        color: #f1f5f9;
        background-color: rgba(30, 41, 59, 0.5);
    }
    
    QGroupBox::title {
        subcontrol-origin: margin;
        left: 16px;
        padding: 0 10px;
        color: #818cf8;
    }
    
    QCheckBox {
        color: #f1f5f9;
        spacing: 10px;
        font-size: 13px;
    }
    
    QCheckBox:hover {
        color: #e2e8f0;
    }
    
    QCheckBox::indicator {
        width: 20px;
        height: 20px;
        border-radius: 6px;
        border: 2px solid #475569;
        background-color: transparent;
    }
    
    QCheckBox::indicator:hover {
        border-color: #6366f1;
    }
    
    QCheckBox::indicator:checked {
        background-color: qlineargradient(x1:0, y1:0, x2:1, y2:1, stop:0 #6366f1, stop:1 #8b5cf6);
        border-color: #6366f1;
    }
    
    QLabel {
        color: #f1f5f9;
        font-size: 13px;
    }
    
    QLabel#subtitle {
        color: #94a3b8;
        font-size: 12px;
        font-weight: 400;
    }
    
    QLabel#statusLabel {
        font-size: 20px;
        font-weight: 700;
    }
    
    QLabel#urlLabel {
        background-color: #1e293b;
        border: 2px solid #334155;
        border-radius: 12px;
        padding: 16px;
        font-family: 'Cascadia Code', 'Consolas', monospace;
        font-size: 15px;
        color: #818cf8;
    }
    
    QLabel#urlLabel:hover {
        border-color: #475569;
        background-color: rgba(30, 41, 59, 0.8);
    }
    
    QScrollArea {
        border: none;
        background-color: transparent;
    }
    
    QScrollBar:vertical {
        background-color: #1e293b;
        width: 12px;
        border-radius: 6px;
        margin: 4px 2px 4px 2px;
    }
    
    QScrollBar::handle:vertical {
        background-color: #475569;
        border-radius: 4px;
        min-height: 30px;
    }
    
    QScrollBar::handle:vertical:hover {
        background-color: #64748b;
    }
    
    QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {
        height: 0px;
    }
    
    QScrollBar:horizontal {
        background-color: #1e293b;
        height: 12px;
        border-radius: 6px;
        margin: 2px 4px 2px 4px;
    }
    
    QScrollBar::handle:horizontal {
        background-color: #475569;
        border-radius: 4px;
        min-width: 30px;
    }
    
    QScrollBar::handle:horizontal:hover {
        background-color: #64748b;
    }
    
    QScrollBar::add-line:horizontal, QScrollBar::sub-line:horizontal {
        width: 0px;
    }
    """
    
    class WebShareGUI(QMainWindow):
        # 스레드 안전한 UI 업데이트를 위한 시그널 정의
        server_update_signal = pyqtSignal(bool)

        def __init__(self):
            super().__init__()
            self.setWindowTitle(APP_TITLE)
            
            # HiDPI 지원: 스케일링된 창 크기
            base_w, base_h = 650, 700
            self.setMinimumSize(base_w, base_h)
            self.resize(base_w, base_h + 50)
            self.setStyleSheet(STYLESHEET)
            
            self.is_closing = False
            self.log_timer = QTimer()
            self.log_timer.timeout.connect(self.process_logs)
            self.log_timer.start(200)
            
            # 시그널 연결
            self.server_update_signal.connect(self.update_ui)

            
            # v4: 실시간 통계 타이머
            self.stats_timer = QTimer()
            self.stats_timer.timeout.connect(self.update_stats)
            self.stats_timer.start(5000)  # 5초마다 업데이트
            
            # v4: 시스템 트레이 설정
            self.setup_tray()
            
            self.init_ui()
        
        def setup_tray(self):
            """v4: 시스템 트레이 아이콘 설정"""
            self.tray_icon = QSystemTrayIcon(self)
            self.tray_icon.setToolTip(APP_TITLE)
            
            # 트레이 메뉴
            tray_menu = QMenu()
            show_action = QAction("프로그램 열기", self)
            show_action.triggered.connect(self.show_normal)
            tray_menu.addAction(show_action)
            
            browser_action = QAction("브라우저로 열기", self)
            browser_action.triggered.connect(self.open_browser)
            tray_menu.addAction(browser_action)
            
            tray_menu.addSeparator()
            
            quit_action = QAction("완전 종료", self)
            quit_action.triggered.connect(self.force_quit)
            tray_menu.addAction(quit_action)
            
            self.tray_icon.setContextMenu(tray_menu)
            self.tray_icon.activated.connect(self.tray_activated)
            self.tray_icon.show()
        
        def tray_activated(self, reason):
            if reason == QSystemTrayIcon.ActivationReason.DoubleClick:
                self.show_normal()
        
        def show_normal(self):
            self.show()
            self.activateWindow()
        
        def force_quit(self):
            self.is_closing = True
            if server_thread and server_thread.is_alive():
                server_thread.shutdown()
            QApplication.quit()
            # 강제 프로세스 종료 (남은 스레드 정리)
            import os
            os._exit(0)
        
        # closeEvent는 라인 6263에 정의됨 (중복 방지)
        
        def show_notification(self, title, message):
            """v4: 시스템 알림 표시"""
            if conf.get('enable_notifications') and self.tray_icon.isVisible():
                self.tray_icon.showMessage(title, message, QSystemTrayIcon.MessageIcon.Information, 3000)
        
        def update_stats(self):
            """v4: 실시간 통계 업데이트"""
            if hasattr(self, 'stats_requests'):
                self.stats_requests.setText(f"요청: {STATS['requests']}")
            if hasattr(self, 'stats_connections'):
                self.stats_connections.setText(f"접속: {STATS['active_connections']}")
            if hasattr(self, 'stats_traffic'):
                # 트래픽 포맷팅
                total_bytes = STATS['bytes_sent'] + STATS['bytes_received']
                if total_bytes < 1024:
                    traffic_str = f"{total_bytes} B"
                elif total_bytes < 1024 * 1024:
                    traffic_str = f"{total_bytes / 1024:.1f} KB"
                elif total_bytes < 1024 * 1024 * 1024:
                    traffic_str = f"{total_bytes / 1024 / 1024:.1f} MB"
                else:
                    traffic_str = f"{total_bytes / 1024 / 1024 / 1024:.2f} GB"
                self.stats_traffic.setText(f"트래픽: {traffic_str}")
            
        def init_ui(self):
            central = QWidget()
            self.setCentralWidget(central)
            layout = QVBoxLayout(central)
            layout.setContentsMargins(20, 20, 20, 20)
            layout.setSpacing(0)
            
            # Header
            header = QHBoxLayout()
            title = QLabel("🚀 WebShare Pro")
            title.setStyleSheet("font-size: 24px; font-weight: bold; color: #818cf8;")
            header.addWidget(title)
            header.addStretch()
            version = QLabel("v5.2")
            version.setObjectName("subtitle")
            header.addWidget(version)
            layout.addLayout(header)
            layout.addSpacing(20)
            
            # Tabs
            tabs = QTabWidget()
            tabs.addTab(self.build_home_tab(), "🏠 홈")
            tabs.addTab(self.build_settings_tab(), "⚙️ 설정")
            tabs.addTab(self.build_logs_tab(), "📝 로그")
            layout.addWidget(tabs)
            
        def build_home_tab(self):
            widget = QWidget()
            layout = QVBoxLayout(widget)
            layout.setAlignment(Qt.AlignmentFlag.AlignCenter)
            layout.setSpacing(20)
            layout.setContentsMargins(40, 40, 40, 40)
            
            # Status indicator
            self.status_label = QLabel("⏹ 서버 중지됨")
            self.status_label.setObjectName("statusLabel")
            self.status_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.status_label.setStyleSheet("color: #94a3b8;")
            layout.addWidget(self.status_label)
            
            layout.addSpacing(20)
            
            # Start/Stop button
            self.toggle_btn = QPushButton("▶  서버 시작")
            self.toggle_btn.setFixedHeight(60)
            self.toggle_btn.setStyleSheet("""
                QPushButton {
                    font-size: 16px;
                    font-weight: bold;
                }
            """)
            self.toggle_btn.clicked.connect(self.toggle_server)
            layout.addWidget(self.toggle_btn)
            
            layout.addSpacing(30)
            
            # Connection info
            info_group = QGroupBox(" 📡 접속 정보")
            info_layout = QVBoxLayout(info_group)
            
            self.url_label = QLabel("-")
            self.url_label.setObjectName("urlLabel")
            self.url_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.url_label.setTextInteractionFlags(Qt.TextInteractionFlag.TextSelectableByMouse)
            info_layout.addWidget(self.url_label)
            
            btn_layout = QHBoxLayout()
            
            browser_btn = QPushButton("🌐 브라우저 열기")
            browser_btn.setObjectName("outlineBtn")
            browser_btn.clicked.connect(self.open_browser)
            btn_layout.addWidget(browser_btn)
            
            qr_btn = QPushButton("📱 QR 코드")
            qr_btn.setObjectName("outlineBtn")
            qr_btn.clicked.connect(self.show_qr)
            btn_layout.addWidget(qr_btn)
            
            # v5: 공유 폴더 열기 버튼
            folder_btn = QPushButton("📂 폴더 열기")
            folder_btn.setObjectName("outlineBtn")
            folder_btn.clicked.connect(self.open_shared_folder)
            btn_layout.addWidget(folder_btn)
            
            info_layout.addLayout(btn_layout)
            layout.addWidget(info_group)
            
            # v4.2: 실시간 통계 패널
            stats_group = QGroupBox(" 📊 실시간 통계")
            stats_layout = QHBoxLayout(stats_group)
            
            self.stats_requests = QLabel("요청: 0")
            self.stats_requests.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.stats_requests.setStyleSheet("font-size: 13px; color: #818cf8; font-weight: bold;")
            stats_layout.addWidget(self.stats_requests)
            
            self.stats_connections = QLabel("접속: 0")
            self.stats_connections.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.stats_connections.setStyleSheet("font-size: 13px; color: #22c55e; font-weight: bold;")
            stats_layout.addWidget(self.stats_connections)
            
            self.stats_traffic = QLabel("트래픽: 0 B")
            self.stats_traffic.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.stats_traffic.setStyleSheet("font-size: 13px; color: #f59e0b; font-weight: bold;")
            stats_layout.addWidget(self.stats_traffic)
            
            layout.addWidget(stats_group)
            
            layout.addStretch()
            return widget
            
        def build_settings_tab(self):
            # QScrollArea로 감싸기
            scroll = QScrollArea()
            scroll.setWidgetResizable(True)
            scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
            scroll.setStyleSheet("QScrollArea { border: none; background-color: transparent; }")
            
            widget = QWidget()
            scroll.setWidget(widget)
            
            layout = QVBoxLayout(widget)
            layout.setContentsMargins(30, 30, 30, 30)
            layout.setSpacing(15)
            
            # v5.1 UI 개선: 가독성 확보
            
            # Folder settings
            folder_label = QLabel("📂 공유 폴더")
            layout.addWidget(folder_label)
            
            folder_layout = QHBoxLayout()
            self.folder_input = QLineEdit(conf.get('folder'))
            self.folder_input.setMinimumWidth(300) # 최소 너비 설정
            folder_layout.addWidget(self.folder_input)
            
            folder_btn = QPushButton("선택")
            folder_btn.setObjectName("outlineBtn")
            folder_btn.setCursor(Qt.CursorShape.PointingHandCursor)
            folder_btn.setFixedWidth(80)
            folder_btn.clicked.connect(self.choose_folder)
            folder_layout.addWidget(folder_btn)
            layout.addLayout(folder_layout)
            
            layout.addSpacing(10)
            
            # Network settings
            net_label = QLabel("🌐 네트워크 (IP / Port)")
            layout.addWidget(net_label)
            
            net_layout = QHBoxLayout()
            self.ip_combo = QComboBox()
            ips = self.get_ip_list()
            self.ip_combo.addItems(ips)
            current = conf.get('display_host')
            if current in ips:
                self.ip_combo.setCurrentText(current)
            self.ip_combo.setMinimumWidth(200)
            net_layout.addWidget(self.ip_combo, 3)
            
            self.port_input = QLineEdit(str(conf.get('port')))
            self.port_input.setFixedWidth(80)
            self.port_input.setAlignment(Qt.AlignmentFlag.AlignCenter)
            net_layout.addWidget(self.port_input, 1)
            layout.addLayout(net_layout)
            
            layout.addSpacing(10)
            
            # Password settings
            pw_label = QLabel("🔐 비밀번호 (관리자 / 게스트)")
            layout.addWidget(pw_label)
            
            pw_layout = QHBoxLayout()
            self.admin_pw = QLineEdit(conf.get('admin_pw'))
            self.admin_pw.setEchoMode(QLineEdit.EchoMode.Password)
            self.admin_pw.setPlaceholderText("관리자 암호")
            pw_layout.addWidget(self.admin_pw)
            
            self.guest_pw = QLineEdit(conf.get('guest_pw'))
            self.guest_pw.setEchoMode(QLineEdit.EchoMode.Password)
            self.guest_pw.setPlaceholderText("게스트 암호")
            pw_layout.addWidget(self.guest_pw)
            layout.addLayout(pw_layout)
            
            layout.addSpacing(15)
            
            # Checkboxes Group
            group_box = QGroupBox("기본 설정")
            group_layout = QVBoxLayout()
            
            self.guest_upload_check = QCheckBox("게스트 업로드 허용")
            self.guest_upload_check.setChecked(conf.get('allow_guest_upload'))
            group_layout.addWidget(self.guest_upload_check)
            
            self.https_check = QCheckBox("HTTPS 사용 (자체 서명 인증서)")
            self.https_check.setChecked(conf.get('use_https'))
            group_layout.addWidget(self.https_check)
            
            group_box.setLayout(group_layout)
            layout.addWidget(group_box)
            
            layout.addSpacing(10)
            
            # Advanced Settings Group
            adv_group = QGroupBox("🔧 고급 설정")
            adv_layout = QVBoxLayout()
            
            self.versioning_check = QCheckBox("파일 버전 관리 활성화")
            self.versioning_check.setChecked(conf.get('enable_versioning'))
            adv_layout.addWidget(self.versioning_check)
            
            self.notification_check = QCheckBox("시스템 알림 활성화")
            self.notification_check.setChecked(conf.get('enable_notifications'))
            adv_layout.addWidget(self.notification_check)
            
            # Tray options
            self.tray_check = QCheckBox("최소화 버튼 시 트레이로 이동")
            self.tray_check.setChecked(conf.get('minimize_to_tray'))
            adv_layout.addWidget(self.tray_check)
            
            self.close_tray_check = QCheckBox("닫기(X) 버튼 시 트레이로 이동")
            self.close_tray_check.setChecked(conf.get('close_to_tray'))
            adv_layout.addWidget(self.close_tray_check)
            
            self.autostart_check = QCheckBox("윈도우 시작 시 자동 실행")
            self.autostart_check.setChecked(conf.get('autostart'))
            adv_layout.addWidget(self.autostart_check)

            # Session Timeout
            timeout_layout = QHBoxLayout()
            timeout_label = QLabel("세션 타임아웃 (분):")
            timeout_layout.addWidget(timeout_label)
            self.timeout_input = QLineEdit(str(conf.get('session_timeout')))
            self.timeout_input.setFixedWidth(80)
            self.timeout_input.setAlignment(Qt.AlignmentFlag.AlignCenter)
            timeout_layout.addWidget(self.timeout_input)
            timeout_layout.addStretch()
            adv_layout.addLayout(timeout_layout)
            
            adv_group.setLayout(adv_layout)
            layout.addWidget(adv_group)
            
            layout.addSpacing(20)
            
            # Save button container
            btn_layout = QHBoxLayout()
            btn_layout.addStretch()
            save_btn = QPushButton("💾 설정 저장")
            save_btn.setCursor(Qt.CursorShape.PointingHandCursor)
            save_btn.setFixedWidth(120)
            save_btn.setFixedHeight(40)
            save_btn.clicked.connect(self.save_settings)
            btn_layout.addWidget(save_btn)
            btn_layout.addStretch()
            layout.addLayout(btn_layout)
            
            layout.addStretch()
            
            return scroll
            
        def build_logs_tab(self):
            widget = QWidget()
            layout = QVBoxLayout(widget)
            layout.setContentsMargins(20, 20, 20, 20)
            
            # v5: 필터 및 도구 바
            toolbar = QHBoxLayout()
            
            filter_label = QLabel("필터:")
            toolbar.addWidget(filter_label)
            
            self.log_filter = QComboBox()
            self.log_filter.addItems(["전체", "INFO", "WARN", "ERROR"])
            self.log_filter.currentTextChanged.connect(self.filter_logs)
            self.log_filter.setFixedWidth(100)
            toolbar.addWidget(self.log_filter)
            
            toolbar.addStretch()
            
            export_btn = QPushButton("📄 내보내기")
            export_btn.setObjectName("outlineBtn")
            export_btn.clicked.connect(self.export_logs)
            toolbar.addWidget(export_btn)
            
            layout.addLayout(toolbar)
            
            self.log_text = QTextEdit()
            self.log_text.setReadOnly(True)
            self.log_text.setPlaceholderText("서버 로그가 여기에 표시됩니다...")
            layout.addWidget(self.log_text)
            
            # v5: 전체 로그 저장 (필터링용)
            self.all_logs = []
            
            btn_layout = QHBoxLayout()
            
            clear_btn = QPushButton("🗑 로그 클리어")
            clear_btn.setObjectName("outlineBtn")
            clear_btn.clicked.connect(self.clear_logs)
            btn_layout.addWidget(clear_btn)
            
            layout.addLayout(btn_layout)
            
            return widget
        
        def filter_logs(self, level):
            """v5: 로그 레벨별 필터링"""
            self.log_text.clear()
            for log in self.all_logs:
                if level == "전체" or f"[{level}]" in log:
                    self.log_text.append(log)
        
        def export_logs(self):
            """v5: 로그 파일로 내보내기"""
            from datetime import datetime
            filename = f"webshare_log_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
            filepath, _ = QFileDialog.getSaveFileName(self, "로그 저장", filename, "Text Files (*.txt)")
            if filepath:
                try:
                    with open(filepath, 'w', encoding='utf-8') as f:
                        f.write('\n'.join(self.all_logs))
                    QMessageBox.information(self, "저장 완료", f"로그가 저장되었습니다.\n{filepath}")
                except IOError as e:
                    QMessageBox.critical(self, "오류", f"저장 실패: {e}")
        
        def clear_logs(self):
            """v5: 로그 클리어"""
            self.log_text.clear()
            self.all_logs.clear()
            
        def get_ip_list(self):
            ips = ['127.0.0.1']
            try:
                s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
                s.settimeout(0.1)
                s.connect(("8.8.8.8", 80))
                ip = s.getsockname()[0]
                if ip and not ip.startswith('127.'):
                    ips.insert(0, ip)
                s.close()
            except Exception: pass
            
            try:
                host_name = socket.gethostname()
                for ip in socket.gethostbyname_ex(host_name)[2]:
                    if ip and not ip.startswith("127.") and ip not in ips:
                        ips.append(ip)
            except Exception: pass
            
            ips.append('0.0.0.0')
            return ips
            
        def choose_folder(self):
            path = QFileDialog.getExistingDirectory(self, "공유 폴더 선택")
            if path:
                self.folder_input.setText(os.path.abspath(path))
                
        def save_settings(self):
            conf.set('folder', self.folder_input.text())
            conf.set('display_host', self.ip_combo.currentText())
            try:
                conf.set('port', int(self.port_input.text()))
            except Exception: pass
            conf.set('admin_pw', self.admin_pw.text())
            conf.set('guest_pw', self.guest_pw.text())
            conf.set('allow_guest_upload', self.guest_upload_check.isChecked())
            
            try:
                conf.set('session_timeout', int(self.timeout_input.text()))
                conf.set('use_https', self.https_check.isChecked())
                conf.set('enable_versioning', self.versioning_check.isChecked())
                conf.set('enable_notifications', self.notification_check.isChecked())
                conf.set('minimize_to_tray', self.tray_check.isChecked())
                # v5.1 추가
                conf.set('close_to_tray', self.close_tray_check.isChecked())
                conf.set('autostart', self.autostart_check.isChecked())
                
                conf.save()
                
                # 자동 실행 설정 적용
                set_autostart(self.autostart_check.isChecked())
                
                QMessageBox.information(self, "저장됨", "설정이 저장되었습니다.\n일부 설정은 재시작 후 적용됩니다.")
            except ValueError:
                QMessageBox.warning(self, "오류", "세션 타임아웃은 숫자여야 합니다.")
            
        def toggle_server(self):
            global server_thread
            
            if server_thread and server_thread.is_alive():
                self.toggle_btn.setEnabled(False)
                self.toggle_btn.setText("⏳ 중지 중...")
                threading.Thread(target=self._stop_server, daemon=True).start()
            else:
                self.save_settings()
                if not os.path.exists(conf.get('folder')):
                    QMessageBox.critical(self, "오류", "공유 폴더 경로가 잘못되었습니다.")
                    return
                    
                server_thread = ServerThread(use_https=conf.get('use_https'))
                server_thread.start()
                self.update_ui(True)
                
        def _stop_server(self):
            global server_thread
            try:
                if server_thread:
                    server_thread.shutdown()
                    server_thread.join(timeout=2.0)
            except Exception as e:
                logger.add(f"서버 종료 중 오류: {e}", "ERROR")
            finally:
                server_thread = None
                if not self.is_closing:
                    # 시그널 통해 UI 업데이트 호출 (스레드 안전)
                    self.server_update_signal.emit(False)
                
        def update_ui(self, running):
            if self.is_closing:
                return
            self.toggle_btn.setEnabled(True)
            
            if running:
                self.toggle_btn.setText("⏹  서버 중지")
                self.toggle_btn.setObjectName("stopBtn")
                self.toggle_btn.setStyleSheet("""
                    QPushButton {
                        background-color: #ef4444;
                        font-size: 16px;
                        font-weight: bold;
                    }
                    QPushButton:hover { background-color: #f87171; }
                """)
                self.status_label.setText("🟢 서버 실행 중")
                self.status_label.setStyleSheet("color: #22c55e;")
                
                proto = "https" if conf.get('use_https') else "http"
                url = f"{proto}://{conf.get('display_host')}:{conf.get('port')}"
                self.url_label.setText(url)
            else:
                self.toggle_btn.setText("▶  서버 시작")
                self.toggle_btn.setObjectName("")
                self.toggle_btn.setStyleSheet("""
                    QPushButton {
                        background-color: #4f46e5;
                        font-size: 16px;
                        font-weight: bold;
                    }
                    QPushButton:hover { background-color: #6366f1; }
                """)
                self.status_label.setText("⏹ 서버 중지됨")
                self.status_label.setStyleSheet("color: #94a3b8;")
                self.url_label.setText("-")
                
        def open_browser(self):
            url = self.url_label.text()
            if url != "-":
                webbrowser.open(url)
        
        def open_shared_folder(self):
            """v5: 공유 폴더를 파일 탐색기에서 열기"""
            folder = conf.get('folder')
            if folder and os.path.exists(folder):
                try:
                    os.startfile(folder)
                except AttributeError:
                    # macOS/Linux
                    import subprocess
                    subprocess.Popen(['open' if sys.platform == 'darwin' else 'xdg-open', folder])
            else:
                QMessageBox.warning(self, "경고", "공유 폴더가 존재하지 않습니다.")
                
        def show_qr(self):
            url = self.url_label.text()
            if url == "-":
                return
            try:
                import qrcode
                qr = qrcode.make(url)
                
                dialog = QDialog(self)
                dialog.setWindowTitle("QR Code")
                dialog.setFixedSize(300, 340)
                dialog.setStyleSheet("background-color: white;")
                
                layout = QVBoxLayout(dialog)
                
                # Convert PIL image to QPixmap
                qr_bytes = io.BytesIO()
                qr.save(qr_bytes, format='PNG')
                qr_bytes.seek(0)
                
                pixmap = QPixmap()
                pixmap.loadFromData(qr_bytes.read())
                pixmap = pixmap.scaled(250, 250, Qt.AspectRatioMode.KeepAspectRatio)
                
                qr_label = QLabel()
                qr_label.setPixmap(pixmap)
                qr_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
                layout.addWidget(qr_label)
                
                text_label = QLabel("모바일로 스캔하여 접속하세요")
                text_label.setStyleSheet("color: #333; font-size: 12px;")
                text_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
                layout.addWidget(text_label)
                
                dialog.exec()
            except ImportError:
                QMessageBox.critical(self, "오류", "qrcode 라이브러리가 설치되지 않았습니다.\npip install qrcode")
                
        def process_logs(self):
            if self.is_closing:
                return
            try:
                current_filter = self.log_filter.currentText() if hasattr(self, 'log_filter') else "전체"
                while not logger.queue.empty():
                    msg = logger.queue.get()
                    # v5: 모든 로그 저장
                    if hasattr(self, 'all_logs'):
                        self.all_logs.append(msg)
                        # 최대 로그 수 제한
                        if len(self.all_logs) > MAX_LOG_LINES:
                            self.all_logs = self.all_logs[-MAX_LOG_LINES:]
                    
                    # 필터 적용
                    if current_filter == "전체" or f"[{current_filter}]" in msg:
                        self.log_text.append(msg)
                    
                    # Limit log lines in display
                    doc = self.log_text.document()
                    if doc.blockCount() > MAX_LOG_LINES:
                        cursor = self.log_text.textCursor()
                        cursor.movePosition(cursor.MoveOperation.Start)
                        cursor.movePosition(cursor.MoveOperation.Down, cursor.MoveMode.KeepAnchor, 
                                          doc.blockCount() - MAX_LOG_LINES)
                        cursor.removeSelectedText()
            except Exception: pass
            
        def closeEvent(self, event):
            global server_thread
            
            # v5.1: 완전 종료 시 트레이 로직 우회
            if self.is_closing:
                event.accept()
                return
            
            # v4: 서버 실행 중이면 트레이로 최소화 (설정에 따라)
            should_minimize = False
            if conf.get('close_to_tray'): # 항상 최소화
                should_minimize = True
            elif server_thread and server_thread.is_alive() and conf.get('minimize_to_tray'): # 서버 실행 중일 때만
                should_minimize = True
                
            if should_minimize:
                event.ignore()
                self.hide()
                self.tray_icon.showMessage(
                    "WebShare Pro",
                    "서버가 백그라운드에서 계속 실행 중입니다.",
                    QSystemTrayIcon.MessageIcon.Information,
                    2000
                )
                return
            
            self.is_closing = True
            if server_thread and server_thread.is_alive():
                reply = QMessageBox.question(self, "종료", "서버가 실행 중입니다. 종료하시겠습니까?",
                                            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
                if reply == QMessageBox.StandardButton.Yes:
                    threading.Thread(target=self._stop_server, daemon=True).start()
                    self.tray_icon.hide()
                    event.accept()
                else:
                    self.is_closing = False
                    event.ignore()
            else:
                self.tray_icon.hide()
                event.accept()

else:
    # ==========================================
    # Fallback: Tkinter GUI (if PyQt6 not available)
    # ==========================================
    class WebShareGUI:
        def __init__(self, root):
            self.root = root
            self.root.title(APP_TITLE)
            self.root.geometry("600x750")
            
            style = ttk.Style()
            style.theme_use('clam')
            style.configure("TFrame", background="#f8fafc")
            style.configure("TLabel", background="#f8fafc", font=("맑은 고딕", 10))
            style.configure("TButton", font=("맑은 고딕", 10), padding=5)
            
            self.root.configure(bg="#f8fafc")
            self.root.protocol("WM_DELETE_WINDOW", self.on_close)
            
            self.is_closing = False
            self.init_ui()
            self.process_logs()

        def init_ui(self):
            tabs = ttk.Notebook(self.root)
            tabs.pack(fill='both', expand=True, padx=15, pady=15)
            
            tab_home = ttk.Frame(tabs); tabs.add(tab_home, text="  🏠 홈  ")
            tab_set = ttk.Frame(tabs); tabs.add(tab_set, text="  ⚙️ 설정  ")
            tab_log = ttk.Frame(tabs); tabs.add(tab_log, text="  📝 로그  ")
            
            self.build_home(tab_home)
            self.build_settings(tab_set)
            self.build_logs(tab_log)

        def build_home(self, parent):
            frame = ttk.Frame(parent)
            frame.pack(fill='both', expand=True, padx=20, pady=20)
            
            self.status_lbl = ttk.Label(frame, text="서버 중지됨", font=("맑은 고딕", 16, "bold"), foreground="#64748b")
            self.status_lbl.pack(pady=20)

            self.btn_toggle = tk.Button(frame, text="서버 시작", bg="#4f46e5", fg="white", 
                                      font=("맑은 고딕", 14, "bold"), relief="flat", cursor="hand2",
                                      command=self.toggle_server)
            self.btn_toggle.pack(fill='x', pady=30, ipady=10)

            info_frame = ttk.LabelFrame(frame, text=" 접속 정보 ", padding=15)
            info_frame.pack(fill='x')
            
            self.url_var = tk.StringVar(value="-")
            url_ent = ttk.Entry(info_frame, textvariable=self.url_var, state="readonly", font=("Consolas", 12), justify="center")
            url_ent.pack(fill='x', pady=5)
            
            btn_box = ttk.Frame(info_frame)
            btn_box.pack(fill='x', pady=5)
            ttk.Button(btn_box, text="브라우저 열기", command=self.open_browser).pack(side='left', expand=True, fill='x', padx=2)
            ttk.Button(btn_box, text="QR 코드", command=self.show_qr).pack(side='right', expand=True, fill='x', padx=2)

        def build_settings(self, parent):
            frame = ttk.Frame(parent)
            frame.pack(fill='both', expand=True, padx=20, pady=20)

            ttk.Label(frame, text="공유 폴더").pack(anchor='w')
            f_box = ttk.Frame(frame); f_box.pack(fill='x', pady=5)
            self.ent_folder = ttk.Entry(f_box)
            self.ent_folder.insert(0, conf.get('folder'))
            self.ent_folder.pack(side='left', fill='x', expand=True)
            ttk.Button(f_box, text="선택", command=self.choose_folder).pack(side='right', padx=5)

            ttk.Label(frame, text="네트워크 (IP / Port)").pack(anchor='w', pady=(15, 0))
            net_box = ttk.Frame(frame); net_box.pack(fill='x', pady=5)
            
            ips = self.get_ip_list()
            self.cb_ip = ttk.Combobox(net_box, values=ips, state="readonly")
            current_host = conf.get('display_host')
            if current_host in ips: self.cb_ip.set(current_host)
            elif ips: self.cb_ip.current(0)
            
            self.cb_ip.pack(side='left', fill='x', expand=True)
            
            self.ent_port = ttk.Entry(net_box, width=8)
            self.ent_port.insert(0, conf.get('port'))
            self.ent_port.pack(side='right', padx=5)

            ttk.Label(frame, text="비밀번호 설정 (관리자 / 게스트)").pack(anchor='w', pady=(15, 0))
            pw_box = ttk.Frame(frame); pw_box.pack(fill='x', pady=5)
            self.ent_admin_pw = ttk.Entry(pw_box, show="*")
            self.ent_admin_pw.insert(0, conf.get('admin_pw'))
            self.ent_admin_pw.pack(side='left', fill='x', expand=True, padx=(0, 5))
            
            self.ent_guest_pw = ttk.Entry(pw_box, show="*")
            self.ent_guest_pw.insert(0, conf.get('guest_pw'))
            self.ent_guest_pw.pack(side='right', fill='x', expand=True)
            
            self.var_upload = tk.BooleanVar(value=conf.get('allow_guest_upload'))
            ttk.Checkbutton(frame, text="게스트 업로드 허용", variable=self.var_upload).pack(anchor='w', pady=(10, 5))

            self.var_https = tk.BooleanVar(value=conf.get('use_https'))
            ttk.Checkbutton(frame, text="HTTPS 사용 (자체 서명 인증서)", variable=self.var_https).pack(anchor='w', pady=5)
            
            ttk.Button(frame, text="설정 저장", command=self.save_settings).pack(fill='x', pady=10)

        def build_logs(self, parent):
            frame = ttk.Frame(parent)
            frame.pack(fill='both', expand=True, padx=10, pady=10)
            self.txt_log = scrolledtext.ScrolledText(frame, state='disabled', font=("Consolas", 9))
            self.txt_log.pack(fill='both', expand=True)
            ttk.Button(frame, text="로그 클리어", command=lambda: self.txt_log.configure(state='normal') or self.txt_log.delete(1.0, tk.END) or self.txt_log.configure(state='disabled')).pack(anchor='e', pady=5)

        def get_ip_list(self):
            ips = set(['0.0.0.0', '127.0.0.1'])
            try:
                s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
                s.settimeout(0.1)
                s.connect(("8.8.8.8", 80))
                ip = s.getsockname()[0]
                if ip and not ip.startswith('127.'):
                    ips.add(ip)
                s.close()
            except Exception: pass
            try:
                host_name = socket.gethostname()
                for ip in socket.gethostbyname_ex(host_name)[2]:
                    if ip and not ip.startswith("127."):
                        ips.add(ip)
            except Exception: pass
            sorted_ips = sorted(list(ips))
            if '0.0.0.0' in sorted_ips:
                sorted_ips.remove('0.0.0.0')
                sorted_ips.append('0.0.0.0')
            return sorted_ips

        def choose_folder(self):
            path = filedialog.askdirectory()
            if path:
                self.ent_folder.delete(0, tk.END)
                self.ent_folder.insert(0, os.path.abspath(path))

        def save_settings(self):
            conf.set('folder', self.ent_folder.get())
            conf.set('display_host', self.cb_ip.get())
            try: conf.set('port', int(self.ent_port.get()))
            except Exception: pass
            conf.set('admin_pw', self.ent_admin_pw.get())
            conf.set('guest_pw', self.ent_guest_pw.get())
            conf.set('allow_guest_upload', self.var_upload.get())
            conf.set('use_https', self.var_https.get())
            conf.save()
            messagebox.showinfo("저장", "설정이 저장되었습니다.")

        def toggle_server(self):
            global server_thread
            if server_thread and server_thread.is_alive():
                self.btn_toggle.config(state='disabled', text="중지 중...")
                threading.Thread(target=self._stop_server_task, daemon=True).start()
            else:
                self.save_settings()
                if not os.path.exists(conf.get('folder')):
                    messagebox.showerror("오류", "공유 폴더 경로가 잘못되었습니다.")
                    return
                server_thread = ServerThread(use_https=conf.get('use_https'))
                server_thread.start()
                self.update_ui_state(True)

        def _stop_server_task(self):
            global server_thread
            if server_thread:
                server_thread.shutdown()
                server_thread.join(timeout=2.0)
                server_thread = None
            if not self.is_closing:
                self.root.after(0, lambda: self.update_ui_state(False))

        def update_ui_state(self, running):
            if self.is_closing: return
            self.btn_toggle.config(state='normal')
            if running:
                self.btn_toggle.config(text="서버 중지", bg="#ef4444")
                self.status_lbl.config(text="서버 실행 중", foreground="#22c55e")
                proto = "https" if conf.get('use_https') else "http"
                url = f"{proto}://{conf.get('display_host')}:{conf.get('port')}"
                self.url_var.set(url)
            else:
                self.btn_toggle.config(text="서버 시작", bg="#4f46e5")
                self.status_lbl.config(text="서버 중지됨", foreground="#64748b")
                self.url_var.set("-")

        def open_browser(self):
            url = self.url_var.get()
            if url != "-": webbrowser.open(url)
        
        def show_qr(self):
            url = self.url_var.get()
            if url == "-": return
            try:
                import qrcode
                qr = qrcode.make(url)
                win = tk.Toplevel(self.root)
                win.title("QR Code")
                win.geometry("300x300")
                img_tk = ImageTk.PhotoImage(qr)
                lbl = tk.Label(win, image=img_tk)
                lbl.image = img_tk
                lbl.pack(expand=True)
                tk.Label(win, text="모바일로 스캔하여 접속하세요").pack(pady=10)
            except ImportError:
                messagebox.showerror("오류", "qrcode/pillow 라이브러리가 설치되지 않았습니다.")

        def process_logs(self):
            if self.is_closing: return
            try:
                while not logger.queue.empty():
                    msg = logger.queue.get()
                    self.txt_log.configure(state='normal')
                    self.txt_log.insert(tk.END, msg + "\n")
                    num_lines = float(self.txt_log.index('end-1c'))
                    if num_lines > MAX_LOG_LINES:
                        self.txt_log.delete('1.0', f'{num_lines - MAX_LOG_LINES + 1}.0')
                    self.txt_log.see(tk.END)
                    self.txt_log.configure(state='disabled')
            except tk.TclError:
                pass 
            self.root.after(200, self.process_logs)

        def on_close(self):
            global server_thread
            self.is_closing = True
            if server_thread and server_thread.is_alive():
                if messagebox.askokcancel("종료", "서버가 실행 중입니다. 종료하시겠습니까?"):
                    try:
                        server_thread.shutdown()
                    except Exception as e:
                        logger.add(f"서버 종료 중 오류: {e}", "ERROR")
                    finally:
                        server_thread = None
                    self.root.destroy()
                    sys.exit(0)
                else:
                    self.is_closing = False
            else:
                self.root.destroy()
                sys.exit(0)


# ==========================================
# 7. Main Entry Point
# ==========================================

def get_dpi_scale() -> float:
    """현재 디스플레이의 DPI 스케일 계수를 반환합니다."""
    try:
        from ctypes import windll
        # GetDeviceCaps를 사용하여 DPI 가져오기
        hdc = windll.user32.GetDC(0)
        dpi = windll.gdi32.GetDeviceCaps(hdc, 88)  # LOGPIXELSX
        windll.user32.ReleaseDC(0, hdc)
        return dpi / 96.0  # 96 DPI가 기준 (100%)
    except Exception:
        return 1.0

def scaled_size(base_size: int) -> int:
    """기본 크기를 DPI 스케일에 맞게 조정합니다."""
    return int(base_size * get_dpi_scale())

def cleanup_temp_files():
    """시작 시 임시 업로드 폴더 정리"""
    try:
        temp_dir = os.path.join(conf.get('folder'), '.webshare_uploads')
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
            logger.add("임시 업로드 파일 정리 완료")
    except Exception as e:
        logger.add(f"임시 파일 정리 실패: {e}", "WARN")

if __name__ == '__main__':
    # 임시 파일 정리
    cleanup_temp_files()
    
    # 메타데이터 로드 (태그, 즐겨찾기, 메모 등)
    load_metadata()

    # v7.1: 주기적 정리 스레드 (5분 간격)
    def periodic_cleanup():
        while True:
            time.sleep(300)
            try:
                # 각종 만료 리소스 정리
                cleanup_expired_sessions()
                cleanup_expired_share_links()
                cleanup_expired_upload_sessions()
                auto_cleanup_trash()
            except Exception as e:
                pass
                # logger.add(f"주기적 정리 오류: {e}", "ERROR") 
                # 메인 GUI 종료 시점 등에 오류 발생 가능성 있으므로 단순화
    
    cleanup_thread = threading.Thread(target=periodic_cleanup, daemon=True)
    cleanup_thread.start()

    # ==========================================
    # HiDPI 지원 설정
    # ==========================================
    
    # 환경 변수 설정 (Qt 초기화 전에 설정해야 함)
    os.environ['QT_ENABLE_HIGHDPI_SCALING'] = '1'
    os.environ['QT_AUTO_SCREEN_SCALE_FACTOR'] = '1'
    os.environ['QT_SCALE_FACTOR_ROUNDING_POLICY'] = 'PassThrough'
    
    # Windows DPI 인식 설정
    try:
        from ctypes import windll
        # SetProcessDpiAwareness: 1 = System DPI Aware, 2 = Per Monitor DPI Aware
        windll.shcore.SetProcessDpiAwareness(2)
    except Exception:
        try:
            windll.user32.SetProcessDPIAware()
        except Exception:
            pass

    if PYQT6_AVAILABLE:
        # PyQt6 High DPI 정책 설정
        from PyQt6.QtCore import Qt as QtCore_Qt
        from PyQt6.QtGui import QGuiApplication
        
        # 고DPI 스케일링 활성화
        QGuiApplication.setHighDpiScaleFactorRoundingPolicy(
            QtCore_Qt.HighDpiScaleFactorRoundingPolicy.PassThrough
        )
        
        qt_app = QApplication(sys.argv)
        qt_app.setStyle('Fusion')
        
        # DPI 정보 로깅
        screen = qt_app.primaryScreen()
        dpi = screen.logicalDotsPerInch()
        scale = screen.devicePixelRatio()
        logger.add(f"Display: DPI={dpi:.0f}, Scale={scale:.1f}x")
        
        window = WebShareGUI()
        window.show()
        sys.exit(qt_app.exec())
    else:
        # Tkinter HiDPI 설정
        root = tk.Tk()
        try:
            from ctypes import windll
            windll.shcore.SetProcessDpiAwareness(1)
            # Tkinter DPI 스케일링
            root.tk.call('tk', 'scaling', get_dpi_scale() * 1.5)
        except Exception:
            pass
        app_gui = WebShareGUI(root)
        root.mainloop()

